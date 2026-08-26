#!/usr/bin/env python3
from pathlib import Path


def replace_once(text: str, old: str, new: str, label: str) -> str:
    if old not in text:
        raise SystemExit(f"v7 office patch anchor missing: {label}")
    return text.replace(old, new, 1)


# ---------------------------------------------------------------------------
# Expand existing Office modify functions instead of replacing them with many
# one-off tool names. This keeps small-model schemas compact while each Skill is
# genuinely capable inside its domain.
# ---------------------------------------------------------------------------
p = Path('python/navixmind/tools/documents.py')
text = p.read_text()

# DOCX actions.
old = '''            elif action == "update_table_cell":
                table_idx = params.get("table", 0)
                row_idx = params.get("row", 0)
                col_idx = params.get("col", 0)
                text = params.get("text", "")
                if table_idx < len(doc.tables):
                    table = doc.tables[table_idx]
                    if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                        table.rows[row_idx].cells[col_idx].text = text
                        applied += 1

        output_dir = os.path.dirname(output_path)
'''
new = '''            elif action == "update_table_cell":
                table_idx = params.get("table", 0)
                row_idx = params.get("row", 0)
                col_idx = params.get("col", 0)
                text = params.get("text", "")
                if table_idx < len(doc.tables):
                    table = doc.tables[table_idx]
                    if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                        table.rows[row_idx].cells[col_idx].text = text
                        applied += 1

            elif action == "add_heading":
                doc.add_heading(params.get("text", ""), level=int(params.get("level", 1)))
                applied += 1

            elif action == "add_page_break":
                doc.add_page_break()
                applied += 1

            elif action == "add_table":
                rows = params.get("rows", []) or []
                if not isinstance(rows, list) or not rows:
                    raise ToolError("add_table requires params.rows as a non-empty 2D list")
                col_count = max(len(row) if isinstance(row, list) else 1 for row in rows)
                table = doc.add_table(rows=len(rows), cols=col_count)
                for r, row in enumerate(rows):
                    values = row if isinstance(row, list) else [row]
                    for c, value in enumerate(values):
                        table.cell(r, c).text = str(value)
                applied += 1

            elif action == "add_image":
                image_path = params.get("image_path")
                if not image_path or not os.path.isfile(image_path):
                    raise ToolError(f"add_image image not found: {image_path}")
                width_inches = params.get("width_inches")
                if width_inches is None:
                    doc.add_picture(image_path)
                else:
                    from docx.shared import Inches
                    doc.add_picture(image_path, width=Inches(float(width_inches)))
                applied += 1

        output_dir = os.path.dirname(output_path)
'''
text = replace_once(text, old, new, 'DOCX advanced actions')

# PPTX actions.
old = '''            elif action == "set_notes":
                slide_num = params.get("slide", 1) - 1
                text = params.get("text", "")
                if 0 <= slide_num < len(prs.slides):
                    slide = prs.slides[slide_num]
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = text
                    applied += 1

        output_dir = os.path.dirname(output_path)
'''
new = '''            elif action == "set_notes":
                slide_num = params.get("slide", 1) - 1
                text = params.get("text", "")
                if 0 <= slide_num < len(prs.slides):
                    slide = prs.slides[slide_num]
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = text
                    applied += 1

            elif action == "add_textbox":
                from pptx.util import Inches
                slide_num = params.get("slide", 1) - 1
                if 0 <= slide_num < len(prs.slides):
                    slide = prs.slides[slide_num]
                    box = slide.shapes.add_textbox(
                        Inches(float(params.get("left", 1))),
                        Inches(float(params.get("top", 1))),
                        Inches(float(params.get("width", 6))),
                        Inches(float(params.get("height", 1))),
                    )
                    box.text_frame.text = str(params.get("text", ""))
                    applied += 1

            elif action == "add_image":
                from pptx.util import Inches
                slide_num = params.get("slide", 1) - 1
                image_path = params.get("image_path")
                if not image_path or not os.path.isfile(image_path):
                    raise ToolError(f"add_image image not found: {image_path}")
                if 0 <= slide_num < len(prs.slides):
                    kwargs = {
                        "left": Inches(float(params.get("left", 1))),
                        "top": Inches(float(params.get("top", 1))),
                    }
                    if params.get("width") is not None:
                        kwargs["width"] = Inches(float(params["width"]))
                    if params.get("height") is not None:
                        kwargs["height"] = Inches(float(params["height"]))
                    prs.slides[slide_num].shapes.add_picture(image_path, **kwargs)
                    applied += 1

            elif action == "delete_slide":
                slide_num = params.get("slide", 1) - 1
                if 0 <= slide_num < len(prs.slides):
                    slide_id = prs.slides._sldIdLst[slide_num]
                    prs.part.drop_rel(slide_id.rId)
                    prs.slides._sldIdLst.remove(slide_id)
                    applied += 1

        output_dir = os.path.dirname(output_path)
'''
text = replace_once(text, old, new, 'PPTX advanced actions')

# XLSX action chain: insert before output save anchor, using an anchor which is
# stable in the current v6 implementation.
old = '''            elif action == "delete_sheet":
                name = params.get("name")
                if name in wb.sheetnames and len(wb.sheetnames) > 1:
                    del wb[name]
                    applied += 1

        output_dir = os.path.dirname(output_path)
'''
new = '''            elif action == "delete_sheet":
                name = params.get("name")
                if name in wb.sheetnames and len(wb.sheetnames) > 1:
                    del wb[name]
                    applied += 1

            elif action == "rename_sheet":
                old_name = params.get("old_name") or params.get("sheet")
                new_name = params.get("new_name")
                if old_name in wb.sheetnames and new_name:
                    wb[old_name].title = str(new_name)[:31]
                    applied += 1

            elif action in {"insert_row", "delete_row", "insert_column", "delete_column"}:
                sheet_name = params.get("sheet") or wb.active.title
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {sheet_name}")
                ws = wb[sheet_name]
                index = int(params.get("index", 1))
                amount = int(params.get("amount", 1))
                if action == "insert_row":
                    ws.insert_rows(index, amount)
                elif action == "delete_row":
                    ws.delete_rows(index, amount)
                elif action == "insert_column":
                    ws.insert_cols(index, amount)
                else:
                    ws.delete_cols(index, amount)
                applied += 1

        output_dir = os.path.dirname(output_path)
'''
text = replace_once(text, old, new, 'XLSX advanced actions')
p.write_text(text)


# ---------------------------------------------------------------------------
# Update Office schemas and skill prompt hints to advertise the real actions.
# ---------------------------------------------------------------------------
p = Path('python/navixmind/tools/__init__.py')
text = p.read_text()
text = text.replace(
    '["replace_text", "add_paragraph", "update_table_cell"]',
    '["replace_text", "add_paragraph", "update_table_cell", "add_heading", "add_page_break", "add_table", "add_image"]',
    1,
)
text = text.replace(
    '"replace_text: {old, new}. add_paragraph: {text, style?}. update_table_cell: {table, row, col, text}."',
    '"replace_text {old,new}; add_paragraph {text,style?}; update_table_cell {table,row,col,text}; add_heading {text,level?}; add_page_break {}; add_table {rows:[[...]]}; add_image {image_path,width_inches?}."',
    1,
)
text = text.replace(
    '["replace_text", "add_slide", "update_slide_text", "set_notes"]',
    '["replace_text", "add_slide", "update_slide_text", "set_notes", "add_textbox", "add_image", "delete_slide"]',
    1,
)
text = text.replace(
    '"replace_text: {old, new}. add_slide: {layout_index?, title?, content?}. update_slide_text: {slide, shape_name, text}. set_notes: {slide, text}."',
    '"replace_text {old,new}; add_slide {layout_index?,title?,content?}; update_slide_text {slide,shape_name,text}; set_notes {slide,text}; add_textbox {slide,text,left?,top?,width?,height?}; add_image {slide,image_path,left?,top?,width?,height?}; delete_slide {slide}."',
    1,
)
text = text.replace(
    '["set_cell", "set_formula", "add_row", "add_sheet", "delete_sheet"]',
    '["set_cell", "set_formula", "add_row", "add_sheet", "delete_sheet", "rename_sheet", "insert_row", "delete_row", "insert_column", "delete_column"]',
    1,
)
text = text.replace(
    '"set_cell: {sheet?, cell, value}. set_formula: {sheet?, cell, formula}. add_row: {sheet?, values: []}. add_sheet: {name}. delete_sheet: {name}."',
    '"set_cell {sheet?,cell,value}; set_formula {sheet?,cell,formula}; add_row {sheet?,values}; add_sheet/delete_sheet; rename_sheet {old_name,new_name}; insert/delete row/column {sheet?,index,amount?}."',
    1,
)
p.write_text(text)


# ---------------------------------------------------------------------------
# Calendar: add update, which is supported by the same calendar.events scope.
# Gmail stays intentionally read-only because the app currently requests only
# gmail.readonly; exposing send without changing OAuth would be a fake tool.
# ---------------------------------------------------------------------------
p = Path('python/navixmind/tools/google_api.py')
text = p.read_text()
text = replace_once(
    text,
    '        elif action == "delete":\n            return _delete_event(base_url, headers, event_id)\n',
    '        elif action == "delete":\n            return _delete_event(base_url, headers, event_id)\n'
    '        elif action == "update":\n            return _update_event(base_url, headers, event_id, event)\n',
    'calendar update dispatch',
)
insert = r'''

def _update_event(base_url: str, headers: dict, event_id: Optional[str], event: Optional[dict]) -> dict:
    """Update selected fields of one Calendar event."""
    if not event_id:
        raise ToolError("event_id required for update action")
    if not event:
        raise ToolError("event details required for update action")
    body = {}
    if event.get("title") is not None:
        body["summary"] = event["title"]
    if event.get("description") is not None:
        body["description"] = event["description"]
    if event.get("location") is not None:
        body["location"] = event["location"]
    if event.get("start") is not None:
        body["start"] = {"dateTime": event["start"], "timeZone": "UTC"}
    if event.get("end") is not None:
        body["end"] = {"dateTime": event["end"], "timeZone": "UTC"}
    if not body:
        raise ToolError("No supported event fields supplied for update")
    response = requests.patch(
        f"{base_url}/calendars/primary/events/{event_id}",
        headers=headers,
        json=body,
        timeout=30,
    )
    response.raise_for_status()
    updated = response.json()
    return {"success": True, "event_id": updated.get("id"), "updated": True}
'''
text = replace_once(text, '\n\ndef _delete_event(', insert + '\n\ndef _delete_event(', 'calendar update implementation')
p.write_text(text)

p = Path('python/navixmind/tools/__init__.py')
text = p.read_text()
# Update the first calendar schema and the dynamic v7 prompt hint.
text = text.replace('["list", "create", "delete"]', '["list", "create", "delete", "update"]', 1)
text = text.replace(
    '"google_calendar": "google_calendar(action, date_range?, event?, event_id?) ; action=list|create|delete",',
    '"google_calendar": "google_calendar(action, date_range?, event?, event_id?) ; action=list|create|delete|update",',
    1,
)
p.write_text(text)

p = Path('lib/core/models/tool_skill.dart')
text = p.read_text()
text = text.replace(
    "capabilities: ['今日/本周/日期范围查询', '创建日程', '删除日程', '标题/时间/描述/地点'],",
    "capabilities: ['今日/本周/日期范围查询', '创建日程', '更新日程', '删除日程', '标题/时间/描述/地点'],",
)
p.write_text(text)

print('Applied RastaCoder v7 expanded Office and Calendar actions')
