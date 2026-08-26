#!/usr/bin/env python3
from pathlib import Path


def patch(path: str, old: str, new: str, count: int = 1) -> None:
    p = Path(path)
    text = p.read_text(encoding='utf-8')
    actual = text.count(old)
    if actual != count:
        raise SystemExit(f'{path}: expected {count} occurrence(s), found {actual}: {old[:120]!r}')
    p.write_text(text.replace(old, new, count), encoding='utf-8')


def replace_between(path: str, start: str, end: str, replacement: str) -> None:
    p = Path(path)
    text = p.read_text(encoding='utf-8')
    si = text.find(start)
    ei = text.find(end, si + len(start))
    if si < 0 or ei < 0:
        raise SystemExit(f'{path}: block anchors not found: {start!r} -> {end!r}')
    p.write_text(text[:si] + replacement.rstrip() + '\n\n' + text[ei:], encoding='utf-8')


# ---------------------------------------------------------------------------
# Search execution: model supplies query only; provider settings come from
# private execution context configured by the user.
# ---------------------------------------------------------------------------
Path('python/navixmind/tools/search_tools.py').write_text(r'''"""Independent web-search provider tools for RastaCoder v9.

The on-device model only supplies the search query. Provider-specific tuning is
loaded from private per-provider settings in the execution context, keeping
small-model tool calls short and deterministic.
"""
from __future__ import annotations

from typing import Any, Dict, Iterable, Optional
import requests
from ..bridge import ToolError

_TIMEOUT = 45
_PROVIDER_LABELS = {
    "anysearch": "AnySearch",
    "exa": "Exa",
    "langsearch": "LangSearch",
    "tavily": "Tavily",
}


def _api_key(context: Optional[Dict[str, Any]], provider: str) -> str:
    keys = (context or {}).get("search_api_keys", {})
    key = keys.get(provider) if isinstance(keys, dict) else None
    if not isinstance(key, str) or not key.strip():
        label = _PROVIDER_LABELS.get(provider, provider)
        raise ToolError(f"{label} API Key 未配置。请在工具管理页面配置该搜索技能的 API Key。")
    return key.strip()


def _settings(context: Optional[Dict[str, Any]], provider: str) -> Dict[str, Any]:
    all_settings = (context or {}).get("search_settings", {})
    if not isinstance(all_settings, dict):
        return {}
    value = all_settings.get(provider, {})
    return dict(value) if isinstance(value, dict) else {}


def _as_int(value: Any, default: int, low: int = 1, high: int = 10) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = default
    return max(low, min(parsed, high))


def _as_bool(value: Any, default: bool) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        lowered = value.strip().lower()
        if lowered in {"true", "1", "yes", "on"}:
            return True
        if lowered in {"false", "0", "no", "off"}:
            return False
    return default


def _domains(value: Any) -> list[str]:
    if isinstance(value, str):
        values = value.split(',')
    elif isinstance(value, Iterable):
        values = list(value)
    else:
        return []
    return [str(x).strip() for x in values if str(x).strip()][:10]


def _raise_http(label: str, response: requests.Response) -> None:
    if response.ok:
        return
    detail = (response.text or "").strip().replace("\n", " ")[:1000]
    raise ToolError(f"{label} HTTP {response.status_code}: {detail}")


def _compact(value: Any, max_chars: int = 16000) -> str:
    text = str(value or "").strip()
    return text if len(text) <= max_chars else text[:max_chars] + "\n...[TRUNCATED]"


def anysearch_search(
    query: str,
    max_results: int = 5,
    domain: str = "",
    sub_domain: str = "",
    sub_domain_params: Optional[Dict[str, Any]] = None,
    _context: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    clean_query = (query or "").strip()
    if not clean_query:
        raise ToolError("AnySearch query is required.")
    cfg = _settings(_context, "anysearch")
    max_results = _as_int(cfg.get("max_results", max_results), 5)
    domain = str(cfg.get("domain", domain) or "").strip()
    sub_domain = str(cfg.get("sub_domain", sub_domain) or "").strip()
    cfg_sub = cfg.get("sub_domain_params")
    if isinstance(cfg_sub, dict):
        sub_domain_params = cfg_sub
    arguments: Dict[str, Any] = {"query": clean_query, "max_results": max_results}
    if domain:
        arguments["domain"] = domain.lower()
        if sub_domain:
            arguments["sub_domain"] = sub_domain
        if isinstance(sub_domain_params, dict) and sub_domain_params:
            arguments["sub_domain_params"] = sub_domain_params
    text = _anysearch_call("search", arguments, _context)
    return {"success": True, "provider": "anysearch", "operation": "search", "query": clean_query,
            "settings_used": {k: v for k, v in arguments.items() if k != "query"},
            "content": _compact(text, 12000)}


def anysearch_extract(url: str, _context: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    clean_url = (url or "").strip()
    if not clean_url.lower().startswith(("http://", "https://")):
        raise ToolError("AnySearch extract only supports http(s) URLs.")
    text = _anysearch_call("extract", {"url": clean_url}, _context)
    return {"success": True, "provider": "anysearch", "operation": "extract", "url": clean_url,
            "content": _compact(text, 16000)}


def anysearch_get_sub_domains(
    domain: str = "", domains: Optional[Iterable[str]] = None,
    _context: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    clean_domains = _domains(domains)
    arguments: Dict[str, Any] = {}
    if clean_domains:
        arguments["domains"] = clean_domains[:5]
    else:
        clean_domain = (domain or "").strip().lower()
        if not clean_domain:
            raise ToolError("AnySearch domain is required.")
        arguments["domain"] = clean_domain
    text = _anysearch_call("get_sub_domains", arguments, _context)
    return {"success": True, "provider": "anysearch", "operation": "get_sub_domains",
            "content": _compact(text, 10000)}


def _anysearch_call(tool_name: str, arguments: Dict[str, Any], context: Optional[Dict[str, Any]]) -> str:
    key = _api_key(context, "anysearch")
    payload = {"jsonrpc": "2.0", "id": 1, "method": "tools/call",
               "params": {"name": tool_name, "arguments": arguments}}
    try:
        response = requests.post(
            "https://api.anysearch.com/mcp", json=payload,
            headers={"Content-Type": "application/json", "X-Anysearch-Client": "rastacoder/9.0",
                     "Authorization": f"Bearer {key}"}, timeout=_TIMEOUT)
    except requests.RequestException as exc:
        raise ToolError(f"AnySearch request failed: {exc}") from exc
    _raise_http("AnySearch", response)
    try:
        body = response.json()
    except ValueError as exc:
        raise ToolError("AnySearch returned a non-JSON response.") from exc
    if isinstance(body, dict) and body.get("error"):
        error = body["error"]
        message = error.get("message") if isinstance(error, dict) else str(error)
        raise ToolError(f"AnySearch: {message}")
    result = body.get("result", {}) if isinstance(body, dict) else {}
    content = result.get("content", []) if isinstance(result, dict) else []
    if isinstance(content, list):
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                return str(item.get("text", ""))
    return _compact(result)


def exa_search(
    query: str, num_results: int = 5, topic: str = "general", search_type: str = "auto",
    start_published_date: str = "", include_domains: Optional[Iterable[str]] = None,
    exclude_domains: Optional[Iterable[str]] = None, include_text: bool = True,
    include_summary: bool = True, include_highlights: bool = False,
    _context: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    clean_query = (query or "").strip()
    if not clean_query:
        raise ToolError("Exa query is required.")
    cfg = _settings(_context, "exa")
    num_results = _as_int(cfg.get("num_results", num_results), 5)
    topic = str(cfg.get("topic", topic) or "general").strip()
    search_type = str(cfg.get("search_type", search_type) or "auto").strip()
    start_published_date = str(cfg.get("start_published_date", start_published_date) or "").strip()
    include_domains = _domains(cfg.get("include_domains", include_domains))
    exclude_domains = _domains(cfg.get("exclude_domains", exclude_domains))
    include_text = _as_bool(cfg.get("include_text", include_text), True)
    include_summary = _as_bool(cfg.get("include_summary", include_summary), True)
    include_highlights = _as_bool(cfg.get("include_highlights", include_highlights), False)
    key = _api_key(_context, "exa")
    contents: Dict[str, Any] = {}
    if include_text: contents["text"] = {"maxCharacters": 5000}
    if include_summary: contents["summary"] = True
    if include_highlights: contents["highlights"] = True
    payload: Dict[str, Any] = {"query": clean_query, "topic": topic, "type": search_type,
                               "numResults": num_results}
    if contents: payload["contents"] = contents
    if start_published_date: payload["startPublishedDate"] = start_published_date
    if include_domains: payload["includeDomains"] = include_domains
    if exclude_domains: payload["excludeDomains"] = exclude_domains
    try:
        response = requests.post("https://api.exa.ai/search", json=payload,
                                 headers={"Content-Type": "application/json", "x-api-key": key},
                                 timeout=_TIMEOUT)
    except requests.RequestException as exc:
        raise ToolError(f"Exa request failed: {exc}") from exc
    _raise_http("Exa", response)
    try:
        body = response.json()
    except ValueError as exc:
        raise ToolError("Exa returned a non-JSON response.") from exc
    return {"success": True, "provider": "exa", "query": clean_query,
            "settings_used": {"num_results": num_results, "topic": topic, "search_type": search_type,
                              "start_published_date": start_published_date,
                              "include_domains": include_domains, "exclude_domains": exclude_domains,
                              "include_text": include_text, "include_summary": include_summary,
                              "include_highlights": include_highlights},
            "results": body.get("results", []) if isinstance(body, dict) else [],
            "autoprompt_string": body.get("autopromptString") if isinstance(body, dict) else None}


def langsearch_search(
    query: str, count: int = 5, freshness: str = "noLimit", summary: bool = True,
    _context: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    clean_query = (query or "").strip()
    if not clean_query:
        raise ToolError("LangSearch query is required.")
    cfg = _settings(_context, "langsearch")
    count = _as_int(cfg.get("count", count), 5)
    freshness = str(cfg.get("freshness", freshness) or "noLimit").strip()
    summary = _as_bool(cfg.get("summary", summary), True)
    key = _api_key(_context, "langsearch")
    payload = {"query": clean_query, "freshness": freshness, "summary": summary, "count": count}
    try:
        response = requests.post("https://api.langsearch.com/v1/web-search", json=payload,
                                 headers={"Content-Type": "application/json", "Authorization": f"Bearer {key}"},
                                 timeout=_TIMEOUT)
    except requests.RequestException as exc:
        raise ToolError(f"LangSearch request failed: {exc}") from exc
    _raise_http("LangSearch", response)
    try:
        body = response.json()
    except ValueError as exc:
        raise ToolError("LangSearch returned a non-JSON response.") from exc
    if isinstance(body, dict) and body.get("code") not in (None, 200):
        raise ToolError(f"LangSearch: {body.get('message') or body.get('code')}")
    data = body.get("data", {}) if isinstance(body, dict) else {}
    pages = data.get("webPages", {}) if isinstance(data, dict) else {}
    return {"success": True, "provider": "langsearch", "query": clean_query,
            "settings_used": {"count": count, "freshness": freshness, "summary": summary},
            "results": pages.get("value", []) if isinstance(pages, dict) else []}


def tavily_search(
    query: str, max_results: int = 5, topic: str = "general", search_depth: str = "basic",
    include_answer: bool = True, time_range: str = "", include_domains: Optional[Iterable[str]] = None,
    exclude_domains: Optional[Iterable[str]] = None, include_raw_content: bool = False,
    _context: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    clean_query = (query or "").strip()
    if not clean_query:
        raise ToolError("Tavily query is required.")
    cfg = _settings(_context, "tavily")
    max_results = _as_int(cfg.get("max_results", max_results), 5)
    topic = str(cfg.get("topic", topic) or "general").strip()
    search_depth = str(cfg.get("search_depth", search_depth) or "basic").strip()
    include_answer = _as_bool(cfg.get("include_answer", include_answer), True)
    time_range = str(cfg.get("time_range", time_range) or "").strip()
    include_domains = _domains(cfg.get("include_domains", include_domains))
    exclude_domains = _domains(cfg.get("exclude_domains", exclude_domains))
    include_raw_content = _as_bool(cfg.get("include_raw_content", include_raw_content), False)
    key = _api_key(_context, "tavily")
    payload: Dict[str, Any] = {"query": clean_query, "topic": topic, "search_depth": search_depth,
                               "max_results": max_results, "include_answer": include_answer,
                               "include_raw_content": include_raw_content}
    if time_range: payload["time_range"] = time_range
    if include_domains: payload["include_domains"] = include_domains
    if exclude_domains: payload["exclude_domains"] = exclude_domains
    try:
        response = requests.post("https://api.tavily.com/search", json=payload,
                                 headers={"Content-Type": "application/json", "Authorization": f"Bearer {key}"},
                                 timeout=_TIMEOUT)
    except requests.RequestException as exc:
        raise ToolError(f"Tavily request failed: {exc}") from exc
    _raise_http("Tavily", response)
    try:
        body = response.json()
    except ValueError as exc:
        raise ToolError("Tavily returned a non-JSON response.") from exc
    return {"success": True, "provider": "tavily", "query": clean_query,
            "settings_used": {"max_results": max_results, "topic": topic, "search_depth": search_depth,
                              "include_answer": include_answer, "time_range": time_range,
                              "include_domains": include_domains, "exclude_domains": exclude_domains,
                              "include_raw_content": include_raw_content},
            "answer": body.get("answer") if isinstance(body, dict) else None,
            "results": body.get("results", []) if isinstance(body, dict) else []}
''', encoding='utf-8')


# ---------------------------------------------------------------------------
# Office mutation tools: transactional validation + actual-change accounting +
# reopen verification. A silent no-op is an error, never success.
# ---------------------------------------------------------------------------
DOCX = r'''def modify_docx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify a DOCX transactionally and verify the saved mutation."""
    from docx import Document
    validate_file_for_processing(input_path, 'document')
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_docx requires at least one operation")
    try:
        doc = Document(input_path)
        applied = 0
        checks = []
        before_paragraphs = len(doc.paragraphs)
        before_tables = len(doc.tables)
        before_images = len(doc.inline_shapes)
        for op in operations:
            if not isinstance(op, dict):
                raise ToolError("Each DOCX operation must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params", {})
            if not isinstance(params, dict):
                raise ToolError(f"{action or 'DOCX operation'} params must be an object")
            if action == "replace_text":
                old, new = str(params.get("old") or ""), str(params.get("new") or "")
                if not old:
                    raise ToolError("replace_text requires non-empty params.old")
                hits = 0
                containers = list(doc.paragraphs)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            containers.extend(cell.paragraphs)
                for para in containers:
                    if old not in para.text:
                        continue
                    run_hits = 0
                    for run in para.runs:
                        if old in run.text:
                            run_hits += run.text.count(old)
                            run.text = run.text.replace(old, new)
                    if run_hits == 0:
                        count = para.text.count(old)
                        para.text = para.text.replace(old, new)
                        run_hits = count
                    hits += run_hits
                if hits == 0:
                    raise ToolError(f"replace_text made no change; text not found: {old!r}")
                applied += 1
                checks.append(("contains", new))
            elif action == "add_paragraph":
                text = params.get("text")
                if text is None or not str(text).strip():
                    raise ToolError("add_paragraph requires non-empty params.text")
                style = params.get("style")
                doc.add_paragraph(str(text), style=style)
                applied += 1
                checks.append(("last_paragraph", str(text)))
            elif action == "update_table_cell":
                try:
                    ti, ri, ci = int(params.get("table", 0)), int(params.get("row", 0)), int(params.get("col", 0))
                except (TypeError, ValueError):
                    raise ToolError("update_table_cell table/row/col must be integers")
                if ti < 0 or ti >= len(doc.tables):
                    raise ToolError(f"DOCX table index out of range: {ti}")
                table = doc.tables[ti]
                if ri < 0 or ri >= len(table.rows) or ci < 0 or ci >= len(table.rows[ri].cells):
                    raise ToolError(f"DOCX table cell out of range: table={ti}, row={ri}, col={ci}")
                value = str(params.get("text", ""))
                table.rows[ri].cells[ci].text = value
                applied += 1
                checks.append(("table_cell", (ti, ri, ci, value)))
            elif action == "add_heading":
                text = str(params.get("text") or "")
                if not text.strip(): raise ToolError("add_heading requires non-empty params.text")
                level = int(params.get("level", 1))
                if level < 0 or level > 9: raise ToolError("add_heading level must be 0..9")
                doc.add_heading(text, level=level)
                applied += 1
                checks.append(("contains", text))
            elif action == "add_page_break":
                doc.add_page_break(); applied += 1; checks.append(("paragraph_count_gt", before_paragraphs))
            elif action == "add_table":
                rows = params.get("rows", []) or []
                if not isinstance(rows, list) or not rows:
                    raise ToolError("add_table requires params.rows as a non-empty 2D list")
                col_count = max(len(row) if isinstance(row, list) else 1 for row in rows)
                table = doc.add_table(rows=len(rows), cols=col_count)
                for r, row in enumerate(rows):
                    values = row if isinstance(row, list) else [row]
                    for c, value in enumerate(values): table.cell(r, c).text = str(value)
                applied += 1; checks.append(("table_count_gt", before_tables))
            elif action == "add_image":
                image_path = params.get("image_path")
                if not image_path or not os.path.isfile(image_path):
                    raise ToolError(f"add_image image not found: {image_path}")
                width_inches = params.get("width_inches")
                if width_inches is None: doc.add_picture(image_path)
                else:
                    from docx.shared import Inches
                    doc.add_picture(image_path, width=Inches(float(width_inches)))
                applied += 1; checks.append(("image_count_gt", before_images))
            else:
                raise ToolError(f"Unknown DOCX modification action: {action}")
        output_dir = os.path.dirname(output_path)
        if output_dir: os.makedirs(output_dir, exist_ok=True)
        doc.save(output_path)
        if not os.path.isfile(output_path) or os.path.getsize(output_path) <= 0:
            raise ToolError("DOCX save verification failed: output file missing or empty")
        verified = Document(output_path)
        all_text = "\n".join(p.text for p in verified.paragraphs)
        all_text += "\n" + "\n".join(cell.text for t in verified.tables for row in t.rows for cell in row.cells)
        for kind, expected in checks:
            if kind == "contains" and expected and str(expected) not in all_text:
                raise ToolError(f"DOCX post-verification failed; expected text missing: {expected!r}")
            if kind == "last_paragraph" and (not verified.paragraphs or verified.paragraphs[-1].text != expected):
                raise ToolError("DOCX post-verification failed; appended paragraph is not at document end")
            if kind == "table_cell":
                ti, ri, ci, value = expected
                if verified.tables[ti].rows[ri].cells[ci].text != value:
                    raise ToolError("DOCX post-verification failed for updated table cell")
            if kind == "paragraph_count_gt" and len(verified.paragraphs) <= int(expected):
                raise ToolError("DOCX post-verification failed for page break")
            if kind == "table_count_gt" and len(verified.tables) <= int(expected):
                raise ToolError("DOCX post-verification failed for added table")
            if kind == "image_count_gt" and len(verified.inline_shapes) <= int(expected):
                raise ToolError("DOCX post-verification failed for added image")
        return {"output_path": output_path, "success": True, "operations_applied": applied,
                "verified": True, "size_bytes": os.path.getsize(output_path)}
    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to modify DOCX: {str(e)}")
'''
replace_between('python/navixmind/tools/documents.py', 'def modify_docx(', '# ---------------------------------------------------------------------------\n# PPTX read/write', DOCX)

PPTX = r'''def modify_pptx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify a PPTX transactionally and verify the saved mutation."""
    from pptx import Presentation
    validate_file_for_processing(input_path, 'document')
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_pptx requires at least one operation")
    try:
        prs = Presentation(input_path)
        applied = 0
        checks = []
        before_slides = len(prs.slides)
        for op in operations:
            if not isinstance(op, dict): raise ToolError("Each PPTX operation must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params", {})
            if not isinstance(params, dict): raise ToolError(f"{action or 'PPTX operation'} params must be an object")
            if action == "replace_text":
                old, new = str(params.get("old") or ""), str(params.get("new") or "")
                if not old: raise ToolError("replace_text requires non-empty params.old")
                hits = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame: continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if old in run.text:
                                    hits += run.text.count(old); run.text = run.text.replace(old, new)
                if hits == 0: raise ToolError(f"replace_text made no change; text not found: {old!r}")
                applied += 1; checks.append(("contains", new))
            elif action == "add_slide":
                layout_idx = int(params.get("layout_index", 1))
                if layout_idx < 0 or layout_idx >= len(prs.slide_layouts):
                    raise ToolError(f"PPTX layout index out of range: {layout_idx}")
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                title, content = str(params.get("title") or ""), str(params.get("content") or "")
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 0 and title: ph.text = title
                    elif ph.placeholder_format.idx == 1 and content: ph.text = content
                applied += 1; checks.append(("slide_count", before_slides + 1))
                if title: checks.append(("contains", title))
                if content: checks.append(("contains", content))
            elif action == "update_slide_text":
                slide_num = int(params.get("slide", 1)) - 1
                shape_name, value = str(params.get("shape_name") or ""), str(params.get("text") or "")
                if not (0 <= slide_num < len(prs.slides)): raise ToolError(f"PPTX slide out of range: {slide_num + 1}")
                match = next((s for s in prs.slides[slide_num].shapes if s.name == shape_name and s.has_text_frame), None)
                if match is None: raise ToolError(f"PPTX text shape not found: {shape_name!r}")
                match.text_frame.paragraphs[0].text = value
                applied += 1; checks.append(("contains", value))
            elif action == "set_notes":
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)): raise ToolError(f"PPTX slide out of range: {slide_num + 1}")
                value = str(params.get("text") or "")
                prs.slides[slide_num].notes_slide.notes_text_frame.text = value
                applied += 1; checks.append(("notes", (slide_num, value)))
            elif action == "add_textbox":
                from pptx.util import Inches
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)): raise ToolError(f"PPTX slide out of range: {slide_num + 1}")
                value = str(params.get("text") or "")
                box = prs.slides[slide_num].shapes.add_textbox(Inches(float(params.get("left", 1))), Inches(float(params.get("top", 1))), Inches(float(params.get("width", 6))), Inches(float(params.get("height", 1))))
                box.text_frame.text = value
                applied += 1; checks.append(("contains", value))
            elif action == "add_image":
                from pptx.util import Inches
                slide_num = int(params.get("slide", 1)) - 1
                image_path = params.get("image_path")
                if not (0 <= slide_num < len(prs.slides)): raise ToolError(f"PPTX slide out of range: {slide_num + 1}")
                if not image_path or not os.path.isfile(image_path): raise ToolError(f"add_image image not found: {image_path}")
                before_shapes = len(prs.slides[slide_num].shapes)
                kwargs = {"left": Inches(float(params.get("left", 1))), "top": Inches(float(params.get("top", 1)))}
                if params.get("width") is not None: kwargs["width"] = Inches(float(params["width"]))
                if params.get("height") is not None: kwargs["height"] = Inches(float(params["height"]))
                prs.slides[slide_num].shapes.add_picture(image_path, **kwargs)
                applied += 1; checks.append(("shape_count_gt", (slide_num, before_shapes)))
            elif action == "delete_slide":
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)): raise ToolError(f"PPTX slide out of range: {slide_num + 1}")
                slide_id = prs.slides._sldIdLst[slide_num]
                prs.part.drop_rel(slide_id.rId); prs.slides._sldIdLst.remove(slide_id)
                applied += 1; checks.append(("slide_count", len(prs.slides)))
            else:
                raise ToolError(f"Unknown PPTX modification action: {action}")
        output_dir = os.path.dirname(output_path)
        if output_dir: os.makedirs(output_dir, exist_ok=True)
        prs.save(output_path)
        if not os.path.isfile(output_path) or os.path.getsize(output_path) <= 0:
            raise ToolError("PPTX save verification failed: output file missing or empty")
        verified = Presentation(output_path)
        text = "\n".join(shape.text for slide in verified.slides for shape in slide.shapes if getattr(shape, 'has_text_frame', False))
        for kind, expected in checks:
            if kind == "contains" and expected and str(expected) not in text: raise ToolError(f"PPTX post-verification failed; text missing: {expected!r}")
            if kind == "slide_count" and len(verified.slides) != int(expected): raise ToolError("PPTX post-verification failed for slide count")
            if kind == "notes":
                si, value = expected
                if verified.slides[si].notes_slide.notes_text_frame.text != value: raise ToolError("PPTX post-verification failed for notes")
            if kind == "shape_count_gt":
                si, before = expected
                if len(verified.slides[si].shapes) <= before: raise ToolError("PPTX post-verification failed for added image")
        return {"output_path": output_path, "success": True, "operations_applied": applied,
                "verified": True, "size_bytes": os.path.getsize(output_path)}
    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to modify PPTX: {str(e)}")
'''
replace_between('python/navixmind/tools/documents.py', 'def modify_pptx(', '# ---------------------------------------------------------------------------\n# XLSX read/write', PPTX)

XLSX = r'''def modify_xlsx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify an XLSX transactionally and verify the saved mutation."""
    from openpyxl import load_workbook
    validate_file_for_processing(input_path, 'document')
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_xlsx requires at least one operation")
    try:
        wb = load_workbook(input_path)
        applied = 0
        checks = []
        for op in operations:
            if not isinstance(op, dict): raise ToolError("Each XLSX operation must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params", {})
            if not isinstance(params, dict): raise ToolError(f"{action or 'XLSX operation'} params must be an object")
            if action in {"set_cell", "set_formula"}:
                sheet_name = str(params.get("sheet") or wb.sheetnames[0])
                if sheet_name not in wb.sheetnames: raise ToolError(f"Sheet '{sheet_name}' not found.")
                cell = str(params.get("cell") or "").strip()
                if not cell: raise ToolError(f"{action} requires params.cell")
                value = params.get("value") if action == "set_cell" else params.get("formula", "")
                ws = wb[sheet_name]
                if ws[cell].value == value: raise ToolError(f"{action} made no change; {sheet_name}!{cell} already has requested value")
                ws[cell] = value; applied += 1; checks.append(("cell", (sheet_name, cell, value)))
            elif action == "add_row":
                sheet_name = str(params.get("sheet") or wb.sheetnames[0])
                if sheet_name not in wb.sheetnames: raise ToolError(f"Sheet '{sheet_name}' not found.")
                values = params.get("values")
                if not isinstance(values, list) or not values: raise ToolError("add_row requires non-empty params.values")
                ws = wb[sheet_name]; before = ws.max_row; ws.append(values); applied += 1
                checks.append(("row_count_gt", (sheet_name, before)))
            elif action == "add_sheet":
                name = str(params.get("name") or "Sheet").strip()[:31]
                if not name: raise ToolError("add_sheet requires params.name")
                if name in wb.sheetnames: raise ToolError(f"add_sheet made no change; sheet already exists: {name}")
                wb.create_sheet(title=name); applied += 1; checks.append(("sheet_exists", name))
            elif action == "delete_sheet":
                name = str(params.get("name") or "").strip()
                if name not in wb.sheetnames: raise ToolError(f"delete_sheet target not found: {name}")
                if len(wb.sheetnames) <= 1: raise ToolError("delete_sheet cannot delete the last worksheet")
                del wb[name]; applied += 1; checks.append(("sheet_missing", name))
            elif action == "rename_sheet":
                old_name = str(params.get("old_name") or params.get("sheet") or "").strip()
                new_name = str(params.get("new_name") or "").strip()[:31]
                if old_name not in wb.sheetnames: raise ToolError(f"rename_sheet source not found: {old_name}")
                if not new_name: raise ToolError("rename_sheet requires new_name")
                if new_name == old_name: raise ToolError("rename_sheet made no change; names are identical")
                if new_name in wb.sheetnames: raise ToolError(f"rename_sheet destination already exists: {new_name}")
                wb[old_name].title = new_name; applied += 1; checks += [("sheet_missing", old_name), ("sheet_exists", new_name)]
            elif action in {"insert_row", "delete_row", "insert_column", "delete_column"}:
                sheet_name = str(params.get("sheet") or wb.active.title)
                if sheet_name not in wb.sheetnames: raise ToolError(f"Sheet not found: {sheet_name}")
                ws = wb[sheet_name]; index = int(params.get("index", 1)); amount = int(params.get("amount", 1))
                if index < 1 or amount < 1: raise ToolError(f"{action} index and amount must be positive")
                if action == "delete_row" and index > ws.max_row: raise ToolError(f"delete_row starts beyond used rows: {index} > {ws.max_row}")
                if action == "delete_column" and index > ws.max_column: raise ToolError(f"delete_column starts beyond used columns: {index} > {ws.max_column}")
                before_rows, before_cols = ws.max_row, ws.max_column
                if action == "insert_row": ws.insert_rows(index, amount)
                elif action == "delete_row": ws.delete_rows(index, amount)
                elif action == "insert_column": ws.insert_cols(index, amount)
                else: ws.delete_cols(index, amount)
                applied += 1; checks.append(("dimensions_changed", (sheet_name, before_rows, before_cols, action)))
            else:
                raise ToolError(f"Unknown XLSX modification action: {action}")
        output_dir = os.path.dirname(output_path)
        if output_dir: os.makedirs(output_dir, exist_ok=True)
        wb.save(output_path); wb.close()
        if not os.path.isfile(output_path) or os.path.getsize(output_path) <= 0:
            raise ToolError("XLSX save verification failed: output file missing or empty")
        verified = load_workbook(output_path, data_only=False)
        for kind, expected in checks:
            if kind == "cell":
                sheet, cell, value = expected
                if verified[sheet][cell].value != value: raise ToolError(f"XLSX post-verification failed for {sheet}!{cell}")
            if kind == "row_count_gt":
                sheet, before = expected
                if verified[sheet].max_row <= before: raise ToolError("XLSX post-verification failed for added row")
            if kind == "sheet_exists" and expected not in verified.sheetnames: raise ToolError(f"XLSX post-verification failed; sheet missing: {expected}")
            if kind == "sheet_missing" and expected in verified.sheetnames: raise ToolError(f"XLSX post-verification failed; sheet still exists: {expected}")
            if kind == "dimensions_changed":
                sheet, br, bc, action = expected
                ws = verified[sheet]
                if action in {"insert_row", "delete_row"} and ws.max_row == br: raise ToolError(f"XLSX post-verification failed for {action}")
                if action in {"insert_column", "delete_column"} and ws.max_column == bc: raise ToolError(f"XLSX post-verification failed for {action}")
        verified.close()
        return {"output_path": output_path, "success": True, "operations_applied": applied,
                "verified": True, "size_bytes": os.path.getsize(output_path)}
    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to modify XLSX: {str(e)}")
'''
replace_between('python/navixmind/tools/documents.py', 'def modify_xlsx(', '# ---------------------------------------------------------------------------\n#', XLSX)


# ---------------------------------------------------------------------------
# Media download: remove an undefined fallback variable and require a real,
# non-empty downloaded file.
# ---------------------------------------------------------------------------
patch('python/navixmind/tools/media.py',
      "            request_headers = best_format.get('http_headers') or headers\n",
      "            request_headers = best_format.get('http_headers') or info.get('http_headers') or {}\n")
patch('python/navixmind/tools/media.py',
      "            return {\n                \"title\": title,\n",
      "            if not os.path.isfile(final_path) or os.path.getsize(final_path) <= 0:\n"
      "                raise ToolError('Media download verification failed: output file is missing or empty')\n"
      "            return {\n                \"title\": title,\n")


# ---------------------------------------------------------------------------
# Compatibility layer: common natural-language aliases, scalar speed factor,
# simple search query repair, private settings boundary, and in-place Office
# updates for files already living in the output workspace.
# ---------------------------------------------------------------------------
patch('python/navixmind/tools/compat.py',
'''                "transcode": "convert", "conversion": "convert",
                "convert_audio": "extract_audio", "audio_convert": "extract_audio", "audio_conversion": "extract_audio",
''',
'''                "transcode": "convert", "conversion": "convert",
                "convert_audio": "extract_audio", "audio_convert": "extract_audio", "audio_conversion": "extract_audio",
                "speed": "speed", "tempo": "speed", "playback_speed": "speed", "rate": "speed",
                "speed_up": "speed", "speedup": "speed", "slow_down": "speed", "slowdown": "speed",
''')
patch('python/navixmind/tools/compat.py',
'''        params = args.get("params")
        if not isinstance(params, dict):
            params = {}
        for key in ("start", "end", "duration", "width", "height", "x", "y", "vf", "af", "video_filter", "audio_filter", "format", "bitrate", "timestamp", "codec", "quality", "args"):
''',
'''        raw_params = args.get("params")
        if args.get("operation") == "speed" and raw_params is not None and not isinstance(raw_params, dict):
            params = {"factor": raw_params}
            notes.append("params:scalar->params.factor")
        else:
            params = raw_params if isinstance(raw_params, dict) else {}
        for key in ("start", "end", "duration", "width", "height", "x", "y", "vf", "af", "video_filter", "audio_filter", "format", "bitrate", "timestamp", "codec", "quality", "args", "factor", "speed", "rate"):
''')
patch('python/navixmind/tools/compat.py',
'''            if key in args and key not in params:
                params[key] = args.pop(key)
                notes.append(f"top-level:{key}->params.{key}")
        if "codec" in params:
''',
'''            if key in args and key not in params:
                params[key] = args.pop(key)
                notes.append(f"top-level:{key}->params.{key}")
        if args.get("operation") == "speed":
            if "factor" not in params:
                for alias in ("speed", "rate"):
                    if alias in params:
                        params["factor"] = params.pop(alias)
                        notes.append(f"params.{alias}->params.factor")
                        break
            if "factor" in params:
                try:
                    params["factor"] = float(str(params["factor"]).rstrip("xX"))
                except (TypeError, ValueError):
                    pass
        if "codec" in params:
''')

# Search aliases must be repaired before generic keys are stripped.
patch('python/navixmind/tools/compat.py',
'''    # Generic free-form keys are compatibility scaffolding, never canonical
    # tool arguments. Remove them after extracting deterministic information.
    for key in ("param", "request", "instruction", "command", "query"):
        if key == "query" and name == "gmail":
            continue
''',
'''    search_tools = {"anysearch_search", "exa_search", "langsearch_search", "tavily_search"}
    if name in search_tools:
        _move_alias(args, "query", ["q", "keyword", "keywords", "search_query", "topic"], notes)
        # For an on-device/manual-Skill request, provider tuning belongs to the
        # user's private settings, never to the model-visible call.
        if isinstance(context, dict) and context.get("_allowed_tools") is not None:
            for key in list(args.keys()):
                if key not in {"query"}:
                    args.pop(key, None)
                    notes.append(f"search_setting_removed:{key}")

    # Generic free-form keys are compatibility scaffolding, never canonical
    # tool arguments. Remove them after extracting deterministic information.
    for key in ("param", "request", "instruction", "command", "query"):
        if key == "query" and name in ({"gmail"} | search_tools):
            continue
''')

# Office default output: same visible output file when it is already in the
# app output workspace; external/internal attachments become a visible copy
# with the same basename in output.
patch('python/navixmind/tools/compat.py',
'''    if name in {"modify_docx", "modify_pptx", "modify_xlsx"} and not args.get("output_path") and isinstance(args.get("input_path"), str):
        ext = {"modify_docx": "docx", "modify_pptx": "pptx", "modify_xlsx": "xlsx"}[name]
        args["output_path"] = _derive_output(args["input_path"], "modified", ext)
        notes.append("derived:output_path")
''',
'''    if name in {"modify_docx", "modify_pptx", "modify_xlsx"} and not args.get("output_path") and isinstance(args.get("input_path"), str):
        source = args["input_path"]
        file_map = context.get("_file_map", {}) if isinstance(context, dict) else {}
        resolved = file_map.get(os.path.basename(source)) if isinstance(file_map, dict) else None
        if isinstance(resolved, str) and resolved:
            source = resolved
        output_root = context.get("output_dir") if isinstance(context, dict) else None
        if isinstance(output_root, str) and os.path.isabs(source):
            try:
                inside_output = os.path.commonpath([os.path.realpath(source), os.path.realpath(output_root)]) == os.path.realpath(output_root)
            except ValueError:
                inside_output = False
        else:
            inside_output = False
        args["output_path"] = source if inside_output else os.path.basename(source)
        notes.append("default:office_output=in_place_or_visible_copy")
''')

# Avoid semantic fake success from a same-format audio "convert" retry.
patch('python/navixmind/tools/compat.py',
'''        if target_audio and (in_ext in AUDIO_EXTS or out_ext in AUDIO_EXTS):
            if args.get("operation") in (None, "convert", "extract_audio"):
                if args.get("operation") != "extract_audio":
                    notes.append(f"operation:{args.get('operation')}->extract_audio")
                args["operation"] = "extract_audio"
                params["format"] = target_audio
''',
'''        if target_audio and (in_ext in AUDIO_EXTS or out_ext in AUDIO_EXTS):
            op = args.get("operation")
            should_extract = op == "extract_audio" or (op in (None, "convert") and (in_ext not in AUDIO_EXTS or target_audio != in_ext))
            if should_extract:
                if op != "extract_audio": notes.append(f"operation:{op}->extract_audio")
                args["operation"] = "extract_audio"
                params["format"] = target_audio
''')


# ---------------------------------------------------------------------------
# FFmpeg: first-class speed operation and semantic duration verification.
# ---------------------------------------------------------------------------
patch('lib/core/services/native_tool_executor.dart',
'''    // Build FFmpeg command based on operation
    // -y flag to auto-overwrite existing files
    String command;
    switch (operation) {
''',
'''    // Build FFmpeg command based on operation
    // -y flag to auto-overwrite existing files
    String command;
    String buildAtempoFilter(double factor) {
      if (factor <= 0) throw ArgumentError('speed factor must be greater than 0');
      var remaining = factor;
      final pieces = <double>[];
      while (remaining > 2.0) { pieces.add(2.0); remaining /= 2.0; }
      while (remaining < 0.5) { pieces.add(0.5); remaining /= 0.5; }
      pieces.add(remaining);
      return pieces.map((v) => 'atempo=${v.toStringAsFixed(6)}').join(',');
    }
    switch (operation) {
      case 'speed':
        final rawFactor = params['factor'] ?? params['speed'] ?? params['rate'];
        final factor = rawFactor is num ? rawFactor.toDouble() : double.tryParse(rawFactor?.toString().replaceAll('x', '') ?? '');
        if (factor == null || factor <= 0 || factor > 16) {
          throw ArgumentError('speed requires params.factor between 0 and 16, e.g. 1.5');
        }
        final atempo = buildAtempoFilter(factor);
        final dot = inputPath.lastIndexOf('.');
        final ext = dot >= 0 ? inputPath.substring(dot).toLowerCase() : '';
        const audioExts = <String>{'.mp3', '.wav', '.m4a', '.aac', '.flac', '.ogg', '.opus', '.wma', '.amr'};
        if (audioExts.contains(ext)) {
          command = '-y -i "$inputPath" -vn -filter:a "$atempo" "$outputPath"';
        } else {
          command = '-y -i "$inputPath" -filter:v "setpts=PTS/$factor" -filter:a "$atempo" -c:v libx264 -pix_fmt yuv420p -c:a aac "$outputPath"';
        }
        break;
''')
patch('lib/core/services/native_tool_executor.dart',
'''      final outputFile = File(outputPath);
      final outputSize = await outputFile.length();

      // Probe output file for actual media duration
''',
'''      final outputFile = File(outputPath);
      if (!await outputFile.exists()) {
        throw Exception('FFmpeg reported success but output file does not exist: $outputPath');
      }
      final outputSize = await outputFile.length();
      if (outputSize <= 0) {
        throw Exception('FFmpeg reported success but output file is empty: $outputPath');
      }

      // Probe output file for actual media duration
''')
patch('lib/core/services/native_tool_executor.dart',
'''      if (mediaDuration != null) {
        result['media_duration_seconds'] = mediaDuration;
      }
      return result;
''',
'''      if (mediaDuration != null) {
        result['media_duration_seconds'] = mediaDuration;
      }
      if (operation == 'speed' && mediaDuration != null) {
        final rawFactor = params['factor'] ?? params['speed'] ?? params['rate'];
        final factor = rawFactor is num ? rawFactor.toDouble() : double.tryParse(rawFactor?.toString().replaceAll('x', '') ?? '');
        try {
          final inputProbe = await FFprobeKit.getMediaInformation(inputPath);
          final inputDuration = double.tryParse(inputProbe.getMediaInformation()?.getDuration() ?? '');
          if (factor != null && inputDuration != null && inputDuration > 0) {
            final expected = inputDuration / factor;
            final tolerance = expected * 0.12 + 0.25;
            if ((mediaDuration - expected).abs() > tolerance) {
              try { await outputFile.delete(); } catch (_) {}
              throw Exception('Speed verification failed: expected duration about ${expected.toStringAsFixed(3)}s, got ${mediaDuration.toStringAsFixed(3)}s');
            }
            result['input_duration_seconds'] = inputDuration;
            result['speed_factor'] = factor;
            result['duration_verified'] = true;
          }
        } catch (e) {
          if (e.toString().contains('Speed verification failed')) rethrow;
        }
      }
      return result;
''')


# ---------------------------------------------------------------------------
# Tool registry: simplify local schemas, add speed enum, generic artifact
# verification, and concise prompt hints. Cloud schemas retain richer options.
# ---------------------------------------------------------------------------
registry = Path('python/navixmind/tools/__init__.py')
text = registry.read_text(encoding='utf-8')
anchor = '# Skill IDs are UI-only. They are deliberately never shown to the model.\n'
if anchor not in text:
    raise SystemExit('__init__.py: local skill anchor missing')
hardening_block = r'''# RASTACODER_V9_SYSTEMIC_HARDENING
# Keep the local 3B-4B contract short. Compatibility/context inference fills
# deterministic paths and private provider settings before strict execution.
for _catalog in (TOOLS_SCHEMA, OFFLINE_TOOLS_SCHEMA):
    for _tool in _catalog:
        if _tool.get("name") == "ffmpeg_process":
            _op = _tool.get("input_schema", {}).get("properties", {}).get("operation", {})
            _enum = _op.get("enum", [])
            if "speed" not in _enum:
                _enum.append("speed")
            _op["description"] = "trim/crop/resize/filter/speed/custom/extract_audio/extract_frame/convert/concat/mix_audio/merge_av"

_search_local = {"anysearch_search", "exa_search", "langsearch_search", "tavily_search"}
for _tool in OFFLINE_TOOLS_SCHEMA:
    if _tool.get("name") in _search_local:
        _tool["description"] = "Search the web. Supply only the search query; result count/type/filter settings are configured by the user."
        _tool["input_schema"] = {
            "type": "object",
            "properties": {"query": {"type": "string", "description": "Search keywords/query"}},
            "required": ["query"],
            "additionalProperties": False,
        }
    if _tool.get("name") == "ffmpeg_process":
        _tool["description"] = "Process audio/video. Common operations include speed, trim, convert, resize, concat, mix and merge. output_path can be chosen automatically."
        _tool["input_schema"]["required"] = ["operation"]
        _tool["input_schema"]["properties"]["params"]["description"] = "speed {factor}; trim {start,end/duration}; resize {width,height}; filter {vf,af}; extract_audio {format,bitrate}; other operation-specific parameters."
    if _tool.get("name") in {"modify_docx", "modify_pptx", "modify_xlsx"}:
        _actions = {
            "modify_docx": ["replace_text", "add_paragraph", "update_table_cell", "add_heading", "add_page_break", "add_table", "add_image"],
            "modify_pptx": ["replace_text", "add_slide", "update_slide_text", "set_notes", "add_textbox", "add_image", "delete_slide"],
            "modify_xlsx": ["set_cell", "set_formula", "add_row", "add_sheet", "delete_sheet", "rename_sheet", "insert_row", "delete_row", "insert_column", "delete_column"],
        }[_tool["name"]]
        _tool["description"] = "Modify the existing file. Supply input_path, one action and params; output is updated in place when safe or saved as a visible copy."
        _tool["input_schema"] = {
            "type": "object",
            "properties": {
                "input_path": {"type": "string"},
                "action": {"type": "string", "enum": _actions},
                "params": {"type": "object"},
                "output_path": {"type": "string", "description": "Optional explicit destination"},
            },
            "required": ["action"],
        }
'''
text = text.replace(anchor, hardening_block + '\n\n' + anchor, 1)
registry.write_text(text, encoding='utf-8')

# Prompt hints are user/model-facing and must match the simplified contract.
p = registry
text = p.read_text(encoding='utf-8')
for old, new in {
    '"modify_docx": "modify_docx(input_path, output_path, operations)",': '"modify_docx": "modify_docx(input_path?, action, params?, output_path?) ; common action=add_paragraph for appending text",',
    '"modify_pptx": "modify_pptx(input_path, output_path, operations)",': '"modify_pptx": "modify_pptx(input_path?, action, params?, output_path?)",',
    '"modify_xlsx": "modify_xlsx(input_path, output_path, operations)",': '"modify_xlsx": "modify_xlsx(input_path?, action, params?, output_path?)",',
    '"ffmpeg_process": "ffmpeg_process(input_path?, input_paths?, output_path, operation, params?) ; operation=trim|crop|resize|filter|custom|extract_audio|extract_frame|convert|concat|mix_audio|merge_av",': '"ffmpeg_process": "ffmpeg_process(operation, input_path?, input_paths?, output_path?, params?) ; speed uses params={factor:1.5}; operation=trim|crop|resize|filter|speed|custom|extract_audio|extract_frame|convert|concat|mix_audio|merge_av",',
}.items():
    if old not in text: raise SystemExit(f'__init__.py prompt hint missing: {old[:80]}')
    text = text.replace(old, new, 1)
# Search hints varied between source generations; replace the whole known four-line run by prefix matching.
lines = text.splitlines()
for i, line in enumerate(lines):
    if line.strip().startswith('"anysearch_search":'):
        lines[i] = '    "anysearch_search": "anysearch_search(query)",'
    elif line.strip().startswith('"exa_search":'):
        lines[i] = '    "exa_search": "exa_search(query)",'
    elif line.strip().startswith('"langsearch_search":'):
        lines[i] = '    "langsearch_search": "langsearch_search(query)",'
    elif line.strip().startswith('"tavily_search":'):
        lines[i] = '    "tavily_search": "tavily_search(query)",'
p.write_text('\n'.join(lines) + '\n', encoding='utf-8')

# Generic output-artifact verification for every local/cloud tool returning files.
patch('python/navixmind/tools/__init__.py',
'''def execute_tool(
''',
'''def _verify_tool_result_artifacts(tool_name: str, result: Any) -> Any:
    """Reject success claims when declared output artifacts are absent/empty."""
    import os
    if not isinstance(result, dict) or result.get("success") is not True:
        return result
    paths = []
    if isinstance(result.get("output_path"), str): paths.append(result["output_path"])
    if isinstance(result.get("output_paths"), list): paths.extend(p for p in result["output_paths"] if isinstance(p, str))
    for path in paths:
        if not os.path.exists(path):
            raise ToolError(f"[TOOL_POST_VERIFY_ERROR] {tool_name} reported success but output is missing: {path}")
        if os.path.isfile(path) and os.path.getsize(path) <= 0:
            raise ToolError(f"[TOOL_POST_VERIFY_ERROR] {tool_name} reported success but output is empty: {path}")
    if tool_name in {"modify_docx", "modify_pptx", "modify_xlsx"} and int(result.get("operations_applied", 0)) <= 0:
        raise ToolError(f"[TOOL_POST_VERIFY_ERROR] {tool_name} reported success without applying any operation")
    return result


def execute_tool(
''')
patch('python/navixmind/tools/__init__.py',
'''    return tool_func(**args)
''',
'''    result = tool_func(**args)
    return _verify_tool_result_artifacts(tool_name, result)
''')


# ---------------------------------------------------------------------------
# PDF no-op protection + output validity.
# ---------------------------------------------------------------------------
patch('python/navixmind/tools/extended_tools.py',
'''    selected = _parse_pages(pages, total)
    if action == "delete_pages":
''',
'''    if action == "delete_pages" and pages is None:
        raise ToolError("delete_pages requires explicit pages; refusing to interpret omitted pages as delete all")
    if action == "rotate" and int(rotation) % 360 == 0:
        raise ToolError("rotate made no change; rotation must not be a multiple of 360")
    selected = _parse_pages(pages, total)
    if action == "delete_pages":
''')
patch('python/navixmind/tools/extended_tools.py',
'''        selected = [i for i in range(total) if i not in excluded]
    elif action not in {"extract_pages", "reorder", "rotate"}:
''',
'''        selected = [i for i in range(total) if i not in excluded]
        if not selected:
            raise ToolError("delete_pages would remove every page; refusing to create an empty PDF")
    elif action not in {"extract_pages", "reorder", "rotate"}:
''')


# ---------------------------------------------------------------------------
# Calendar success must contain the server-side identity that proves mutation.
# ---------------------------------------------------------------------------
patch('python/navixmind/tools/google_api.py',
'''    updated = response.json()
    return {"success": True, "event_id": updated.get("id"), "updated": True}
''',
'''    updated = response.json()
    returned_id = updated.get("id") if isinstance(updated, dict) else None
    if not returned_id or returned_id != event_id:
        raise ToolError("Calendar update verification failed: server did not return the requested event ID")
    return {"success": True, "event_id": returned_id, "updated": True, "verified": True}
''')
patch('python/navixmind/tools/google_api.py',
'''    created = response.json()
    return {
        "success": True,
        "event_id": created.get("id"),
''',
'''    created = response.json()
    if not isinstance(created, dict) or not created.get("id"):
        raise ToolError("Calendar create verification failed: server returned no event ID")
    return {
        "success": True,
        "verified": True,
        "event_id": created.get("id"),
''')

print('Applied RastaCoder v9 systemic core hardening')
