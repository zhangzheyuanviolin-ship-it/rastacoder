#!/usr/bin/env python3
from __future__ import annotations

import sys
import tempfile
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / 'python'))


def require(condition, message):
    if not condition:
        raise AssertionError(message)


def source(path):
    return (ROOT / path).read_text(encoding='utf-8')


from navixmind.tools import TOOLS_SCHEMA, _verify_tool_result
from navixmind.tools.extended_tools import extract_zip, image_compose
from navixmind.tools.documents import read_xlsx
from navixmind.tools.code_executor import execute_python, validate_code

names = [tool.get('name') for tool in TOOLS_SCHEMA]
require(len(set(names)) == 37, f'Expected 37 canonical local functions, got {len(set(names))}')

# 1: python-pptx notes API + Android resource packaging.
from pptx import Presentation
with tempfile.TemporaryDirectory() as td:
    pptx_path = Path(td) / 'notes.pptx'
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.notes_slide.notes_text_frame.text = 'V15 speaker notes'
    prs.save(pptx_path)
    check = Presentation(pptx_path)
    require(check.slides[0].notes_slide.notes_text_frame.text == 'V15 speaker notes', 'PPTX notes roundtrip failed')

gradle = source('android/app/build.gradle')
root_gradle = source('android/build.gradle')
require('version "3.13"' in gradle, 'Android Python 3.13 runtime missing')
require('extractPackages("pptx")' in gradle, 'pptx package extraction missing')
require('curl-cffi==0.16.2' in gradle, 'curl-cffi Android runtime missing')
require('com.chaquo.python:gradle:16.1.0' in root_gradle, 'Chaquopy 16.1 runtime missing')

# 2: legal zero-byte ZIP members survive extraction and common postconditions.
with tempfile.TemporaryDirectory() as td:
    root = Path(td)
    archive = root / 'empty_member.zip'
    with zipfile.ZipFile(archive, 'w') as zf:
        zf.writestr('empty.txt', b'')
        zf.writestr('nonempty.txt', b'ok')
    result = extract_zip(str(archive), output_dir=str(root / 'out'), overwrite=True)
    empty_path = str(root / 'out' / 'empty.txt')
    require(result['output_sizes'][empty_path] == 0, 'ZIP did not preserve zero-byte metadata')
    verified = _verify_tool_result('extract_zip', {'zip_path': str(archive)}, result)
    require(verified.get('verified_output') is True, 'Common verifier rejected a legal empty ZIP member')
    require((root / 'out' / 'empty.txt').is_file(), 'Empty ZIP member was not extracted')

# 5: image format aliases accept lower-case jpg.
from PIL import Image
with tempfile.TemporaryDirectory() as td:
    root = Path(td)
    src = root / 'source.png'
    out = root / 'converted.jpg'
    Image.new('RGB', (8, 8), 'white').save(src)
    result = image_compose([str(src)], str(out), 'convert', {'format': 'jpg'})
    require(result['format'] == 'JPEG', f'jpg alias did not normalize to JPEG: {result}')
    with Image.open(out) as image:
        require(image.format == 'JPEG', 'Converted image is not JPEG')

# 9: formula-only XLSX without cached values must not become an all-null matrix.
from openpyxl import Workbook
with tempfile.TemporaryDirectory() as td:
    path = Path(td) / 'formula_only.xlsx'
    wb = Workbook()
    ws = wb.active
    ws['A1'] = '=1+1'
    ws['B1'] = '=SUM(2,3)'
    wb.save(path)
    values = read_xlsx(str(path), range='A1:B1', extract='values')
    sheet = values['sheets'][values['sheet_names'][0]]
    require(sheet['rows'][0] == ['=1+1', '=SUM(2,3)'], f'Formula fallback missing: {sheet}')
    require(sheet['has_uncached_formulas'] is True, 'Uncached formula flag missing')
    require(values['uncached_formula_count'] == 2, 'Unexpected uncached formula count')
    all_data = read_xlsx(str(path), range='A1:B1', extract='all')
    all_sheet = all_data['sheets'][all_data['sheet_names'][0]]
    require(all_sheet['formula_rows'][0] == ['=1+1', '=SUM(2,3)'], 'extract=all lost formula rows')

# 6/7/11: OUTPUT_DIR write->read, safe os.path, version diagnostics and dunder boundary.
with tempfile.TemporaryDirectory() as td:
    code = '''
from os import path
p = path.join(OUTPUT_DIR, "roundtrip.bin")
with open(p, "wb") as f:
    f.write(b"v15")
with open(p, "rb") as f:
    data = f.read()
print(path.basename(p), data.decode())
'''
    result = execute_python(code, output_dir=td)
    require(result['success'] is True, result.get('error'))
    require('roundtrip.bin v15' in result['output'], f'OUTPUT_DIR readback failed: {result}')

valid, errors = validate_code('import os.path as path\nprint(path.join("a", "b"))')
require(valid, f'os.path safe facade rejected: {errors}')
valid, errors = validate_code('from importlib.metadata import version\nprint(version("pip"))')
require(valid, f'importlib.metadata safe diagnostic rejected: {errors}')
valid, errors = validate_code('import pandas as pd\nprint(pd.__version__)')
require(valid, f'module.__version__ rejected: {errors}')
valid, errors = validate_code('import os\nos.system("id")')
require(not valid, 'Full os import/system call must remain blocked')
valid, errors = validate_code('x = (1).__class__')
require(not valid, 'Dangerous __class__ access must remain blocked')

# 3/4: FFmpeg output-aware codec and numeric mix duration.
native = source('lib/core/services/native_tool_executor.dart')
require("case '.mp3':" in native and 'libmp3lame' in native, 'MP3 output codec inference missing')
require("case '.opus':" in native and 'libopus' in native, 'Opus output codec inference missing')
require("case '.wav':" in native and 'pcm_s16le' in native, 'WAV output codec inference missing')
require("durationLimit = '-t $seconds'" in native, 'Numeric mix duration is not converted to output time limit')
require("duration=$durationMode:normalize=0" in native, 'amix duration mode missing')
require("if (_isAudioOnlyOutput(outputPath))" in native, 'Audio-only convert path missing')
require("-vn -af \"$af\" $audioCodec" in native, 'Audio-only filter path missing')
require("case 'extract_frame':" in native and '-vframes 1' in native, 'extract_frame regression path missing')
require("case 'convert':" in native, 'convert regression path missing')

# 8: OCR no-text outcome has an explicit signal.
require("'text_detected': textDetected" in native, 'OCR text_detected flag missing')
require("'reason': 'no_text_detected'" in native, 'OCR no_text_detected reason missing')

# 10: media downloader has actual browser impersonation and clear non-media handling.
media = source('python/navixmind/tools/media.py')
require('from curl_cffi import requests as browser_requests' in media, 'curl-cffi transfer path missing')
require("impersonate='chrome'" in media, 'Browser impersonation is not enabled')
require('download_media only supports video/audio URLs' in media, 'Non-media error contract missing')
require('formats[-1]' not in media, 'Unsafe formats[-1] fallback remains')

# Unlimited tool-call mode: zero sentinel crosses UI/storage/agent and bypasses
# the tool counter plus tool-driven iteration ceiling.
agent = source('python/navixmind/agent.py')
settings = source('lib/features/settings/settings_screen.dart')
storage = source('lib/core/services/storage_service.dart')
require('unlimited_tool_calls = max_tool_calls <= 0' in agent, 'Unlimited sentinel missing in agent')
require('while unlimited_tool_calls or iteration < max_iterations:' in agent, 'Iteration ceiling still blocks unlimited tool mode')
require('if not unlimited_tool_calls and tool_call_count > max_tool_calls:' in agent, 'Tool-call counter still caps unlimited mode')
require("DropdownMenuItem(value: 0, child: Text('不限次数'))" in settings, 'Unlimited UI option missing')
require('不限制本轮工具调用与工具驱动步骤' in settings, 'Unlimited UI explanation missing')
require('zero means unlimited' in storage.lower(), 'Storage contract for zero sentinel missing')

# Previously uncompleted coverage actions remain present.
extended = source('python/navixmind/tools/extended_tools.py')
require('"reorder"' in extended and 'pdf_manage' in extended, 'pdf_manage reorder path missing')
require('smart_crop' in names, 'smart_crop function missing')

print('V15 validation passed: document/runtime, sandbox, media, OCR, XLSX, ZIP, and unlimited-tool gates are green.')
