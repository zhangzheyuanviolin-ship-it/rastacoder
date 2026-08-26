#!/usr/bin/env python3
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path
from unittest.mock import patch as mock_patch

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / 'python'))


def require(condition, message):
    if not condition:
        raise AssertionError(message)


def expect_tool_error(fn, contains=None):
    from navixmind.bridge import ToolError
    try:
        fn()
    except ToolError as exc:
        if contains is not None:
            require(contains in str(exc), f'expected ToolError containing {contains!r}, got {exc!r}')
        return str(exc)
    raise AssertionError('expected ToolError, call succeeded')


# ---------------------------------------------------------------------------
# Static systemic invariants.
# ---------------------------------------------------------------------------
registry_text = (ROOT / 'python/navixmind/tools/__init__.py').read_text(encoding='utf-8')
compat_text = (ROOT / 'python/navixmind/tools/compat.py').read_text(encoding='utf-8')
native_text = (ROOT / 'lib/core/services/native_tool_executor.dart').read_text(encoding='utf-8')
media_text = (ROOT / 'python/navixmind/tools/media.py').read_text(encoding='utf-8')
storage_text = (ROOT / 'lib/core/services/storage_service.dart').read_text(encoding='utf-8')
bridge_text = (ROOT / 'lib/core/bridge/bridge.dart').read_text(encoding='utf-8')
ui_text = (ROOT / 'lib/features/settings/tool_skills_screen.dart').read_text(encoding='utf-8')
doc_text = (ROOT / 'python/navixmind/tools/documents.py').read_text(encoding='utf-8')
pdf_text = (ROOT / 'python/navixmind/tools/extended_tools.py').read_text(encoding='utf-8')
calendar_text = (ROOT / 'python/navixmind/tools/google_api.py').read_text(encoding='utf-8')

for marker in (
    'RASTACODER_V9_SYSTEMIC_HARDENING',
    '_verify_tool_result_artifacts',
    'search_setting_removed:',
    'params:scalar->params.factor',
    'default:office_output=in_place_or_visible_copy',
):
    require(marker in registry_text + compat_text, f'missing systemic marker: {marker}')

require("case 'speed':" in native_text, 'native FFmpeg speed operation missing')
require('atempo=' in native_text and 'setpts=PTS/$factor' in native_text, 'speed audio/video filters missing')
require('duration_verified' in native_text and 'Speed verification failed' in native_text,
        'FFmpeg semantic duration post-verification missing')
require('output file is empty' in native_text, 'native FFmpeg zero-byte output gate missing')
require("or headers" not in media_text, 'media downloader still has undefined headers fallback')
require('Media download verification failed' in media_text, 'download post-verification missing')
require('setSearchProviderSettings' in storage_text and 'getConfiguredSearchProviderSettings' in storage_text,
        'private search settings storage missing')
require("'search_settings': searchSettings" in bridge_text, 'search settings are not injected privately')
require('_configureSearchSettings' in ui_text, 'per-provider search settings UI missing')
for provider in ('AnySearch', 'Exa', 'LangSearch', 'Tavily'):
    require(f'配置 $providerLabel 搜索参数' in ui_text or provider in ui_text, f'{provider} search settings UI missing')
require('appended paragraph is not at document end' in doc_text, 'DOCX append post-verification missing')
require('replace_text made no change' in doc_text, 'Office no-op replace gate missing')
require('post-verification failed' in doc_text, 'Office reopen verification missing')
require('delete_pages requires explicit pages' in pdf_text, 'PDF omitted-pages delete safety gate missing')
require('would remove every page' in pdf_text, 'PDF delete-all safety gate missing')
require('Calendar update verification failed' in calendar_text, 'Calendar update server-ID verification missing')


# ---------------------------------------------------------------------------
# Catalog/schema invariants: keep all 25 Skills / 37 canonical local tools.
# Four search functions are query-only to the local model.
# ---------------------------------------------------------------------------
from navixmind.tools import (  # noqa: E402
    OFFLINE_TOOLS_SCHEMA, LOCAL_SKILLS, get_enabled_tool_names,
    get_offline_tools_for_skills, _verify_tool_result_artifacts,
)
from navixmind.tools.compat import normalize_tool_call  # noqa: E402
from navixmind.bridge import ToolError  # noqa: E402

schema = {item['name']: item for item in OFFLINE_TOOLS_SCHEMA}
covered = get_enabled_tool_names(tuple(LOCAL_SKILLS.keys()))
require(len(LOCAL_SKILLS) == 25, f'expected 25 manual Skills, got {len(LOCAL_SKILLS)}')
require(len(schema) == 37, f'expected 37 canonical functions, got {len(schema)}')
require(set(schema) == covered, f'Skill/schema mismatch: schema-only={set(schema)-covered}; skill-only={covered-set(schema)}')

search_names = ('anysearch_search', 'exa_search', 'langsearch_search', 'tavily_search')
for name in search_names:
    entry = schema[name]['input_schema']
    require(set(entry['properties']) == {'query'}, f'{name} leaks provider tuning into local schema: {entry["properties"].keys()}')
    require(entry.get('required') == ['query'], f'{name} must require only query')
    require(entry.get('additionalProperties') is False, f'{name} should reject model-added provider settings')
    prompt_entry = get_offline_tools_for_skills([name])[0] if name in LOCAL_SKILLS else None

ffmpeg_enum = schema['ffmpeg_process']['input_schema']['properties']['operation']['enum']
require('speed' in ffmpeg_enum, f'speed missing from local FFmpeg enum: {ffmpeg_enum}')
require(schema['ffmpeg_process']['input_schema']['required'] == ['operation'],
        'local FFmpeg schema should let path/output be inferred for the small model')
for office in ('modify_docx', 'modify_pptx', 'modify_xlsx'):
    office_schema = schema[office]['input_schema']
    require('action' in office_schema['properties'] and 'params' in office_schema['properties'],
            f'{office} local schema is still nested/complex')
    require(office_schema['required'] == ['action'], f'{office} should require only action locally')


# ---------------------------------------------------------------------------
# Exact reported model-call regressions.
# ---------------------------------------------------------------------------
local_ctx = {
    '_allowed_tools': {'ffmpeg_process', 'file_info', 'file_manage', 'list_files'},
    '_current_files': ['analysis_article.mp3'],
    '_file_map': {'analysis_article.mp3': '/tmp/analysis_article.mp3'},
    'output_dir': '/tmp/output',
}
name, args, notes = normalize_tool_call(
    'ffmpeg_process',
    {'input_path': 'analysis_article.mp3', 'operation': 'speed', 'params': '1.5'},
    context=local_ctx,
)
require(name == 'ffmpeg_process', name)
require(args.get('operation') == 'speed', args)
require(abs(float(args.get('params', {}).get('factor')) - 1.5) < 1e-9, f'scalar speed factor lost: {args}')
require(args.get('output_path', '').endswith('analysis_article_speed.mp3'), f'speed output name not derived: {args}')
require('params:scalar->params.factor' in notes, f'speed scalar repair not diagnosed: {notes}')

search_ctx = {'_allowed_tools': {'exa_search'}}
name, args, notes = normalize_tool_call(
    'exa_search',
    {'num_results': 5, 'topic': '2026年世界杯决赛', 'search_type': 'news', 'start_published_date': '2023-06-01'},
    context=search_ctx,
)
require(name == 'exa_search', name)
require(args == {'query': '2026年世界杯决赛'}, f'Exa local call was not reduced to query-only: {args}')
require(any(n.startswith('search_setting_removed:') for n in notes), f'provider-setting cleanup not diagnosed: {notes}')

name, args, notes = normalize_tool_call('tavily_search', {'q': '韩国最新科技新闻', 'max_results': 9}, context={'_allowed_tools': {'tavily_search'}})
require(args == {'query': '韩国最新科技新闻'}, f'q->query repair failed: {args}')


# ---------------------------------------------------------------------------
# Search provider settings are applied privately, not supplied by the model.
# ---------------------------------------------------------------------------
from navixmind.tools import search_tools  # noqa: E402

class FakeResponse:
    def __init__(self, body, status_code=200):
        self._body = body
        self.status_code = status_code
        self.ok = 200 <= status_code < 300
        self.text = json.dumps(body, ensure_ascii=False)
    def json(self):
        return self._body
    def raise_for_status(self):
        if not self.ok:
            raise RuntimeError(self.status_code)

captured = []
def fake_post(url, **kwargs):
    captured.append((url, kwargs))
    if 'exa.ai' in url:
        return FakeResponse({'results': [{'title': 'ok'}]})
    if 'langsearch.com' in url:
        return FakeResponse({'code': 200, 'data': {'webPages': {'value': [{'name': 'ok'}]}}})
    if 'tavily.com' in url:
        return FakeResponse({'answer': 'ok', 'results': [{'title': 'ok'}]})
    return FakeResponse({'result': {'content': [{'type': 'text', 'text': 'ok'}]}})

search_context = {
    'search_api_keys': {'anysearch': 'a', 'exa': 'e', 'langsearch': 'l', 'tavily': 't'},
    'search_settings': {
        'anysearch': {'max_results': 7, 'domain': 'example.com'},
        'exa': {'num_results': 8, 'topic': 'news', 'search_type': 'auto', 'include_text': False, 'include_summary': True},
        'langsearch': {'count': 6, 'freshness': 'oneWeek', 'summary': False},
        'tavily': {'max_results': 9, 'topic': 'news', 'search_depth': 'advanced', 'include_answer': False},
    },
}
with mock_patch('navixmind.tools.search_tools.requests.post', side_effect=fake_post):
    search_tools.anysearch_search('alpha', _context=search_context)
    search_tools.exa_search('beta', _context=search_context)
    search_tools.langsearch_search('gamma', _context=search_context)
    search_tools.tavily_search('delta', _context=search_context)

require(len(captured) == 4, f'expected four mocked provider calls, got {len(captured)}')
any_payload = captured[0][1]['json']['params']['arguments']
exa_payload = captured[1][1]['json']
lang_payload = captured[2][1]['json']
tav_payload = captured[3][1]['json']
require(any_payload['query'] == 'alpha' and any_payload['max_results'] == 7 and any_payload['domain'] == 'example.com', any_payload)
require(exa_payload['query'] == 'beta' and exa_payload['numResults'] == 8 and exa_payload['topic'] == 'news', exa_payload)
require(lang_payload['query'] == 'gamma' and lang_payload['count'] == 6 and lang_payload['freshness'] == 'oneWeek' and lang_payload['summary'] is False, lang_payload)
require(tav_payload['query'] == 'delta' and tav_payload['max_results'] == 9 and tav_payload['search_depth'] == 'advanced' and tav_payload['include_answer'] is False, tav_payload)


# ---------------------------------------------------------------------------
# Real Office filesystem mutations: valid changes persist; silent no-ops fail.
# ---------------------------------------------------------------------------
from docx import Document  # noqa: E402
from pptx import Presentation  # noqa: E402
from openpyxl import Workbook, load_workbook  # noqa: E402
from navixmind.tools.documents import modify_docx, modify_pptx, modify_xlsx  # noqa: E402

with tempfile.TemporaryDirectory() as td:
    root = Path(td)

    # DOCX: exact user case, append paragraph at the document end, in-place.
    doc_path = root / 'article.docx'
    doc = Document(); doc.add_paragraph('第一段'); doc.save(doc_path)
    result = modify_docx(str(doc_path), str(doc_path), [{'action': 'add_paragraph', 'params': {'text': '这是追加到结尾的新段落。'}}])
    require(result['success'] and result['verified'] and result['operations_applied'] == 1, result)
    reopened = Document(doc_path)
    require(reopened.paragraphs[-1].text == '这是追加到结尾的新段落。', 'DOCX append did not persist at document end')
    expect_tool_error(lambda: modify_docx(str(doc_path), str(doc_path), [{'action': 'replace_text', 'params': {'old': '绝对不存在的文本', 'new': 'x'}}]), 'made no change')
    expect_tool_error(lambda: modify_docx(str(doc_path), str(doc_path), [{'action': 'update_table_cell', 'params': {'table': 9, 'row': 0, 'col': 0, 'text': 'x'}}]), 'out of range')

    # PPTX: no-match and invalid slide must fail, valid slide addition persists.
    ppt_path = root / 'slides.pptx'
    prs = Presentation(); prs.slides.add_slide(prs.slide_layouts[0]); prs.save(ppt_path)
    expect_tool_error(lambda: modify_pptx(str(ppt_path), str(ppt_path), [{'action': 'replace_text', 'params': {'old': 'missing', 'new': 'x'}}]), 'made no change')
    expect_tool_error(lambda: modify_pptx(str(ppt_path), str(ppt_path), [{'action': 'delete_slide', 'params': {'slide': 99}}]), 'out of range')
    result = modify_pptx(str(ppt_path), str(ppt_path), [{'action': 'add_slide', 'params': {'title': 'V9新增页', 'content': '真实写入'}}])
    require(result['verified'] and len(Presentation(ppt_path).slides) == 2, result)

    # XLSX: no-op/invalid mutation must fail; valid cell mutation persists.
    xlsx_path = root / 'data.xlsx'
    wb = Workbook(); wb.active['A1'] = 'old'; wb.save(xlsx_path); wb.close()
    expect_tool_error(lambda: modify_xlsx(str(xlsx_path), str(xlsx_path), [{'action': 'set_cell', 'params': {'cell': 'A1', 'value': 'old'}}]), 'made no change')
    expect_tool_error(lambda: modify_xlsx(str(xlsx_path), str(xlsx_path), [{'action': 'delete_sheet', 'params': {'name': 'missing'}}]), 'not found')
    result = modify_xlsx(str(xlsx_path), str(xlsx_path), [{'action': 'set_cell', 'params': {'cell': 'A1', 'value': 'new'}}])
    require(result['verified'], result)
    check = load_workbook(xlsx_path); require(check.active['A1'].value == 'new', 'XLSX mutation did not persist'); check.close()

    # Generic artifact verification rejects fake success.
    missing = str(root / 'never_created.bin')
    expect_tool_error(lambda: _verify_tool_result_artifacts('fake_tool', {'success': True, 'output_path': missing}), 'reported success')
    empty = root / 'empty.bin'; empty.write_bytes(b'')
    expect_tool_error(lambda: _verify_tool_result_artifacts('fake_tool', {'success': True, 'output_path': str(empty)}), 'empty')


# ---------------------------------------------------------------------------
# Existing file lifecycle and PDF safety regressions stay fixed.
# ---------------------------------------------------------------------------
from navixmind.tools.extended_tools import file_manage, pdf_manage  # noqa: E402
from pypdf import PdfWriter  # noqa: E402

with tempfile.TemporaryDirectory() as td:
    output = Path(td) / 'output'; output.mkdir()
    file_manage('mkdir', path='nested', _output_dir=str(output))
    touched = file_manage('touch', path='nested/test.txt', _output_dir=str(output))['path']
    Path(touched).write_text('delete me', encoding='utf-8')
    deleted = file_manage('delete', path='output/nested', recursive=True, _output_dir=str(output))
    require(deleted['success'] and deleted['exists_after'] is False and not (output / 'nested').exists(), 'recursive delete regression')

    pdf = output / 'two.pdf'
    writer = PdfWriter(); writer.add_blank_page(width=100, height=100); writer.add_blank_page(width=100, height=100)
    with open(pdf, 'wb') as f: writer.write(f)
    expect_tool_error(lambda: pdf_manage('delete_pages', input_path=str(pdf), output_path=str(output / 'out.pdf'), _output_dir=str(output)), 'requires explicit pages')
    expect_tool_error(lambda: pdf_manage('delete_pages', input_path=str(pdf), output_path=str(output / 'out.pdf'), pages='1-2', _output_dir=str(output)), 'remove every page')
    expect_tool_error(lambda: pdf_manage('rotate', input_path=str(pdf), output_path=str(output / 'rot.pdf'), pages='1', rotation=360, _output_dir=str(output)), 'made no change')


# ---------------------------------------------------------------------------
# FFmpeg speed semantics: exercise the same atempo equation on the CI host.
# Flutter Analyze separately compiles the actual Dart/FFmpegKit implementation.
# ---------------------------------------------------------------------------
ffmpeg = subprocess.run(['bash', '-lc', 'command -v ffmpeg'], capture_output=True, text=True).stdout.strip()
ffprobe = subprocess.run(['bash', '-lc', 'command -v ffprobe'], capture_output=True, text=True).stdout.strip()
require(ffmpeg and ffprobe, 'CI image must provide ffmpeg/ffprobe for V9 speed semantic test')
with tempfile.TemporaryDirectory() as td:
    td = Path(td); src = td / 'tone.wav'; dst = td / 'tone_speed.wav'
    subprocess.run([ffmpeg, '-hide_banner', '-loglevel', 'error', '-y', '-f', 'lavfi', '-i', 'sine=frequency=440:duration=3', str(src)], check=True)
    subprocess.run([ffmpeg, '-hide_banner', '-loglevel', 'error', '-y', '-i', str(src), '-vn', '-filter:a', 'atempo=1.500000', str(dst)], check=True)
    value = subprocess.run([ffprobe, '-v', 'error', '-show_entries', 'format=duration', '-of', 'default=noprint_wrappers=1:nokey=1', str(dst)], check=True, capture_output=True, text=True).stdout.strip()
    duration = float(value)
    require(abs(duration - 2.0) < 0.2, f'1.5x speed semantic duration wrong: {duration}')

print('V9 systemic functional validation passed: 25 Skills / 37 tools, query-only search, speed semantics, verified Office/PDF/file mutations, artifact truth gates.')
