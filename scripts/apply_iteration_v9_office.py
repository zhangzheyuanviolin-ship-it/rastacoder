#!/usr/bin/env python3
"""Make Office modifications transactional and reject silent no-op success."""
from pathlib import Path


def replace_between(path: str, start: str, end: str, new_block: str) -> None:
    p = Path(path)
    text = p.read_text(encoding="utf-8")
    si = text.find(start)
    ei = text.find(end, si + len(start))
    if si < 0 or ei < 0:
        raise SystemExit(f"{path}: block anchors missing: {start!r} -> {end!r}")
    p.write_text(text[:si] + new_block.rstrip() + "\n\n" + text[ei:], encoding="utf-8")


replace_between(
    "python/navixmind/tools/documents.py",
    "def modify_docx(",
    "# ---------------------------------------------------------------------------\n# PPTX read/write",
    r'''def modify_docx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify DOCX transactionally. Every requested operation must have a real effect."""
    from docx import Document

    validate_file_for_processing(input_path, "document")
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_docx requires a non-empty operations list")

    try:
        doc = Document(input_path)
        applied = 0
        verification = []

        def paragraphs_with_tables():
            for paragraph in doc.paragraphs:
                yield paragraph
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            yield paragraph

        for index, op in enumerate(operations):
            if not isinstance(op, dict):
                raise ToolError(f"DOCX operation #{index + 1} must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params") if isinstance(op.get("params"), dict) else {}

            if action == "replace_text":
                old = str(params.get("old") or "")
                new = str(params.get("new") or "")
                if not old:
                    raise ToolError("replace_text requires non-empty params.old")
                replacements = 0
                for para in paragraphs_with_tables():
                    if old in para.text:
                        replacements += para.text.count(old)
                        para.text = para.text.replace(old, new)
                if replacements == 0:
                    raise ToolError(f"replace_text found no matches for: {old!r}")
                verification.append(("replace_text", new, replacements))

            elif action == "add_paragraph":
                text = str(params.get("text") or "")
                if not text:
                    raise ToolError("add_paragraph requires non-empty params.text")
                style = params.get("style")
                doc.add_paragraph(text, style=style)
                verification.append(("paragraph", len(doc.paragraphs) - 1, text))

            elif action == "update_table_cell":
                table_idx = int(params.get("table", 0))
                row_idx = int(params.get("row", 0))
                col_idx = int(params.get("col", 0))
                if min(table_idx, row_idx, col_idx) < 0 or table_idx >= len(doc.tables):
                    raise ToolError("update_table_cell table/row/col is out of range")
                table = doc.tables[table_idx]
                if row_idx >= len(table.rows) or col_idx >= len(table.rows[row_idx].cells):
                    raise ToolError("update_table_cell table/row/col is out of range")
                text = str(params.get("text") or "")
                table.rows[row_idx].cells[col_idx].text = text
                verification.append(("cell", table_idx, row_idx, col_idx, text))

            elif action == "add_heading":
                text = str(params.get("text") or "")
                if not text:
                    raise ToolError("add_heading requires non-empty params.text")
                level = int(params.get("level", 1))
                if level < 0 or level > 9:
                    raise ToolError("add_heading level must be between 0 and 9")
                doc.add_heading(text, level=level)
                verification.append(("paragraph", len(doc.paragraphs) - 1, text))

            elif action == "add_page_break":
                before = len(doc.paragraphs)
                doc.add_page_break()
                if len(doc.paragraphs) <= before:
                    raise ToolError("add_page_break produced no structural change")

            elif action == "add_table":
                rows = params.get("rows") or []
                if not isinstance(rows, list) or not rows:
                    raise ToolError("add_table requires params.rows as a non-empty 2D list")
                col_count = max(len(row) if isinstance(row, list) else 1 for row in rows)
                before = len(doc.tables)
                table = doc.add_table(rows=len(rows), cols=col_count)
                for r, row in enumerate(rows):
                    values = row if isinstance(row, list) else [row]
                    for c, value in enumerate(values):
                        table.cell(r, c).text = str(value)
                verification.append(("table_count", before + 1))

            elif action == "add_image":
                image_path = params.get("image_path")
                if not image_path or not os.path.isfile(image_path):
                    raise ToolError(f"add_image image not found: {image_path}")
                before = len(doc.inline_shapes)
                width_inches = params.get("width_inches")
                if width_inches is None:
                    doc.add_picture(image_path)
                else:
                    from docx.shared import Inches
                    doc.add_picture(image_path, width=Inches(float(width_inches)))
                verification.append(("inline_shapes", before + 1))

            else:
                raise ToolError(f"Unknown DOCX modify action: {action}")
            applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        same_path = os.path.abspath(input_path) == os.path.abspath(output_path)
        save_path = output_path + ".rastacoder_tmp.docx" if same_path else output_path
        try:
            doc.save(save_path)
            check = Document(save_path)
            for item in verification:
                kind = item[0]
                if kind == "paragraph":
                    _, para_idx, expected = item
                    if para_idx >= len(check.paragraphs) or check.paragraphs[para_idx].text != expected:
                        raise ToolError("DOCX paragraph verification failed after save")
                elif kind == "cell":
                    _, t, r, c, expected = item
                    if check.tables[t].rows[r].cells[c].text != expected:
                        raise ToolError("DOCX table-cell verification failed after save")
                elif kind == "table_count" and len(check.tables) < item[1]:
                    raise ToolError("DOCX table verification failed after save")
                elif kind == "inline_shapes" and len(check.inline_shapes) < item[1]:
                    raise ToolError("DOCX image verification failed after save")
            if same_path:
                os.replace(save_path, output_path)
        finally:
            if same_path and os.path.exists(save_path):
                try:
                    os.remove(save_path)
                except OSError:
                    pass

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
            "operations_requested": len(operations),
            "verified": True,
            "in_place": same_path,
        }
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"Failed to modify DOCX: {exc}")


# ---------------------------------------------------------------------------
# PPTX read/write''',
)

replace_between(
    "python/navixmind/tools/documents.py",
    "def modify_pptx(",
    "# ---------------------------------------------------------------------------\n# XLSX read/write",
    r'''def modify_pptx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify PPTX transactionally and reject silent no-op operations."""
    from pptx import Presentation

    validate_file_for_processing(input_path, "document")
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_pptx requires a non-empty operations list")

    try:
        prs = Presentation(input_path)
        applied = 0
        expected_slide_count = len(prs.slides)
        verification = []

        for index, op in enumerate(operations):
            if not isinstance(op, dict):
                raise ToolError(f"PPTX operation #{index + 1} must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params") if isinstance(op.get("params"), dict) else {}

            if action == "replace_text":
                old = str(params.get("old") or "")
                new = str(params.get("new") or "")
                if not old:
                    raise ToolError("replace_text requires non-empty params.old")
                replacements = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in shape.text_frame.paragraphs:
                            if old in para.text:
                                replacements += para.text.count(old)
                                para.text = para.text.replace(old, new)
                if replacements == 0:
                    raise ToolError(f"replace_text found no matches for: {old!r}")

            elif action == "add_slide":
                layout_idx = int(params.get("layout_index", 1))
                if layout_idx < 0:
                    raise ToolError("add_slide layout_index must be non-negative")
                layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
                slide = prs.slides.add_slide(layout)
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 0 and params.get("title") is not None:
                        ph.text = str(params.get("title"))
                    elif ph.placeholder_format.idx == 1 and params.get("content") is not None:
                        ph.text = str(params.get("content"))
                expected_slide_count += 1

            elif action == "update_slide_text":
                slide_num = int(params.get("slide", 1)) - 1
                shape_name = str(params.get("shape_name") or "")
                if not (0 <= slide_num < len(prs.slides)) or not shape_name:
                    raise ToolError("update_slide_text requires a valid slide and shape_name")
                found = False
                for shape in prs.slides[slide_num].shapes:
                    if shape.name == shape_name and shape.has_text_frame:
                        expected = str(params.get("text") or "")
                        shape.text_frame.text = expected
                        verification.append(("shape_text", slide_num, shape_name, expected))
                        found = True
                        break
                if not found:
                    raise ToolError(f"Shape not found or not text-editable: {shape_name}")

            elif action == "set_notes":
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("set_notes slide is out of range")
                expected = str(params.get("text") or "")
                prs.slides[slide_num].notes_slide.notes_text_frame.text = expected
                verification.append(("notes", slide_num, expected))

            elif action == "add_textbox":
                from pptx.util import Inches
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("add_textbox slide is out of range")
                expected = str(params.get("text") or "")
                box = prs.slides[slide_num].shapes.add_textbox(
                    Inches(float(params.get("left", 1))),
                    Inches(float(params.get("top", 1))),
                    Inches(float(params.get("width", 6))),
                    Inches(float(params.get("height", 1))),
                )
                box.text_frame.text = expected
                verification.append(("slide_contains", slide_num, expected))

            elif action == "add_image":
                from pptx.util import Inches
                slide_num = int(params.get("slide", 1)) - 1
                image_path = params.get("image_path")
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("add_image slide is out of range")
                if not image_path or not os.path.isfile(image_path):
                    raise ToolError(f"add_image image not found: {image_path}")
                before = len(prs.slides[slide_num].shapes)
                kwargs = {
                    "left": Inches(float(params.get("left", 1))),
                    "top": Inches(float(params.get("top", 1))),
                }
                if params.get("width") is not None:
                    kwargs["width"] = Inches(float(params["width"]))
                if params.get("height") is not None:
                    kwargs["height"] = Inches(float(params["height"]))
                prs.slides[slide_num].shapes.add_picture(image_path, **kwargs)
                verification.append(("shape_count", slide_num, before + 1))

            elif action == "delete_slide":
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("delete_slide slide is out of range")
                slide_id = prs.slides._sldIdLst[slide_num]
                prs.part.drop_rel(slide_id.rId)
                prs.slides._sldIdLst.remove(slide_id)
                expected_slide_count -= 1

            else:
                raise ToolError(f"Unknown PPTX modify action: {action}")
            applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        same_path = os.path.abspath(input_path) == os.path.abspath(output_path)
        save_path = output_path + ".rastacoder_tmp.pptx" if same_path else output_path
        try:
            prs.save(save_path)
            check = Presentation(save_path)
            if len(check.slides) != expected_slide_count:
                raise ToolError("PPTX slide-count verification failed after save")
            for item in verification:
                kind = item[0]
                if kind == "shape_text":
                    _, slide_num, shape_name, expected = item
                    matches = [shape for shape in check.slides[slide_num].shapes if shape.name == shape_name and shape.has_text_frame]
                    if not matches or matches[0].text_frame.text != expected:
                        raise ToolError("PPTX shape-text verification failed after save")
                elif kind == "notes":
                    _, slide_num, expected = item
                    if check.slides[slide_num].notes_slide.notes_text_frame.text != expected:
                        raise ToolError("PPTX notes verification failed after save")
                elif kind == "slide_contains":
                    _, slide_num, expected = item
                    texts = [shape.text for shape in check.slides[slide_num].shapes if getattr(shape, "has_text_frame", False)]
                    if expected not in texts:
                        raise ToolError("PPTX textbox verification failed after save")
                elif kind == "shape_count":
                    _, slide_num, expected = item
                    if len(check.slides[slide_num].shapes) < expected:
                        raise ToolError("PPTX image verification failed after save")
            if same_path:
                os.replace(save_path, output_path)
        finally:
            if same_path and os.path.exists(save_path):
                try:
                    os.remove(save_path)
                except OSError:
                    pass

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
            "operations_requested": len(operations),
            "verified": True,
            "in_place": same_path,
        }
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"Failed to modify PPTX: {exc}")


# ---------------------------------------------------------------------------
# XLSX read/write''',
)

p = Path("python/navixmind/tools/documents.py")
text = p.read_text(encoding="utf-8")
si = text.find("def modify_xlsx(")
if si < 0:
    raise SystemExit("documents.py: modify_xlsx anchor missing")
text = text[:si] + r'''def modify_xlsx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify XLSX transactionally and reject silent no-op operations."""
    from openpyxl import load_workbook

    validate_file_for_processing(input_path, "document")
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_xlsx requires a non-empty operations list")

    try:
        wb = load_workbook(input_path)
        applied = 0
        verification = []

        for index, op in enumerate(operations):
            if not isinstance(op, dict):
                raise ToolError(f"XLSX operation #{index + 1} must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params") if isinstance(op.get("params"), dict) else {}

            if action in {"set_cell", "set_formula"}:
                sheet_name = str(params.get("sheet") or wb.sheetnames[0])
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet '{sheet_name}' not found.")
                cell = str(params.get("cell") or "").strip()
                if not cell:
                    raise ToolError(f"{action} requires params.cell")
                value = params.get("value") if action == "set_cell" else str(params.get("formula") or "")
                wb[sheet_name][cell] = value
                verification.append(("cell", sheet_name, cell, value))

            elif action == "add_row":
                sheet_name = str(params.get("sheet") or wb.sheetnames[0])
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet '{sheet_name}' not found.")
                values = params.get("values")
                if not isinstance(values, list):
                    raise ToolError("add_row requires params.values as a list")
                before = wb[sheet_name].max_row
                wb[sheet_name].append(values)
                verification.append(("min_rows", sheet_name, before + 1))

            elif action == "add_sheet":
                name = str(params.get("name") or "").strip()
                if not name:
                    raise ToolError("add_sheet requires params.name")
                if name in wb.sheetnames:
                    raise ToolError(f"Sheet already exists: {name}")
                wb.create_sheet(title=name)
                verification.append(("sheet_present", name))

            elif action == "delete_sheet":
                name = str(params.get("name") or "").strip()
                if not name or name not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {name}")
                if len(wb.sheetnames) <= 1:
                    raise ToolError("Cannot delete the only worksheet")
                del wb[name]
                verification.append(("sheet_absent", name))

            elif action == "rename_sheet":
                old_name = str(params.get("old_name") or params.get("sheet") or "").strip()
                new_name = str(params.get("new_name") or "").strip()[:31]
                if not old_name or old_name not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {old_name}")
                if not new_name:
                    raise ToolError("rename_sheet requires params.new_name")
                if new_name in wb.sheetnames and new_name != old_name:
                    raise ToolError(f"Destination sheet name already exists: {new_name}")
                wb[old_name].title = new_name
                verification.append(("sheet_absent", old_name))
                verification.append(("sheet_present", new_name))

            elif action in {"insert_row", "delete_row", "insert_column", "delete_column"}:
                sheet_name = str(params.get("sheet") or wb.active.title)
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {sheet_name}")
                index_value = int(params.get("index", 1))
                amount = int(params.get("amount", 1))
                if index_value < 1 or amount < 1:
                    raise ToolError(f"{action} requires index>=1 and amount>=1")
                ws = wb[sheet_name]
                if action == "insert_row":
                    ws.insert_rows(index_value, amount)
                elif action == "delete_row":
                    ws.delete_rows(index_value, amount)
                elif action == "insert_column":
                    ws.insert_cols(index_value, amount)
                else:
                    ws.delete_cols(index_value, amount)

            else:
                raise ToolError(f"Unknown XLSX modify action: {action}")
            applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        same_path = os.path.abspath(input_path) == os.path.abspath(output_path)
        save_path = output_path + ".rastacoder_tmp.xlsx" if same_path else output_path
        try:
            wb.save(save_path)
            wb.close()
            check = load_workbook(save_path, data_only=False)
            for item in verification:
                kind = item[0]
                if kind == "cell":
                    _, sheet_name, cell, expected = item
                    if check[sheet_name][cell].value != expected:
                        raise ToolError("XLSX cell verification failed after save")
                elif kind == "min_rows":
                    _, sheet_name, minimum = item
                    if check[sheet_name].max_row < minimum:
                        raise ToolError("XLSX row verification failed after save")
                elif kind == "sheet_present" and item[1] not in check.sheetnames:
                    raise ToolError("XLSX sheet-add/rename verification failed after save")
                elif kind == "sheet_absent" and item[1] in check.sheetnames:
                    raise ToolError("XLSX sheet-delete/rename verification failed after save")
            check.close()
            if same_path:
                os.replace(save_path, output_path)
        finally:
            if same_path and os.path.exists(save_path):
                try:
                    os.remove(save_path)
                except OSError:
                    pass

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
            "operations_requested": len(operations),
            "verified": True,
            "in_place": same_path,
        }
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"Failed to modify XLSX: {exc}")
''' + "\n"
p.write_text(text, encoding="utf-8")
print("Applied RastaCoder v9 transactional Office patch")
