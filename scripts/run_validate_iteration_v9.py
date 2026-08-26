#!/usr/bin/env python3
"""Systemic v9 regression gate: user-reported cases plus generalized no-op/postcondition checks."""
from __future__ import annotations

import json
import os
import sys
import tempfile
from pathlib import Path
from unittest.mock import patch as mock_patch

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / 'python'))

from navixmind.bridge import ToolError
from navixmind.tools import (
    TOOLS_SCHEMA,
    OFFLINE_TOOLS_SCHEMA,
    ALL_LOCAL_SKILL_IDS,
    LOCAL_SKILLS,
    LOCAL_TOOL_PROMPT_HINTS,
    _verify_tool_result,
)
from navixmind.tools.compat import normalize_tool_call
from navixmind.tools.documents import modify_docx, modify_pptx, modify_xlsx
from navixmind.tools import search_tools


def require(condition, message):
    if not condition:
        raise AssertionError(message)


def expect_tool_error(fn, message):
    try:
        fn()
    except ToolError:
        return
    raise AssertionError(message)


def schema(name: str, source):
    matches = [item for item in source if item.get('name') == name]
    require(len(matches) == 1, f'{name}: expected one schema, got {len(matches)}')
    return matches[0]


# Capability catalogue remains exhaustive after v8 search expansion.
require(len(ALL_LOCAL_SKILL_IDS) == 25, f'expected 25 Skills, got {len(ALL_LOCAL_SKILL_IDS)}')
covered = {tool for spec in LOCAL_SKILLS.values() for tool in spec['tools']}
offline = {item['name'] for item in OFFLINE_TOOLS_SCHEMA}
require(covered == offline, f'Skill/schema mismatch: missing={sorted(offline-covered)}, extra={sorted(covered-offline)}')
require(len(offline) == 37, f'expected 37 canonical local functions, got {len(offline)}')

# Search model contracts are intentionally one-argument only.
for name in ('anysearch_search', 'exa_search', 'langsearch_search', 'tavily_search'):
    for source in (TOOLS_SCHEMA, OFFLINE_TOOLS_SCHEMA):
        props = schema(name, source)['input_schema']['properties']
        require(set(props) == {'query'}, f'{name}: model-visible search params leaked: {sorted(props)}')
        require(schema(name, source)['input_schema']['required'] == ['query'], f'{name}: query not sole required arg')
    require(LOCAL_TOOL_PROMPT_HINTS[name] == f'{name}(query)', f'{name}: prompt hint still exposes provider knobs')

# Reproduce the user's Exa failure: Qwen placed keywords in topic and emitted provider knobs.
name, args, notes = normalize_tool_call('exa_search', {
    'num_results': 5,
    'topic': '2026年世界杯决赛',
    'search_type': 'news',
    'start_published_date': '2023-06-01',
    'include_domains': [],
})
require(name == 'exa_search', name)
require(args == {'query': '2026年世界杯决赛'}, f'Exa repair failed: {args!r}')
require('topic->query' in notes, f'Exa topic fallback not diagnosed: {notes!r}')

# A correct query must survive generic-key cleanup.
for tool in ('anysearch_search', 'exa_search', 'langsearch_search', 'tavily_search'):
    _, args, _ = normalize_tool_call(tool, {'query': '今天的人工智能新闻', 'max_results': 9})
    require(args == {'query': '今天的人工智能新闻'}, f'{tool}: query lost or extra args leaked: {args}')

# Reproduce the user's audio speed call exactly: scalar params must become factor=1.5.
# context={} reproduces the execution-stage second normalization pass where safe output names are synthesized.
name, args, notes = normalize_tool_call('ffmpeg_process', {
    'input_path': 'analysis_article.mp3',
    'operation': 'speed',
    'params': '1.5',
}, context={})
require(name == 'ffmpeg_process', name)
require(args.get('operation') == 'speed', f'speed op changed: {args}')
require(args.get('params', {}).get('factor') == 1.5, f'scalar speed factor lost: {args}')
require(args.get('output_path') == 'analysis_article_speed.mp3', f'speed output naming failed: {args}')
for source in (TOOLS_SCHEMA, OFFLINE_TOOLS_SCHEMA):
    ops = schema('ffmpeg_process', source)['input_schema']['properties']['operation']['enum']
    require('speed' in ops, f'speed absent from schema: {ops}')

# Search execution must merge private user settings, never model-supplied knobs.
class FakeResponse:
    ok = True
    status_code = 200
    text = ''
    def __init__(self, body): self._body = body
    def json(self): return self._body

captured = {}
def fake_post(url, **kwargs):
    captured['url'] = url
    captured['json'] = kwargs.get('json')
    return FakeResponse({'results': [{'title': 'ok'}]})

search_context = {
    'search_api_keys': {'exa': 'secret'},
    'search_provider_settings': {
        'exa': {
            'num_results': 7,
            'topic': 'news',
            'search_type': 'deep',
            'include_text': False,
            'include_summary': True,
            'include_highlights': False,
            'include_domains': ['example.com'],
        }
    },
}
with mock_patch.object(search_tools.requests, 'post', fake_post):
    result = search_tools.exa_search('世界杯', _context=search_context)
require(result['success'] is True, result)
require(captured['json']['query'] == '世界杯', captured)
require(captured['json']['numResults'] == 7, captured)
require(captured['json']['topic'] == 'news' and captured['json']['type'] == 'deep', captured)
require(captured['json']['includeDomains'] == ['example.com'], captured)

# Office edits: default normalization must be in-place, and the physical file must change.
with tempfile.TemporaryDirectory() as td:
    td = Path(td)
    from docx import Document
    docx = td / 'sample.docx'
    d = Document(); d.add_paragraph('第一段'); d.save(docx)
    _, norm, _ = normalize_tool_call('modify_docx', {
        'input_path': str(docx),
        'operations': [{'action': 'add_paragraph', 'params': {'text': '这是追加到结尾的新段落'}}],
    }, context={})
    require(norm['output_path'] == str(docx), f'DOCX edit should default in-place: {norm}')
    result = modify_docx(**norm)
    require(result['success'] and result['verified'] and result['in_place'], result)
    reopened = Document(docx)
    require(reopened.paragraphs[-1].text == '这是追加到结尾的新段落', 'DOCX append was not physically persisted')
    expect_tool_error(
        lambda: modify_docx(str(docx), str(docx), [{'action': 'replace_text', 'params': {'old': '绝对不存在', 'new': 'x'}}]),
        'DOCX zero-match replace reported success',
    )

    from pptx import Presentation
    pptx = td / 'sample.pptx'
    prs = Presentation(); prs.slides.add_slide(prs.slide_layouts[6]); prs.save(pptx)
    result = modify_pptx(str(pptx), str(pptx), [{'action': 'add_textbox', 'params': {'slide': 1, 'text': '新增文本框'}}])
    require(result['success'] and result['verified'], result)
    reopened_prs = Presentation(pptx)
    require(any(getattr(s, 'has_text_frame', False) and s.text == '新增文本框' for s in reopened_prs.slides[0].shapes), 'PPTX edit not persisted')
    expect_tool_error(
        lambda: modify_pptx(str(pptx), str(pptx), [{'action': 'update_slide_text', 'params': {'slide': 1, 'shape_name': '不存在', 'text': 'x'}}]),
        'PPTX missing shape reported success',
    )

    from openpyxl import Workbook, load_workbook
    xlsx = td / 'sample.xlsx'
    wb = Workbook(); wb.save(xlsx); wb.close()
    result = modify_xlsx(str(xlsx), str(xlsx), [{'action': 'set_cell', 'params': {'cell': 'B2', 'value': '真实写入'}}])
    require(result['success'] and result['verified'], result)
    check = load_workbook(xlsx); require(check.active['B2'].value == '真实写入', 'XLSX edit not persisted'); check.close()
    expect_tool_error(
        lambda: modify_xlsx(str(xlsx), str(xlsx), [{'action': 'delete_sheet', 'params': {'name': '不存在'}}]),
        'XLSX missing sheet delete reported success',
    )

    # Universal postcondition verifier catches missing/empty/corrupt or mismatched outputs.
    txt = td / 'a.txt'; txt.write_text('hello', encoding='utf-8')
    verified = _verify_tool_result('write_file', {'content': 'hello'}, {'success': True, 'output_path': str(txt)})
    require(verified.get('verified_output') is True, verified)
    expect_tool_error(
        lambda: _verify_tool_result('write_file', {'content': 'different'}, {'success': True, 'output_path': str(txt)}),
        'write_file content mismatch passed postcondition',
    )
    empty = td / 'empty.bin'; empty.write_bytes(b'')
    expect_tool_error(
        lambda: _verify_tool_result('download_media', {}, {'success': True, 'output_path': str(empty)}),
        'empty generated file passed postcondition',
    )

# Static native/UI guards for behavior that cannot run on Linux CI.
native = (ROOT / 'lib/core/services/native_tool_executor.dart').read_text(encoding='utf-8')
for marker in (
    "case 'speed':",
    '_buildATempoChain',
    "FFmpeg reported success but output file is empty",
    'Speed verification failed',
    "result['duration_verified']",
):
    require(marker in native, f'native FFmpeg v9 marker missing: {marker}')

media = (ROOT / 'python/navixmind/tools/media.py').read_text(encoding='utf-8')
require("or info.get('http_headers') or {}" in media, 'media fallback headers bug remains')
require("or headers" not in media, 'undefined media headers fallback remains')

storage = (ROOT / 'lib/core/services/storage_service.dart').read_text(encoding='utf-8')
bridge = (ROOT / 'lib/core/bridge/bridge.dart').read_text(encoding='utf-8')
skills_ui = (ROOT / 'lib/features/settings/tool_skills_screen.dart').read_text(encoding='utf-8')
settings_ui = (ROOT / 'lib/features/settings/search_provider_settings_screen.dart').read_text(encoding='utf-8')
require('getAllSearchProviderSettings' in storage, 'search settings persistence missing')
require("'search_provider_settings': searchProviderSettings" in bridge, 'search settings not injected privately')
require('SearchProviderSettingsScreen' in skills_ui, 'search settings button not wired to Skill UI')
require('本地模型只需要提供搜索关键词' in settings_ui, 'settings responsibility explanation missing')

print(json.dumps({
    'status': 'PASS',
    'skills': len(ALL_LOCAL_SKILL_IDS),
    'canonical_tools': len(offline),
    'search_model_args': ['query'],
    'ffmpeg_speed_factor_case': 1.5,
    'office_in_place_verified': True,
    'universal_output_postconditions': True,
}, ensure_ascii=False, indent=2))
