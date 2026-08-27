"""Extended local tools for RastaCoder v7 complete-skill coverage.

These tools intentionally expose common user-facing operations as structured,
bounded functions so a small on-device model does not need shell access.
"""

import fnmatch
import os
import shutil
import zipfile
from typing import Any, Dict, List, Optional

from ..bridge import ToolError
# RASTACODER_V12_EXTENDED_PATH_CONTRACT
from .path_contract import resolve_model_path, resolve_list_path


def _ensure_parent(path: str) -> None:
    parent = os.path.dirname(path)
    if parent:
        os.makedirs(parent, exist_ok=True)


def _default_output_dir(_output_dir: Optional[str]) -> str:
    if _output_dir:
        os.makedirs(_output_dir, exist_ok=True)
        return _output_dir
    return os.getcwd()


# RASTACODER_V11_WORKSPACE_ROOT
# RASTACODER_V12_EXTENDED_PATH_CONTRACT
def _resolve_workspace_path(value: str, _output_dir: Optional[str]) -> str:
    """Resolve through the same central model-facing path contract as execute_tool."""
    root = os.path.normpath(_default_output_dir(_output_dir))
    return resolve_model_path(value, root, allow_android_roots=True)


def _resolve_named_directory(directory: str, _output_dir: Optional[str]) -> str:
    key = (directory or "output").strip().lower()
    mapping = {
        "output": _default_output_dir(_output_dir),
        "downloads": "/storage/emulated/0/Download",
        "documents": "/storage/emulated/0/Documents",
        "pictures": "/storage/emulated/0/Pictures",
        "screenshots": "/storage/emulated/0/Pictures/Screenshots",
        "camera": "/storage/emulated/0/DCIM/Camera",
    }
    return mapping.get(key, directory)


def _resolve_list_target(directory: str, path: Optional[str], _output_dir: Optional[str]) -> str:
    """Resolve list target through the central logical namespace."""
    root = os.path.normpath(_default_output_dir(_output_dir))
    return resolve_list_path(path, root, legacy_directory=directory)


def list_files(
    directory: str = "output",
    path: Optional[str] = None,
    recursive: bool = False,
    pattern: Optional[str] = None,
    include_directories: bool = True,
    _output_dir: Optional[str] = None,
) -> dict:
    """List files/directories in app output or common Android folders."""
    target = _resolve_list_target(directory, path, _output_dir)
    if not os.path.isdir(target):
        raise ToolError(f"Directory not found or inaccessible: {target}")

    entries: List[dict] = []
    try:
        if recursive:
            iterator = (
                os.path.join(root, name)
                for root, dirs, files in os.walk(target)
                for name in (dirs + files if include_directories else files)
            )
        else:
            iterator = (os.path.join(target, name) for name in os.listdir(target))

        for full in iterator:
            name = os.path.basename(full)
            if pattern and not fnmatch.fnmatch(name, pattern):
                continue
            try:
                is_dir = os.path.isdir(full)
                stat = os.stat(full)
                entries.append({
                    "name": name,
                    "path": full,
                    "type": "directory" if is_dir else "file",
                    "size_bytes": 0 if is_dir else stat.st_size,
                    "modified_epoch": int(stat.st_mtime),
                })
            except OSError:
                continue
            if len(entries) >= 1000:
                break
        entries.sort(key=lambda item: (item["type"] != "directory", item["name"].lower()))
        return {
            "directory": target,
            "requested_path": path if path not in (None, "") else ".",
            "workspace_root": os.path.normpath(_default_output_dir(_output_dir)),
            "count": len(entries),
            "recursive": recursive,
            "pattern": pattern,
            "entries": entries,
            "truncated": len(entries) >= 1000,
        }
    except Exception as exc:
        raise ToolError(f"Failed to list files: {exc}")


def file_manage(
    action: str,
    path: Optional[str] = None,
    source_path: Optional[str] = None,
    destination_path: Optional[str] = None,
    recursive: bool = False,
    overwrite: bool = False,
    _output_dir: Optional[str] = None,
) -> dict:
    """Manage files relative to the real app output root and verify mutations."""
    action = (action or "").strip().lower()
    source_raw = source_path or path

    try:
        if action == "list":
            resolved = _resolve_workspace_path(path, _output_dir) if path else _default_output_dir(_output_dir)
            return list_files(path=resolved, directory="output", recursive=recursive, _output_dir=_output_dir)

        if action == "mkdir":
            target_raw = path or destination_path
            if not target_raw:
                raise ToolError("mkdir requires path")
            target = _resolve_workspace_path(target_raw, _output_dir)
            os.makedirs(target, exist_ok=True)
            if not os.path.isdir(target):
                raise ToolError(f"mkdir verification failed: {target}")
            return {"success": True, "action": action, "path": target, "exists_after": True}

        if action == "exists":
            if not source_raw:
                raise ToolError("exists requires path")
            source = _resolve_workspace_path(source_raw, _output_dir)
            exists = os.path.lexists(source)
            return {
                "success": True,
                "action": action,
                "path": source,
                "exists": exists,
                "is_file": os.path.isfile(source),
                "is_directory": os.path.isdir(source),
            }

        if action == "touch":
            target_raw = source_raw or destination_path
            if not target_raw:
                raise ToolError("touch requires path")
            target = _resolve_workspace_path(target_raw, _output_dir)
            _ensure_parent(target)
            with open(target, "a", encoding="utf-8"):
                os.utime(target, None)
            if not os.path.isfile(target):
                raise ToolError(f"touch verification failed: {target}")
            return {"success": True, "action": action, "path": target, "exists_after": True}

        if action in {"copy", "move", "rename"}:
            if not source_raw or not destination_path:
                raise ToolError(f"{action} requires source_path and destination_path")
            source = _resolve_workspace_path(source_raw, _output_dir)
            destination = _resolve_workspace_path(destination_path, _output_dir)
            if not os.path.lexists(source):
                raise ToolError(f"Source not found: {source}")
            if os.path.lexists(destination):
                if not overwrite:
                    raise ToolError(f"Destination already exists: {destination}")
                if os.path.isdir(destination) and not os.path.islink(destination):
                    shutil.rmtree(destination)
                else:
                    os.remove(destination)
            _ensure_parent(destination)
            if action == "copy":
                if os.path.isdir(source) and not os.path.islink(source):
                    shutil.copytree(source, destination)
                else:
                    shutil.copy2(source, destination)
                if not os.path.lexists(source) or not os.path.lexists(destination):
                    raise ToolError(f"Copy verification failed: {source} -> {destination}")
            else:
                shutil.move(source, destination)
                if os.path.lexists(source) or not os.path.lexists(destination):
                    raise ToolError(f"{action} verification failed: {source} -> {destination}")
            return {
                "success": True,
                "action": action,
                "source_path": source,
                "destination_path": destination,
                "destination_exists_after": True,
            }

        if action == "delete":
            if not source_raw:
                raise ToolError("delete requires path")
            source = _resolve_workspace_path(source_raw, _output_dir)
            if not os.path.lexists(source):
                raise ToolError(f"Delete target not found: {source}")
            if os.path.isdir(source) and not os.path.islink(source):
                if recursive:
                    shutil.rmtree(source)
                else:
                    os.rmdir(source)
            else:
                os.remove(source)
            if os.path.lexists(source):
                raise ToolError(f"Delete verification failed; target still exists: {source}")
            return {
                "success": True,
                "action": action,
                "path": source,
                "deleted": True,
                "exists_after": False,
            }

        raise ToolError(
            "Unknown file_manage action. Use list, mkdir, copy, move, rename, delete, touch, or exists."
        )
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"File operation failed ({action}): {exc}")


def list_zip(zip_path: str) -> dict:
    """List a ZIP archive without extracting it."""
    if not os.path.isfile(zip_path):
        raise ToolError(f"ZIP file not found: {zip_path}")
    try:
        with zipfile.ZipFile(zip_path, "r") as zf:
            entries = [
                {
                    "name": info.filename,
                    "size_bytes": info.file_size,
                    "compressed_bytes": info.compress_size,
                    "is_directory": info.is_dir(),
                }
                for info in zf.infolist()
            ]
        return {"path": zip_path, "count": len(entries), "entries": entries[:2000]}
    except zipfile.BadZipFile as exc:
        raise ToolError(f"Invalid ZIP archive: {exc}")


def extract_zip(
    zip_path: str,
    output_dir: Optional[str] = None,
    overwrite: bool = False,
    _output_dir: Optional[str] = None,
) -> dict:
    """Safely extract ZIP contents, rejecting path traversal entries."""
    if not os.path.isfile(zip_path):
        raise ToolError(f"ZIP file not found: {zip_path}")
    root = output_dir
    if not root:
        base = os.path.splitext(os.path.basename(zip_path))[0] + "_extracted"
        root = os.path.join(_default_output_dir(_output_dir), base)
    elif not os.path.isabs(root):
        root = os.path.join(_default_output_dir(_output_dir), root)
    os.makedirs(root, exist_ok=True)
    root_real = os.path.realpath(root)
    extracted: List[str] = []
    extracted_sizes: Dict[str, int] = {}
    try:
        with zipfile.ZipFile(zip_path, "r") as zf:
            for info in zf.infolist():
                dest = os.path.realpath(os.path.join(root_real, info.filename))
                if dest != root_real and not dest.startswith(root_real + os.sep):
                    raise ToolError(f"Unsafe ZIP entry rejected: {info.filename}")
                if info.is_dir():
                    os.makedirs(dest, exist_ok=True)
                    continue
                if os.path.exists(dest) and not overwrite:
                    raise ToolError(f"Extraction destination exists: {dest}")
                os.makedirs(os.path.dirname(dest), exist_ok=True)
                with zf.open(info, "r") as src, open(dest, "wb") as dst:
                    shutil.copyfileobj(src, dst)
                extracted.append(dest)
                extracted_sizes[dest] = int(info.file_size)
        return {
            "success": True,
            "zip_path": zip_path,
            "output_dir": root_real,
            "file_count": len(extracted),
            "output_paths": extracted[:2000],
            "output_sizes": {p: extracted_sizes[p] for p in extracted[:2000]},
        }
    except ToolError:
        raise
    except zipfile.BadZipFile as exc:
        raise ToolError(f"Invalid ZIP archive: {exc}")
    except Exception as exc:
        raise ToolError(f"ZIP extraction failed: {exc}")


def _parse_pages(pages: Any, total: int) -> List[int]:
    """Return zero-based page indices from 1-based list/range input."""
    if pages is None or pages == "all":
        return list(range(total))
    values: List[int] = []
    if isinstance(pages, list):
        raw_parts = pages
    else:
        raw_parts = str(pages).replace(" ", "").split(",")
    for part in raw_parts:
        if isinstance(part, int):
            nums = [part]
        else:
            token = str(part)
            if "-" in token:
                a, b = token.split("-", 1)
                nums = list(range(int(a), int(b) + 1))
            else:
                nums = [int(token)]
        for n in nums:
            idx = n - 1
            if idx < 0 or idx >= total:
                raise ToolError(f"PDF page {n} out of range 1..{total}")
            values.append(idx)
    return values


def pdf_manage(
    action: str,
    input_path: Optional[str] = None,
    input_paths: Optional[List[str]] = None,
    output_path: Optional[str] = None,
    pages: Any = None,
    rotation: int = 90,
    _output_dir: Optional[str] = None,
) -> dict:
    """Manage PDF pages: merge/split/extract/reorder/delete/rotate."""
    from pypdf import PdfReader, PdfWriter

    action = (action or "").strip().lower()
    output_base = _default_output_dir(_output_dir)

    if action == "merge":
        paths = input_paths or ([] if input_path is None else [input_path])
        if len(paths) < 2:
            raise ToolError("PDF merge requires at least two input_paths")
        out = output_path or os.path.join(output_base, "merged.pdf")
        if not os.path.isabs(out):
            out = os.path.join(output_base, out)
        _ensure_parent(out)
        writer = PdfWriter()
        for path in paths:
            if not os.path.isfile(path):
                raise ToolError(f"PDF not found: {path}")
            for page in PdfReader(path).pages:
                writer.add_page(page)
        with open(out, "wb") as f:
            writer.write(f)
        return {"success": True, "action": action, "output_path": out, "input_count": len(paths)}

    if not input_path or not os.path.isfile(input_path):
        raise ToolError("pdf_manage requires a valid input_path for this action")
    reader = PdfReader(input_path)
    total = len(reader.pages)

    if action == "split":
        base = os.path.splitext(os.path.basename(input_path))[0]
        folder = output_path or os.path.join(output_base, f"{base}_pages")
        if not os.path.isabs(folder):
            folder = os.path.join(output_base, folder)
        os.makedirs(folder, exist_ok=True)
        outputs = []
        for i, page in enumerate(reader.pages, start=1):
            writer = PdfWriter()
            writer.add_page(page)
            out = os.path.join(folder, f"{base}_page_{i}.pdf")
            with open(out, "wb") as f:
                writer.write(f)
            outputs.append(out)
        return {"success": True, "action": action, "output_paths": outputs, "page_count": total}

    selected = _parse_pages(pages, total)
    if action == "delete_pages":
        excluded = set(selected)
        selected = [i for i in range(total) if i not in excluded]
    elif action not in {"extract_pages", "reorder", "rotate"}:
        raise ToolError("Unknown pdf_manage action. Use merge, split, extract_pages, reorder, delete_pages, or rotate.")

    out = output_path or os.path.join(
        output_base,
        f"{os.path.splitext(os.path.basename(input_path))[0]}_{action}.pdf",
    )
    if not os.path.isabs(out):
        out = os.path.join(output_base, out)
    _ensure_parent(out)
    writer = PdfWriter()
    for idx in selected:
        page = reader.pages[idx]
        if action == "rotate":
            page.rotate(int(rotation))
        writer.add_page(page)
    with open(out, "wb") as f:
        writer.write(f)
    return {
        "success": True,
        "action": action,
        "output_path": out,
        "page_count": len(selected),
        "source_page_count": total,
    }


def create_pptx(
    output_path: str,
    title: Optional[str] = None,
    slides: Optional[List[Dict[str, Any]]] = None,
) -> dict:
    """Create a PowerPoint presentation from structured slide data."""
    from pptx import Presentation

    _ensure_parent(output_path)
    prs = Presentation()
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = str(title)
    for item in slides or []:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(item.get("title", ""))
        body = "\n".join(item.get("bullets", [])) if isinstance(item.get("bullets"), list) else str(item.get("content", ""))
        for ph in slide.placeholders:
            if getattr(ph.placeholder_format, "idx", -1) == 1:
                ph.text = body
                break
        notes = item.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
    prs.save(output_path)
    return {"success": True, "output_path": output_path, "slide_count": len(prs.slides)}


def create_xlsx(
    output_path: str,
    sheets: Optional[List[Dict[str, Any]]] = None,
) -> dict:
    """Create XLSX from canonical or compatibility row structures, with verification."""
    import json
    from openpyxl import Workbook, load_workbook

    def scalar(value):
        if value is None or isinstance(value, (str, int, float, bool)):
            return value
        return json.dumps(value, ensure_ascii=False, sort_keys=False)

    def normalize_rows(spec):
        raw = spec.get('rows')
        if raw is None:
            raw = spec.get('data')
        if raw is None:
            return []
        if isinstance(raw, tuple):
            raw = list(raw)
        if not isinstance(raw, list):
            raise ToolError('create_xlsx sheet rows/data must be an array')

        # A list of ordinary objects is a record set: write deterministic
        # headers once, then one value row per object.
        if raw and all(isinstance(row, dict) and not (len(row) == 1 and 'item' in row) for row in raw):
            headers = []
            for row in raw:
                for key in row.keys():
                    key = str(key)
                    if key not in headers:
                        headers.append(key)
            return [headers] + [[scalar(row.get(key)) for key in headers] for row in raw]

        matrix = []
        for row in raw:
            # Compatibility adapters and some models emit {"item": [...]}
            # for an array item. Unwrap it instead of writing the key itself.
            if isinstance(row, dict) and len(row) == 1 and 'item' in row:
                row = row['item']
            if isinstance(row, (list, tuple)):
                matrix.append([scalar(value) for value in row])
            elif isinstance(row, dict):
                matrix.append([scalar(value) for value in row.values()])
            else:
                matrix.append([scalar(row)])
        return matrix

    _ensure_parent(output_path)
    configured = sheets or [{"name": "Sheet1", "rows": []}]
    if not isinstance(configured, list) or not all(isinstance(spec, dict) for spec in configured):
        raise ToolError('create_xlsx sheets must be an array of objects')

    wb = Workbook()
    default = wb.active
    expected = {}
    for i, spec in enumerate(configured):
        ws = default if i == 0 else wb.create_sheet()
        ws.title = str(spec.get('name') or spec.get('sheet_name') or f"Sheet{i + 1}")[:31]
        matrix = normalize_rows(spec)
        expected[ws.title] = matrix
        for row in matrix:
            ws.append(row)
    wb.save(output_path)
    wb.close()

    check = load_workbook(output_path, data_only=False)
    try:
        for sheet_name, matrix in expected.items():
            ws = check[sheet_name]
            for r, row in enumerate(matrix, start=1):
                for c, value in enumerate(row, start=1):
                    actual = ws.cell(row=r, column=c).value
                    if actual != value:
                        raise ToolError(
                            f"XLSX verification failed at {sheet_name}!{ws.cell(r, c).coordinate}: "
                            f"expected {value!r}, got {actual!r}"
                        )
    finally:
        check.close()

    return {
        "success": True,
        "output_path": output_path,
        "sheet_names": list(expected.keys()),
        "row_counts": {name: len(rows) for name, rows in expected.items()},
        "verified": True,
    }


def image_compose(
    input_paths: List[str],
    output_path: str,
    operation: str,
    params: Optional[Dict[str, Any]] = None,
) -> dict:
    """Structured image editing using Pillow; supports common composition and conversion tasks."""
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps

    if not input_paths:
        raise ToolError("image_compose requires at least one input image")
    for path in input_paths:
        if not os.path.isfile(path):
            raise ToolError(f"Image not found: {path}")
    params = params or {}
    images = [Image.open(path).convert("RGBA") for path in input_paths]
    op = (operation or "").strip().lower()

    if op == "concat_horizontal":
        h = max(im.height for im in images)
        resized = [im if im.height == h else im.resize((round(im.width * h / im.height), h), Image.Resampling.LANCZOS) for im in images]
        result = Image.new("RGBA", (sum(im.width for im in resized), h), (0, 0, 0, 0))
        x = 0
        for im in resized:
            result.alpha_composite(im, (x, 0)); x += im.width
    elif op == "concat_vertical":
        w = max(im.width for im in images)
        resized = [im if im.width == w else im.resize((w, round(im.height * w / im.width)), Image.Resampling.LANCZOS) for im in images]
        result = Image.new("RGBA", (w, sum(im.height for im in resized)), (0, 0, 0, 0))
        y = 0
        for im in resized:
            result.alpha_composite(im, (0, y)); y += im.height
    elif op == "overlay":
        if len(images) < 2:
            raise ToolError("overlay requires two input images")
        result = images[0].copy()
        result.alpha_composite(images[1], (int(params.get("x", 0)), int(params.get("y", 0))))
    elif op == "resize":
        im = images[0]
        width = params.get("width")
        height = params.get("height")
        if width is None and height is None:
            raise ToolError("resize requires width or height")
        if width is None:
            width = round(im.width * int(height) / im.height)
        if height is None:
            height = round(im.height * int(width) / im.width)
        result = im.resize((int(width), int(height)), Image.Resampling.LANCZOS)
    elif op == "crop":
        im = images[0]
        x, y = int(params.get("x", 0)), int(params.get("y", 0))
        width, height = params.get("width"), params.get("height")
        if width is None or height is None:
            raise ToolError("crop requires width and height")
        result = im.crop((x, y, x + int(width), y + int(height)))
    elif op == "adjust":
        result = images[0].convert("RGB")
        if "brightness" in params:
            result = ImageEnhance.Brightness(result).enhance(float(params["brightness"]))
        if "contrast" in params:
            result = ImageEnhance.Contrast(result).enhance(float(params["contrast"]))
        if "saturation" in params:
            result = ImageEnhance.Color(result).enhance(float(params["saturation"]))
        if "sharpness" in params:
            result = ImageEnhance.Sharpness(result).enhance(float(params["sharpness"]))
        if "gamma" in params:
            gamma = max(0.01, float(params["gamma"]))
            lut = [min(255, round((i / 255.0) ** (1.0 / gamma) * 255)) for i in range(256)]
            result = result.point(lut * 3)
    elif op == "grayscale":
        result = ImageOps.grayscale(images[0])
    elif op == "blur":
        result = images[0].filter(ImageFilter.GaussianBlur(radius=float(params.get("radius", 2))))
    elif op == "rotate":
        result = images[0].rotate(float(params.get("degrees", 90)), expand=True)
    elif op == "flip":
        direction = str(params.get("direction", "horizontal")).lower()
        result = ImageOps.flip(images[0]) if direction == "vertical" else ImageOps.mirror(images[0])
    elif op == "convert":
        result = images[0]
    else:
        raise ToolError(
            "Unknown image_compose operation. Use concat_horizontal, concat_vertical, overlay, resize, adjust, crop, grayscale, blur, rotate, flip, or convert."
        )

    _ensure_parent(output_path)
    ext = os.path.splitext(output_path)[1].lower()
    aliases = {
        "jpg": "JPEG", "jpeg": "JPEG", "png": "PNG", "webp": "WEBP",
        "bmp": "BMP", "gif": "GIF", "tif": "TIFF", "tiff": "TIFF",
    }
    requested_format = str(params.get("format", "")).strip().lower().lstrip(".")
    if requested_format:
        save_format = aliases.get(requested_format)
        if save_format is None:
            raise ToolError(
                f"Unsupported image format: {params.get('format')}. "
                "Use jpg/jpeg, png, webp, bmp, gif, tif, or tiff."
            )
    else:
        save_format = aliases.get(ext.lstrip("."), "PNG")
    if save_format == "JPEG":
        result = result.convert("RGB")
    result.save(output_path, format=save_format, quality=int(params.get("quality", 92)))
    for im in images:
        try:
            im.close()
        except Exception:
            pass
    return {
        "success": True,
        "operation": op,
        "output_path": output_path,
        "width": result.width,
        "height": result.height,
        "format": save_format,
    }
