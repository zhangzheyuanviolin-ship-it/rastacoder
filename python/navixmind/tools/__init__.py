"""
Tools - Native tool implementations for the agent

This module provides tool definitions and execution logic.
"""

from typing import Any, Dict
from pathlib import Path
import copy

from .web import web_fetch, headless_browser
from .documents import (
    read_pdf, create_pdf, create_docx, convert_document, create_zip, read_file, write_file,
    read_docx, modify_docx, read_pptx, modify_pptx, read_xlsx, modify_xlsx,
)
from .media import download_media
from .search_tools import (
    anysearch_search, anysearch_extract, anysearch_get_sub_domains,
    exa_search, langsearch_search, tavily_search,
)
from .google_api import google_calendar, gmail
from .code_executor import python_execute
from .extended_tools import (
    list_files, file_manage, list_zip, extract_zip, pdf_manage,
    create_pptx, create_xlsx, image_compose,
)

from ..bridge import ToolError, get_bridge
from .compat import normalize_tool_call
# RASTACODER_V12_PATH_CONTRACT_IMPORT
from .path_contract import resolve_model_path, resolve_output_path

# RASTACODER_V4_TOOL_CONTRACT


# Tool schema for Claude
TOOLS_SCHEMA = [
    {
        "name": "web_fetch",
        "description": "Fetch a webpage and extract its text content.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "URL to fetch"},
                "extract_mode": {
                    "type": "string",
                    "enum": ["text", "html", "links"],
                    "description": "What to extract from the page"
                }
            },
            "required": ["url"]
        }
    },
    {
        "name": "headless_browser",
        "description": "Load a JavaScript-heavy page in a headless browser and extract content.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "URL to load"},
                "wait_seconds": {
                    "type": "integer",
                    "default": 5,
                    "description": "Seconds to wait for JS to render"
                },
                "extract_selector": {
                    "type": "string",
                    "description": "CSS selector to extract, or empty for full page"
                }
            },
            "required": ["url"]
        }
    },
    {
        "name": "read_pdf",
        "description": "Extract text content from a PDF file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pdf_path": {"type": "string", "description": "Path to PDF file"},
                "pages": {
                    "type": "string",
                    "description": "Page range, e.g., '1-5' or 'all'"
                }
            },
            "required": ["pdf_path"]
        }
    },
    {
        "name": "create_pdf",
        "description": "Create a PDF document from text and/or images. Can embed images (JPG, PNG) directly into the PDF.",
        "input_schema": {
            "type": "object",
            "properties": {
                "content": {"type": "string", "description": "Text content for the PDF (optional if images provided)"},
                "title": {"type": "string", "description": "Document title"},
                "output_path": {"type": "string", "description": "Where to save the PDF"},
                "image_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "List of image file paths to embed in the PDF"
                }
            },
            "required": ["output_path"]
        }
    },
    {
        "name": "convert_document",
        "description": "Text-oriented conversion between TXT, DOCX, PDF, and HTML. Complex layout may be simplified.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Path to .txt, .docx, .pdf, .html, or .htm input"},
                "output_format": {"type": "string", "enum": ["pdf", "html", "txt", "docx"], "description": "Target format"},
                "output_path": {"type": "string", "description": "Optional explicit output path"}
            },
            "required": ["input_path", "output_format"]
        }
    },
    {
        "name": "create_docx",
        "description": "Create a new DOCX Word document from text.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "Where to save the .docx file"},
                "content": {"type": "string", "description": "Document text"},
                "title": {"type": "string", "description": "Optional document title"}
            },
            "required": ["output_path", "content"]
        }
    },
    {
        "name": "read_docx",
        "description": "Extract text, tables, and metadata from a DOCX file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "docx_path": {"type": "string", "description": "Path to the DOCX file"},
                "extract": {
                    "type": "string",
                    "enum": ["text", "tables", "all"],
                    "default": "all",
                    "description": "What to extract"
                }
            },
            "required": ["docx_path"]
        }
    },
    {
        "name": "modify_docx",
        "description": "Modify an existing DOCX file. Can replace text, add paragraphs, update table cells, and save back.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Path to the source DOCX"},
                "output_path": {"type": "string", "description": "Where to save the modified DOCX"},
                "operations": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "action": {"type": "string", "enum": ["replace_text", "add_paragraph", "update_table_cell", "add_heading", "add_page_break", "add_table", "add_image"]},
                            "params": {"type": "object"}
                        }
                    },
                    "description": "replace_text {old,new}; add_paragraph {text,style?}; update_table_cell {table,row,col,text}; add_heading {text,level?}; add_page_break {}; add_table {rows:[[...]]}; add_image {image_path,width_inches?}."
                }
            },
            "required": ["input_path", "output_path", "operations"]
        }
    },
    {
        "name": "read_pptx",
        "description": "Extract text, slide content, speaker notes, and metadata from a PPTX file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pptx_path": {"type": "string", "description": "Path to the PPTX file"},
                "extract": {
                    "type": "string",
                    "enum": ["text", "slides", "notes", "all"],
                    "default": "all",
                    "description": "What to extract"
                }
            },
            "required": ["pptx_path"]
        }
    },
    {
        "name": "modify_pptx",
        "description": "Modify an existing PPTX file. Can replace text across slides, add slides, update shape text, set speaker notes, and save back.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Path to the source PPTX"},
                "output_path": {"type": "string", "description": "Where to save the modified PPTX"},
                "operations": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "action": {"type": "string", "enum": ["replace_text", "add_slide", "update_slide_text", "set_notes", "add_textbox", "add_image", "delete_slide"]},
                            "params": {"type": "object"}
                        }
                    },
                    "description": "replace_text {old,new}; add_slide {layout_index?,title?,content?}; update_slide_text {slide,shape_name,text}; set_notes {slide,text}; add_textbox {slide,text,left?,top?,width?,height?}; add_image {slide,image_path,left?,top?,width?,height?}; delete_slide {slide}."
                }
            },
            "required": ["input_path", "output_path", "operations"]
        }
    },
    {
        "name": "read_xlsx",
        "description": "Extract cell data, sheet names, and formulas from an XLSX file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "xlsx_path": {"type": "string", "description": "Path to the XLSX file"},
                "sheet": {"type": "string", "description": "Sheet name or index. Omit for all sheets."},
                "range": {"type": "string", "description": "Cell range, e.g., 'A1:D10'. Omit for all data."},
                "extract": {
                    "type": "string",
                    "enum": ["values", "formulas", "all"],
                    "default": "values",
                    "description": "What to extract"
                }
            },
            "required": ["xlsx_path"]
        }
    },
    {
        "name": "modify_xlsx",
        "description": "Modify an existing XLSX file. Can update cells, set formulas, add rows/sheets, delete sheets, and save back.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Path to the source XLSX"},
                "output_path": {"type": "string", "description": "Where to save the modified XLSX"},
                "operations": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "action": {"type": "string", "enum": ["set_cell", "set_formula", "add_row", "add_sheet", "delete_sheet", "rename_sheet", "insert_row", "delete_row", "insert_column", "delete_column"]},
                            "params": {"type": "object"}
                        }
                    },
                    "description": "set_cell {sheet?,cell,value}; set_formula {sheet?,cell,formula}; add_row {sheet?,values}; add_sheet/delete_sheet; rename_sheet {old_name,new_name}; insert/delete row/column {sheet?,index,amount?}."
                }
            },
            "required": ["input_path", "output_path", "operations"]
        }
    },
    {
        "name": "create_zip",
        "description": "Create a ZIP archive from one or more files. Supports deflated (compressed) and stored (no compression) modes.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "Where to save the ZIP file"},
                "file_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "List of file or directory paths to include recursively in the archive"
                },
                "compression": {
                    "type": "string",
                    "enum": ["deflated", "stored"],
                    "default": "deflated",
                    "description": "Compression method: 'deflated' (smaller) or 'stored' (no compression, faster)"
                }
            },
            "required": ["output_path", "file_paths"]
        }
    },
    {
        "name": "download_media",
        "description": "Download video/audio from supported platforms (TikTok, Instagram, etc.). NOT YouTube.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "URL of the media"},
                "format": {
                    "type": "string",
                    "enum": ["video", "audio"],
                    "description": "Whether to download video or audio only"
                }
            },
            "required": ["url"]
        }
    },
    {
        "name": "ffmpeg_process",
        "description": "Process video/audio with FFmpeg. Operations: trim, crop, resize, filter (brightness/contrast/etc), custom (raw FFmpeg args for complex ops), extract_audio, extract_frame, convert.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Path to input file"},
                "output_path": {"type": "string", "description": "Path for output file"},
                "operation": {
                    "type": "string",
                    "enum": ["trim", "crop", "resize", "filter", "custom", "extract_audio", "extract_frame", "convert"],
                    "description": "Operation to perform"
                },
                "params": {
                    "type": "object",
                    "description": "Operation params. trim: {start, end} or {start, duration}. crop: {width, height, x, y}. resize: {width, height}. filter: {vf} for video filters (e.g. 'eq=brightness=0.3'), {af} for audio filters. custom: {args} raw FFmpeg arguments between -i input and output (for complex multi-filter chains, e.g. \"-vf select='not(mod(floor(t)\\,2))',setpts=N/FRAME_RATE/TB -af aselect='not(mod(floor(t)\\,2))',asetpts=N/SR/TB -c:v libx264 -crf 23 -c:a aac\"). extract_audio: {format, bitrate}. extract_frame: {timestamp}. convert: {codec, quality (int 0-51)}."
                }
            },
            "required": ["input_path", "output_path", "operation"]
        }
    },
    {
        "name": "ocr_image",
        "description": "Extract text from an image using OCR.",
        "input_schema": {
            "type": "object",
            "properties": {
                "image_path": {"type": "string", "description": "Path to image file"}
            },
            "required": ["image_path"]
        }
    },
    {
        "name": "smart_crop",
        "description": "Smart crop video/image to focus on faces. ONLY use for simple face-centered cropping. For TikTok/Reels adaptation with effects, transitions, or custom crop positions, use ffmpeg_process instead.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Path to input video or image"},
                "output_path": {"type": "string", "description": "Path for output file"},
                "aspect_ratio": {
                    "type": "string",
                    "default": "9:16",
                    "description": "Target aspect ratio (e.g., '9:16' for vertical, '16:9' for horizontal)"
                }
            },
            "required": ["input_path", "output_path"]
        }
    },
    {
        "name": "google_calendar",
        "description": "Query or create Google Calendar events. Requires user authorization.",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["list", "create", "delete", "update"],
                    "description": "Action to perform"
                },
                "date_range": {
                    "type": "string",
                    "description": "For list: 'today', 'this_week', or ISO date range"
                },
                "event": {
                    "type": "object",
                    "description": "For create: {title, start, end, description, location?}"
                },
                "event_id": {"type": "string", "description": "For delete: Calendar event ID"}
            },
            "required": ["action"]
        }
    },
    {
        "name": "gmail",
        "description": "Read Gmail messages (read-only). Requires user authorization.",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["list", "read"],
                    "description": "Action to perform"
                },
                "query": {
                    "type": "string",
                    "description": "For list: Gmail search query"
                },
                "message_id": {
                    "type": "string",
                    "description": "For read: message ID"
                }
            },
            "required": ["action"]
        }
    },
    {
        "name": "file_info",
        "description": "Get file metadata (size, name, extension). Use this instead of trying os.path in python_execute.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to the file"}
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "read_file",
        "description": "Read text content from a file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to the file to read"}
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "write_file",
        "description": "Write text content to a file. The created file will be available for download/sharing.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "Filename or path for the output file"},
                "content": {"type": "string", "description": "Text content to write to the file"}
            },
            "required": ["output_path", "content"]
        }
    },
    {
        "name": "python_execute",
        "description": """Execute Python code in a secure sandbox. Use this for:
- Data processing and analysis (pandas DataFrames, CSV, groupby, etc.)
- Mathematical calculations and algorithms
- Statistical analysis (numpy, statistics)
- Charts and plots (matplotlib — figures are auto-saved as PNG and returned)
- Text manipulation and parsing
- JSON/CSV data processing
- Any computation that requires custom logic

Available modules: math, numpy, pandas, matplotlib, json, re, datetime, collections, itertools, statistics, csv, base64, hashlib.
FORBIDDEN: subprocess, os, sys, shutil, socket, http, urllib, pathlib, glob, signal, ctypes, requests, multiprocessing, threading.
To run FFmpeg/FFprobe, use the ffmpeg_process tool instead. To access files, use dedicated tools (read_pdf, ocr_image, etc.).

For plots, use matplotlib — figures are auto-saved as PNG files and returned. An OUTPUT_DIR variable is available for saving files explicitly (e.g., df.to_csv(OUTPUT_DIR + '/data.csv')).

The code runs with a 30-second timeout. Print statements and the last expression's value are captured and returned.""",
        "input_schema": {
            "type": "object",
            "properties": {
                "code": {
                    "type": "string",
                    "description": "Python code to execute. Can be multi-line. Use print() for output."
                },
                "file_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional list of file paths the code is allowed to read (must be files provided by user)"
                }
            },
            "required": ["code"]
        }
    }
]


# RASTACODER_V7_COMPLETE_SKILLS
# Restore post-Qwen3 upstream tools which were omitted by the old 23-tool
# baseline and add structured scene-complete primitives requested for v7.
TOOLS_SCHEMA.extend([
    {
        "name": "list_files",
        "description": "List/discover files and directories. Supports app output and common Android folders, optional recursive traversal and filename pattern filtering.",
        "input_schema": {
            "type": "object",
            "properties": {
                "directory": {"type": "string", "enum": ["output", "downloads", "documents", "pictures", "screenshots", "camera"], "default": "output"},
                "path": {"type": "string", "description": "Optional explicit accessible directory path"},
                "recursive": {"type": "boolean", "default": False},
                "pattern": {"type": "string", "description": "Optional glob such as *.txt"},
                "include_directories": {"type": "boolean", "default": True},
            },
            "required": [],
        },
    },
    {
        "name": "file_manage",
        "description": "Manage files/directories: list, mkdir, copy, move, rename, delete, touch, exists. Use structured paths; no shell commands are needed.",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["list", "mkdir", "copy", "move", "rename", "delete", "touch", "exists"]},
                "path": {"type": "string"},
                "source_path": {"type": "string"},
                "destination_path": {"type": "string"},
                "recursive": {"type": "boolean", "default": False},
                "overwrite": {"type": "boolean", "default": False},
            },
            "required": ["action"],
        },
    },
    {
        "name": "list_zip",
        "description": "List files and metadata inside a ZIP archive without extracting it.",
        "input_schema": {
            "type": "object",
            "properties": {"zip_path": {"type": "string"}},
            "required": ["zip_path"],
        },
    },
    {
        "name": "extract_zip",
        "description": "Safely extract a ZIP archive into a writable directory. Rejects path traversal entries.",
        "input_schema": {
            "type": "object",
            "properties": {
                "zip_path": {"type": "string"},
                "output_dir": {"type": "string", "description": "Optional folder name/path; generated automatically when omitted"},
                "overwrite": {"type": "boolean", "default": False},
            },
            "required": ["zip_path"],
        },
    },
    {
        "name": "pdf_manage",
        "description": "Manage PDF pages: merge PDFs, split into individual pages, extract/reorder/delete pages, or rotate selected pages.",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": ["merge", "split", "extract_pages", "reorder", "delete_pages", "rotate"]},
                "input_path": {"type": "string"},
                "input_paths": {"type": "array", "items": {"type": "string"}},
                "output_path": {"type": "string"},
                "pages": {"description": "1-based pages, e.g. '1-3,5' or [3,1,2] for reorder"},
                "rotation": {"type": "integer", "default": 90},
            },
            "required": ["action"],
        },
    },
    {
        "name": "create_pptx",
        "description": "Create a PowerPoint PPTX from structured slides with titles, content/bullets and optional speaker notes.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "title": {"type": "string"},
                "slides": {"type": "array", "items": {"type": "object"}},
            },
            "required": ["output_path"],
        },
    },
    {
        "name": "create_xlsx",
        "description": "Create an Excel XLSX workbook from structured sheets and rows.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "sheets": {"type": "array", "items": {"type": "object"}, "description": "[{name, rows:[[...], ...]}]"},
            },
            "required": ["output_path"],
        },
    },
    {
        "name": "image_compose",
        "description": "Full image manipulation: horizontal/vertical concat, overlay, resize/upscale/downscale, color adjustment, crop, grayscale, blur, rotate, flip and format conversion.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_paths": {"type": "array", "items": {"type": "string"}},
                "output_path": {"type": "string"},
                "operation": {"type": "string", "enum": ["concat_horizontal", "concat_vertical", "overlay", "resize", "adjust", "crop", "grayscale", "blur", "rotate", "flip", "convert"]},
                "params": {"type": "object", "description": "resize {width,height}; crop {x,y,width,height}; adjust {brightness,contrast,saturation,sharpness,gamma}; rotate {degrees}; flip {direction}; convert {format,quality}."},
            },
            "required": ["input_paths", "output_path", "operation"],
        },
    },
])


# Compact tool schemas for offline (on-device) models with small context windows.
# Includes all tools that work offline. Descriptions kept minimal to conserve tokens.
OFFLINE_TOOLS_SCHEMA = [
    {
        "name": "python_execute",
        "description": "Run Python code. Available: math, numpy, pandas, matplotlib, json, re, datetime, collections, itertools, statistics, csv, base64, hashlib. Use print() for output. FORBIDDEN: os, sys, subprocess.",
        "input_schema": {
            "type": "object",
            "properties": {
                "code": {"type": "string", "description": "Python code to run"},
                "file_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Files the code can read"
                }
            },
            "required": ["code"]
        }
    },
    {
        "name": "ffmpeg_process",
        "description": "Process video/audio with FFmpeg. Operations: trim, crop, resize, filter, extract_audio, extract_frame, convert. Use this for ALL video/audio processing.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Input file path"},
                "output_path": {"type": "string", "description": "Output file path"},
                "operation": {
                    "type": "string",
                    "enum": ["trim", "crop", "resize", "filter", "extract_audio", "extract_frame", "convert"],
                    "description": "Operation to perform"
                },
                "params": {
                    "type": "object",
                    "description": "trim: {start, end/duration}. crop: {width, height, x, y}. resize: {width, height}. filter: {vf, af}. extract_audio: {format, bitrate}. extract_frame: {timestamp}. convert: {codec, quality}."
                }
            },
            "required": ["input_path", "output_path", "operation"]
        }
    },
    {
        "name": "smart_crop",
        "description": "Smart crop video/image to focus on faces. For simple face-centered cropping only.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Input file path"},
                "output_path": {"type": "string", "description": "Output file path"},
                "aspect_ratio": {"type": "string", "default": "9:16", "description": "Target aspect ratio"}
            },
            "required": ["input_path", "output_path"]
        }
    },
    {
        "name": "ocr_image",
        "description": "Extract text from an image using OCR.",
        "input_schema": {
            "type": "object",
            "properties": {
                "image_path": {"type": "string", "description": "Path to image"}
            },
            "required": ["image_path"]
        }
    },
    {
        "name": "read_pdf",
        "description": "Extract text from a PDF file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pdf_path": {"type": "string", "description": "Path to PDF"},
                "pages": {"type": "string", "description": "Page range, e.g. '1-5' or 'all'"}
            },
            "required": ["pdf_path"]
        }
    },
    {
        "name": "create_pdf",
        "description": "Create a PDF from text and/or images.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "Where to save PDF"},
                "content": {"type": "string", "description": "Text content"},
                "title": {"type": "string", "description": "Document title"},
                "image_paths": {"type": "array", "items": {"type": "string"}, "description": "Images to embed"}
            },
            "required": ["output_path"]
        }
    },
    {
        "name": "create_zip",
        "description": "Create ZIP archive from files.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "Where to save ZIP"},
                "file_paths": {"type": "array", "items": {"type": "string"}, "description": "Files to include"},
                "compression": {"type": "string", "enum": ["deflated", "stored"], "default": "deflated"}
            },
            "required": ["output_path", "file_paths"]
        }
    },
    {
        "name": "convert_document",
        "description": "Text-oriented conversion among TXT/DOCX/PDF/HTML.",
        "input_schema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Input file path"},
                "output_format": {"type": "string", "enum": ["pdf", "html", "txt", "docx"], "description": "Target format"},
                "output_path": {"type": "string", "description": "Optional output path"}
            },
            "required": ["input_path", "output_format"]
        }
    },
    {
        "name": "create_docx",
        "description": "Create DOCX Word document from text.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "content": {"type": "string"},
                "title": {"type": "string"}
            },
            "required": ["output_path", "content"]
        }
    },
    {
        "name": "read_docx",
        "description": "Extract text and tables from a DOCX file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "docx_path": {"type": "string", "description": "Path to DOCX"},
                "extract": {"type": "string", "enum": ["text", "tables", "all"], "default": "all"}
            },
            "required": ["docx_path"]
        }
    },
    {
        "name": "read_pptx",
        "description": "Extract text and slide content from a PPTX file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pptx_path": {"type": "string", "description": "Path to PPTX"},
                "extract": {"type": "string", "enum": ["text", "slides", "notes", "all"], "default": "all"}
            },
            "required": ["pptx_path"]
        }
    },
    {
        "name": "read_xlsx",
        "description": "Extract data from an XLSX spreadsheet.",
        "input_schema": {
            "type": "object",
            "properties": {
                "xlsx_path": {"type": "string", "description": "Path to XLSX"},
                "sheet": {"type": "string", "description": "Sheet name or index"},
                "range": {"type": "string", "description": "Cell range, e.g. 'A1:D10'"}
            },
            "required": ["xlsx_path"]
        }
    },
    {
        "name": "read_file",
        "description": "Read text from a file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "File to read"}
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "write_file",
        "description": "Write text to a file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "Output filename"},
                "content": {"type": "string", "description": "Content to write"}
            },
            "required": ["output_path", "content"]
        }
    },
    {
        "name": "file_info",
        "description": "Get file size, name, extension.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "File to inspect"}
            },
            "required": ["file_path"]
        }
    },
]


# RASTACODER_V8_SEARCH_SKILLS
_V8_SEARCH_TOOL_SCHEMAS = [
    {
        "name": "anysearch_search",
        "description": "Search the web with AnySearch. API credential is configured by the user in Tool Management.",
        "input_schema": {"type": "object", "properties": {
            "query": {"type": "string"}, "max_results": {"type": "integer"},
            "domain": {"type": "string"}, "sub_domain": {"type": "string"},
            "sub_domain_params": {"type": "object"}}, "required": ["query"]},
    },
    {
        "name": "anysearch_extract",
        "description": "Extract readable content from one web URL with AnySearch.",
        "input_schema": {"type": "object", "properties": {"url": {"type": "string"}}, "required": ["url"]},
    },
    {
        "name": "anysearch_get_sub_domains",
        "description": "List supported AnySearch sub-domains for one or more domains.",
        "input_schema": {"type": "object", "properties": {
            "domain": {"type": "string"}, "domains": {"type": "array", "items": {"type": "string"}}}},
    },
    {
        "name": "exa_search",
        "description": "Search the web with Exa and optionally return text, summaries, or highlights.",
        "input_schema": {"type": "object", "properties": {
            "query": {"type": "string"}, "num_results": {"type": "integer"},
            "topic": {"type": "string", "enum": ["general", "news"]},
            "search_type": {"type": "string", "enum": ["auto", "neural", "fast", "deep"]},
            "start_published_date": {"type": "string"},
            "include_domains": {"type": "array", "items": {"type": "string"}},
            "exclude_domains": {"type": "array", "items": {"type": "string"}},
            "include_text": {"type": "boolean"}, "include_summary": {"type": "boolean"},
            "include_highlights": {"type": "boolean"}}, "required": ["query"]},
    },
    {
        "name": "langsearch_search",
        "description": "Search the web with LangSearch.",
        "input_schema": {"type": "object", "properties": {
            "query": {"type": "string"}, "count": {"type": "integer"},
            "freshness": {"type": "string", "enum": ["oneDay", "oneWeek", "oneMonth", "oneYear", "noLimit"]},
            "summary": {"type": "boolean"}}, "required": ["query"]},
    },
    {
        "name": "tavily_search",
        "description": "Search the web with Tavily.",
        "input_schema": {"type": "object", "properties": {
            "query": {"type": "string"}, "max_results": {"type": "integer"},
            "topic": {"type": "string", "enum": ["general", "news"]},
            "search_depth": {"type": "string", "enum": ["basic", "advanced"]},
            "include_answer": {"type": "boolean"}, "time_range": {"type": "string", "enum": ["", "day", "week", "month", "year"]},
            "include_domains": {"type": "array", "items": {"type": "string"}},
            "exclude_domains": {"type": "array", "items": {"type": "string"}},
            "include_raw_content": {"type": "boolean"}}, "required": ["query"]},
    },
]
_existing_tool_names = {t["name"] for t in TOOLS_SCHEMA}
TOOLS_SCHEMA.extend(t for t in _V8_SEARCH_TOOL_SCHEMAS if t["name"] not in _existing_tool_names)
_existing_offline_names = {t["name"] for t in OFFLINE_TOOLS_SCHEMA}
OFFLINE_TOOLS_SCHEMA.extend(t for t in _V8_SEARCH_TOOL_SCHEMAS if t["name"] not in _existing_offline_names)


# RASTACODER_V9_SEARCH_INTENT_ONLY
for _schema_list in (TOOLS_SCHEMA, OFFLINE_TOOLS_SCHEMA):
    for _tool in _schema_list:
        if _tool.get("name") in {"anysearch_search", "exa_search", "langsearch_search", "tavily_search"}:
            _tool["input_schema"] = {"type": "object", "properties": {"query": {"type": "string", "description": "Search keywords or natural-language search question"}}, "required": ["query"]}
            _tool["description"] = "Search the web. Only supply query; result count/type/filter settings are configured by the user."


# RASTACODER_V11_CANONICAL_LIST_FILES
# One model-facing path concept prevents directory/path ambiguity. Common Android
# roots are addressed as path prefixes (downloads/, documents/, pictures/, etc.).
for _schema_list in (TOOLS_SCHEMA, OFFLINE_TOOLS_SCHEMA):
    for _tool in _schema_list:
        if _tool.get("name") == "list_files":
            _tool["description"] = (
                "List files/directories. path is relative to the app workspace by default; "
                "use downloads/, documents/, pictures/, screenshots/, or camera/ for common Android folders."
            )
            _tool["input_schema"] = {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "default": ".", "description": "Logical folder path. Use exactly '.' for workspace root; nested workspace paths are relative such as folder/sub"},
                    "recursive": {"type": "boolean", "default": False},
                    "pattern": {"type": "string", "description": "Optional glob such as *.pptx"},
                    "include_directories": {"type": "boolean", "default": True},
                },
                "required": [],
            }

# RASTACODER_V7_COMPLETE_SKILLS
# Every structured v7 utility is available to the local model when its Skill is
# enabled. Keep the schema gated by Skills rather than dumping all tools into
# every request.
_V7_LOCAL_TOOL_NAMES = {
    "list_files", "file_manage", "list_zip", "extract_zip", "pdf_manage",
    "create_pptx", "create_xlsx", "image_compose",
}
_existing_offline_names = {t["name"] for t in OFFLINE_TOOLS_SCHEMA}
OFFLINE_TOOLS_SCHEMA.extend(
    t for t in TOOLS_SCHEMA
    if t["name"] in _V7_LOCAL_TOOL_NAMES and t["name"] not in _existing_offline_names
)

# The native executor already supports raw/custom FFmpeg. V6 accidentally hid
# that escape hatch from the compact local schema. V7 also adds structured
# multi-input operations so common concat/mix tasks do not require raw syntax.
for _schema_list in (TOOLS_SCHEMA, OFFLINE_TOOLS_SCHEMA):
    for _tool in _schema_list:
        if _tool.get("name") != "ffmpeg_process":
            continue
        _props = _tool["input_schema"]["properties"]
        _props["input_paths"] = {
            "type": "array", "items": {"type": "string"},
            "description": "Multiple input media files for concat, mix_audio or merge_av"
        }
        _op = _props["operation"]
        _op["enum"] = [
            "trim", "crop", "resize", "filter", "custom", "extract_audio",
            "extract_frame", "convert", "concat", "mix_audio", "merge_av", "speed"
        ]
        _op["description"] = "Structured media operation; custom is the advanced raw FFmpeg escape hatch"
        _props["params"]["description"] = (
            "trim {start,end/duration}; crop {width,height,x,y}; resize {width,height}; "
            "filter {vf,af}; custom {args}; extract_audio {format,bitrate}; "
            "extract_frame {timestamp}; convert {codec,quality}; concat {media_type:audio|video}; "
            "mix_audio {duration}; merge_av uses first input as video and second as audio; "
            "speed {factor} changes playback speed, for example factor=1.5."
        )
        # input_path is operation-dependent now; executor performs the precise check.
        _tool["input_schema"]["required"] = ["output_path", "operation"]


# RASTACODER_V7_MEDIA_DOWNLOAD
for _tool in TOOLS_SCHEMA:
    if _tool.get("name") == "download_media":
        _tool["description"] = "Resolve and actually download video/audio from supported platforms into the app workspace (not YouTube)."
        _tool["input_schema"]["properties"]["output_path"] = {"type": "string", "description": "Optional output filename/path"}
        break

# Local inference and tool locality are independent. Expose these app-side
# capabilities to on-device models too; each tool still enforces its own
# connectivity/auth requirements.
_LOCAL_EXTRA_TOOL_NAMES = {
    "web_fetch", "headless_browser", "download_media",
    "modify_docx", "modify_pptx", "modify_xlsx",
    "google_calendar", "gmail",
}
_existing_offline_names = {t["name"] for t in OFFLINE_TOOLS_SCHEMA}
OFFLINE_TOOLS_SCHEMA.extend(
    t for t in TOOLS_SCHEMA
    if t["name"] in _LOCAL_EXTRA_TOOL_NAMES and t["name"] not in _existing_offline_names
)


# RASTACODER_V6_TOOL_RELIABILITY
# Skill IDs are UI-only. They are deliberately never shown to the model.
# The model sees canonical callable function names only.
LOCAL_SKILLS = {
    "text_files": {"tools": ("read_file", "write_file", "file_info", "list_files", "file_manage")},
    "zip_archive": {"tools": ("create_zip", "list_zip", "extract_zip", "file_info", "list_files", "file_manage")},
    "pdf_read": {"tools": ("read_pdf", "pdf_manage", "file_info", "list_files")},
    "pdf_create": {"tools": ("create_pdf", "pdf_manage", "image_compose", "file_info", "list_files")},
    "document_convert": {"tools": ("convert_document", "read_file", "read_pdf", "read_docx", "file_info", "list_files")},
    "word": {"tools": ("create_docx", "read_docx", "modify_docx", "convert_document", "file_info", "list_files", "file_manage")},
    "powerpoint": {"tools": ("create_pptx", "read_pptx", "modify_pptx", "file_info", "list_files", "file_manage")},
    "excel": {"tools": ("create_xlsx", "read_xlsx", "modify_xlsx", "file_info", "list_files", "file_manage")},
    "ocr": {"tools": ("ocr_image", "image_compose", "file_info", "list_files")},
    "image_processing": {"tools": ("image_compose", "smart_crop", "file_info", "list_files", "file_manage")},
    "video_processing": {"tools": ("ffmpeg_process", "file_info", "list_files", "file_manage")},
    "audio_processing": {"tools": ("ffmpeg_process", "file_info", "list_files", "file_manage")},
    "media_download": {"tools": ("download_media", "file_info", "list_files", "file_manage")},
    "web_fetch": {"tools": ("web_fetch", "write_file", "file_info", "list_files")},
    "dynamic_web": {"tools": ("headless_browser", "web_fetch", "write_file", "file_info")},
    "basic_calculation": {"tools": ("python_execute", "read_file", "write_file", "file_info")},
    "scientific_calculation": {"tools": ("python_execute", "read_file", "write_file", "file_info")},
    "data_analysis": {"tools": ("python_execute", "read_file", "write_file", "read_xlsx", "create_xlsx", "file_info", "list_files")},
    "charts": {"tools": ("python_execute", "write_file", "image_compose", "file_info", "list_files")},
    "gmail": {"tools": ("gmail",)},
    "google_calendar": {"tools": ("google_calendar",)},
    "anysearch_search": {"tools": ("anysearch_search", "anysearch_extract", "anysearch_get_sub_domains")},
    "exa_search": {"tools": ("exa_search",)},
    "langsearch_search": {"tools": ("langsearch_search",)},
    "tavily_search": {"tools": ("tavily_search",)},
}

ALL_LOCAL_SKILL_IDS = tuple(LOCAL_SKILLS.keys())

LOCAL_TOOL_PROMPT_HINTS = {
    "read_file": "read_file(file_path)",
    "write_file": "write_file(output_path, content)",
    "file_info": "file_info(file_path)",
    "list_files": "list_files(path='.', recursive=false, pattern=null, include_directories=true) ; use path='.' for workspace root",
    "file_manage": "file_manage(action, path, source_path, destination_path, recursive, overwrite) ; action=list|mkdir|copy|move|rename|delete|touch|exists",
    "create_zip": "create_zip(output_path, file_paths, compression)",
    "list_zip": "list_zip(zip_path)",
    "extract_zip": "extract_zip(zip_path, output_dir, overwrite)",
    "read_pdf": "read_pdf(pdf_path, pages)",
    "create_pdf": "create_pdf(output_path, content, title, image_paths)",
    "pdf_manage": "pdf_manage(action, input_path, input_paths, output_path, pages, rotation) ; action=merge|split|extract_pages|reorder|delete_pages|rotate",
    "convert_document": "convert_document(input_path, output_format, output_path) ; output_format=pdf|html|txt|docx",
    "create_docx": "create_docx(output_path, content, title)",
    "read_docx": "read_docx(docx_path) ; ordinary reads need only the file path",
    "modify_docx": "modify_docx(input_path, output_path, operations)",
    "create_pptx": "create_pptx(output_path, title, slides)",
    "read_pptx": "read_pptx(pptx_path) ; ordinary reads need only the file path",
    "modify_pptx": "modify_pptx(input_path, output_path, operations)",
    "create_xlsx": "create_xlsx(output_path, sheets)",
    "read_xlsx": "read_xlsx(xlsx_path, sheet, range) ; omit sheet/range for the whole workbook",
    "modify_xlsx": "modify_xlsx(input_path, output_path, operations)",
    "ocr_image": "ocr_image(image_path)",
    "image_compose": "image_compose(input_paths, output_path, operation, params) ; operation=concat_horizontal|concat_vertical|overlay|resize|adjust|crop|grayscale|blur|rotate|flip|convert",
    "smart_crop": "smart_crop(input_path, output_path, aspect_ratio)",
    "ffmpeg_process": "ffmpeg_process(input_path, input_paths, output_path, operation, params) ; operation=trim|crop|resize|filter|custom|extract_audio|extract_frame|convert|concat|mix_audio|merge_av|speed ; speed params={factor:1.5}",
    "download_media": "download_media(url, format)",
    "web_fetch": "web_fetch(url, extract_mode)",
    "headless_browser": "headless_browser(url, wait_seconds, extract_selector)",
    "python_execute": "python_execute(code, file_paths)",
    "gmail": "gmail(action, query, message_id) ; action=list|read",
    "google_calendar": "google_calendar(action, date_range, event, event_id) ; action=list|create|delete|update",
    "anysearch_search": "anysearch_search(query)",
    "anysearch_extract": "anysearch_extract(url)",
    "anysearch_get_sub_domains": "anysearch_get_sub_domains(domain or domains)",
    "exa_search": "exa_search(query)",
    "langsearch_search": "langsearch_search(query)",
    "tavily_search": "tavily_search(query)",
}


# RASTACODER_V13_SMALL_MODEL_TOOL_ABI
# Strict executor/cloud schemas remain unchanged. Local 3B-4B models receive a
# deep-copied projection which hides deterministic app-defaultable selectors.
_LOCAL_MODEL_HIDDEN_ARGS = {
    "read_docx": {"extract"},
    "read_pptx": {"extract"},
    "read_xlsx": {"extract"},
    "web_fetch": {"extract_mode"},
}


def get_local_tool_argument_classes():
    # Classify every local schema property for audit/fuzz coverage.
    result = {}
    for tool in OFFLINE_TOOLS_SCHEMA:
        schema = tool.get("input_schema") or {}
        props = schema.get("properties") or {}
        required = set(schema.get("required") or [])
        hidden = set(_LOCAL_MODEL_HIDDEN_ARGS.get(tool.get("name"), set()))
        classes = {}
        for key, spec in props.items():
            if key in required:
                classes[key] = "model_essential"
            elif key in hidden or (isinstance(spec, dict) and "default" in spec):
                classes[key] = "app_defaultable"
            else:
                classes[key] = "advanced_optional"
        result[str(tool.get("name"))] = classes
    return result


def _project_tool_for_local_model(tool):
    projected = copy.deepcopy(tool)
    name = str(projected.get("name") or "")
    schema = projected.get("input_schema") or {}
    props = schema.get("properties") or {}
    for key in _LOCAL_MODEL_HIDDEN_ARGS.get(name, set()):
        props.pop(key, None)
        required = schema.get("required")
        if isinstance(required, list) and key in required:
            required.remove(key)
    if name == "read_docx":
        projected["description"] = (
            "Read a DOCX file. Give only docx_path for an ordinary full read; "
            "the app chooses safe extraction defaults."
        )
    elif name == "read_pptx":
        projected["description"] = (
            "Read a PPTX file. Give only pptx_path for an ordinary full read; "
            "the app chooses safe extraction defaults."
        )
    elif name == "read_xlsx":
        projected["description"] = (
            "Read an XLSX workbook. Give xlsx_path; sheet/range are optional targeting controls."
        )
    elif name == "web_fetch":
        projected["description"] = "Fetch the readable text of a webpage. Give the URL."
    return projected


def _offline_tool_names():
    return {tool["name"] for tool in OFFLINE_TOOLS_SCHEMA}


def get_enabled_tool_names(skill_ids=None):
    if skill_ids is None:
        skill_ids = ALL_LOCAL_SKILL_IDS
    enabled = set()
    for skill_id in skill_ids:
        skill = LOCAL_SKILLS.get(str(skill_id))
        if skill:
            enabled.update(skill["tools"])
    return enabled


def get_offline_tools_for_skills(skill_ids=None):
    enabled = get_enabled_tool_names(skill_ids)
    return [
        _project_tool_for_local_model(tool)
        for tool in OFFLINE_TOOLS_SCHEMA
        if tool["name"] in enabled
    ]


def build_offline_skill_prompt(skill_ids=None):
    ids = ALL_LOCAL_SKILL_IDS if skill_ids is None else tuple(str(x) for x in skill_ids)
    selected = [skill_id for skill_id in ids if skill_id in LOCAL_SKILLS]
    base = (
        "You are RastaCoder, an AI assistant on Android. Tool availability is manually selected by the user. "
        "UI Skill/category labels are not callable functions and are intentionally omitted from this prompt."
    )
    if not selected:
        return base + " No tools are enabled. Answer directly and never emit a tool call."

    enabled_tools = get_enabled_tool_names(selected)
    ordered_tools = [t["name"] for t in OFFLINE_TOOLS_SCHEMA if t["name"] in enabled_tools]
    lines = [
        base,
        "When a tool is needed, output ONLY this XML wrapper with valid JSON inside:",
        "<tool_call>",
        '{"name":"CANONICAL_FUNCTION_NAME","arguments":{"exact_parameter_name":"value"}}',
        "</tool_call>",
        "CALLABLE FUNCTIONS (these exact names only):",
    ]
    for tool_name in ordered_tools:
        lines.append(f"- {LOCAL_TOOL_PROMPT_HINTS[tool_name]}")
    lines.extend([
        "STRICT TOOL-CALL RULES:",
        "- The name field MUST be one canonical function name listed above. Never call a Skill/category label.",
        "- arguments MUST use the exact parameter names shown in that function signature.",
        "- Never invent generic argument keys such as param, request, instruction, or command.",
        "- Use attached file basenames exactly as shown in the user message; the app resolves them to real paths.",
        "- WORKSPACE PATH RULE: use path='.' for the workspace root and relative paths like folder/file.txt below it. Do not invent Linux roots such as /workspace or /output.",
        "- Choose a sensible output filename yourself. Do not ask the user for an output path when a filename can be chosen safely.",
        "- Do not place prose before or after a tool call. After the tool result, give the concise final answer.",
        "- Use only the callable functions listed above.",
    ])
    return "\n".join(lines)


# Import-time invariant: the 21-skill catalogue must cover exactly every local
# tool exposed before classification. This intentionally fails fast in CI if a
# future tool is added without assigning it to at least one skill.
_skill_covered_tools = get_enabled_tool_names(ALL_LOCAL_SKILL_IDS)
_offline_tools = _offline_tool_names()
if _skill_covered_tools != _offline_tools:
    missing = sorted(_offline_tools - _skill_covered_tools)
    extra = sorted(_skill_covered_tools - _offline_tools)
    raise RuntimeError(f"Local skill coverage mismatch; missing={missing}, extra={extra}")



def _safe_diag_value(value: Any) -> Any:
    """Redact secrets and bound large diagnostic payloads."""
    secret_words = {"api_key", "access_token", "google_access_token", "authorization", "token", "password"}
    if isinstance(value, dict):
        out = {}
        for key, item in value.items():
            key_s = str(key)
            if key_s.lower() in secret_words or key_s == "_context":
                out[key_s] = "[REDACTED]"
            else:
                out[key_s] = _safe_diag_value(item)
        return out
    if isinstance(value, list):
        return [_safe_diag_value(v) for v in value[:50]]
    if isinstance(value, str) and len(value) > 2000:
        return value[:2000] + "...[truncated]"
    return value


def _record_tool_diag(context: Dict[str, Any], stage: str, **fields: Any) -> None:
    if not isinstance(context, dict):
        return
    events = context.setdefault('_diagnostics', [])
    if not isinstance(events, list):
        return
    event = {"stage": stage}
    event.update({k: _safe_diag_value(v) for k, v in fields.items()})
    events.append(event)


def _verify_output_artifact(path: Any) -> int:
    import os
    if not isinstance(path, str) or not path or not os.path.isfile(path):
        raise ToolError(f"[TOOL_POSTCONDITION_ERROR] Output file missing: {path}")
    size = os.path.getsize(path)
    if size <= 0:
        raise ToolError(f"[TOOL_POSTCONDITION_ERROR] Output file is empty: {path}")
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".zip":
            import zipfile
            with zipfile.ZipFile(path, "r") as zf:
                bad = zf.testzip()
                if bad is not None:
                    raise ToolError(f"[TOOL_POSTCONDITION_ERROR] Corrupt ZIP member: {bad}")
        elif ext == ".pdf":
            from pypdf import PdfReader
            if len(PdfReader(path).pages) < 1:
                raise ToolError("[TOOL_POSTCONDITION_ERROR] Generated PDF has no pages")
        elif ext == ".docx":
            from docx import Document
            Document(path)
        elif ext == ".pptx":
            from pptx import Presentation
            Presentation(path)
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=False)
            if not wb.sheetnames:
                wb.close(); raise ToolError("[TOOL_POSTCONDITION_ERROR] Generated XLSX has no sheets")
            wb.close()
        elif ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}:
            from PIL import Image
            with Image.open(path) as image:
                image.verify()
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"[TOOL_POSTCONDITION_ERROR] Cannot reopen generated artifact {path}: {exc}") from exc
    return size


def _verify_tool_result(tool_name: str, args: Dict[str, Any], result: Any) -> Any:
    if not isinstance(result, dict):
        return result
    if result.get("success") is False:
        raise ToolError(f"[TOOL_POSTCONDITION_ERROR] {tool_name} returned success=false")
    paths = []
    if isinstance(result.get("output_path"), str) and result.get("output_path"):
        paths.append(result["output_path"])
    if isinstance(result.get("output_paths"), list):
        paths.extend(p for p in result["output_paths"] if isinstance(p, str) and p)
    sizes = {p: _verify_output_artifact(p) for p in dict.fromkeys(paths)}
    if tool_name == "write_file" and paths:
        expected = str(args.get("content", ""))
        try:
            actual = Path(paths[0]).read_text(encoding="utf-8")
        except Exception as exc:
            raise ToolError(f"[TOOL_POSTCONDITION_ERROR] Cannot read back text output: {exc}") from exc
        if actual != expected:
            raise ToolError("[TOOL_POSTCONDITION_ERROR] Written text differs from requested content")
    if paths:
        result["verified_output"] = True
        result["verified_size_bytes"] = sum(sizes.values())
    return result


def execute_tool(
    tool_name: str,
    args: Dict[str, Any],
    context: Dict[str, Any]
) -> Any:
    """
    Execute a tool by name.

    Args:
        tool_name: Name of the tool to execute
        args: Tool arguments
        context: Execution context (tokens, etc.)

    Returns:
        Tool result

    Raises:
        ToolError: If tool execution fails
    """
    original_tool_name = tool_name
    tool_name, args, compatibility_notes = normalize_tool_call(tool_name, args, context=context)
    bridge = get_bridge()
    if compatibility_notes:
        bridge.log(
            "Tool compatibility: " + "; ".join(compatibility_notes),
            level="warn",
        )

    tool_map = {
        "web_fetch": web_fetch,
        "headless_browser": headless_browser,
        "read_pdf": read_pdf,
        "create_pdf": create_pdf,
        "create_docx": create_docx,
        "convert_document": convert_document,
        "read_docx": read_docx,
        "modify_docx": modify_docx,
        "read_pptx": read_pptx,
        "modify_pptx": modify_pptx,
        "read_xlsx": read_xlsx,
        "modify_xlsx": modify_xlsx,
        "create_zip": create_zip,
        "download_media": download_media,
        "google_calendar": google_calendar,
        "gmail": gmail,
        "ffmpeg_process": _ffmpeg_process,
        "ocr_image": _ocr_image,
        "smart_crop": _smart_crop,
        "python_execute": python_execute,
        "file_info": _file_info,
        "read_file": read_file,
        "write_file": write_file,
        "list_files": list_files,
        "file_manage": file_manage,
        "list_zip": list_zip,
        "extract_zip": extract_zip,
        "pdf_manage": pdf_manage,
        "create_pptx": create_pptx,
        "create_xlsx": create_xlsx,
        "image_compose": image_compose,
        "anysearch_search": anysearch_search,
        "anysearch_extract": anysearch_extract,
        "anysearch_get_sub_domains": anysearch_get_sub_domains,
        "exa_search": exa_search,
        "langsearch_search": langsearch_search,
        "tavily_search": tavily_search,
    }

    if tool_name not in tool_map:
        raise ToolError(
            f"[MODEL_TOOL_NAME_ERROR] Unknown tool '{original_tool_name}'. "
            f"Normalized name: '{tool_name}'."
        )

    # Manual Skill boundary comes before schema validation. A hallucinated
    # disabled tool must never be repaired into an executable call.
    allowed_tools = context.get('_allowed_tools')
    if allowed_tools is not None and tool_name not in set(allowed_tools):
        _record_tool_diag(context, "disabled", tool=tool_name, original_tool=original_tool_name)
        raise ToolError(
            f"[MODEL_TOOL_DISABLED] Tool '{tool_name}' is not enabled for this conversation."
        )

    _record_tool_diag(
        context, "normalized", tool=tool_name, original_tool=original_tool_name,
        args=_safe_diag_value(args), repairs=compatibility_notes,
    )

    # Validate required parameters and top-level enum values only after the
    # compatibility layer has synthesized deterministic safe defaults.
    schema_entry = next((t for t in TOOLS_SCHEMA if t.get("name") == tool_name), None)
    if schema_entry:
        input_schema = schema_entry.get("input_schema", {})
        missing = [
            key for key in input_schema.get("required", [])
            if key not in args or args.get(key) is None or args.get(key) == ""
        ]
        if missing:
            _record_tool_diag(context, "schema_error", tool=tool_name, missing=missing, args=_safe_diag_value(args))
            raise ToolError(
                f"[MODEL_TOOL_ARGUMENT_ERROR] {tool_name} missing required "
                f"parameter(s): {', '.join(missing)}. Received: {sorted(args.keys())}. "
                "Retry the same enabled tool with corrected arguments; choose a sensible output filename yourself when only output_path is missing."
            )
        for key, spec in input_schema.get("properties", {}).items():
            if key in args and isinstance(spec, dict) and spec.get("enum"):
                if args[key] not in spec["enum"]:
                    _record_tool_diag(context, "enum_error", tool=tool_name, key=key, value=args[key])
                    raise ToolError(
                        f"[MODEL_TOOL_ARGUMENT_ERROR] {tool_name}.{key} received "
                        f"{args[key]!r}; allowed values: {spec['enum']}. Retry with one allowed value."
                    )

    tool_func = tool_map[tool_name]

    # Resolve file paths: if a tool arg is a basename that matches an attached file,
    # replace it with the full path so native tools can find the file
    file_map = context.get('_file_map', {})
    if file_map:
        _resolve_file_paths(args, file_map)

    # Resolve every model-facing relative file path against the same workspace root.
    output_dir = context.get('output_dir')
    if output_dir:
        # RASTACODER_V12_PRESERVE_LIST_LOGICAL_PATH
        # list_files owns its logical-path -> physical-path translation via
        # resolve_list_path(_output_dir). Keeping its path logical here is
        # essential: requested_path is later returned to the model and must
        # remain '.', 'folder/sub', 'downloads/...', etc., never the private
        # Android/app filesystem root. All other tools keep the universal
        # input resolver because their implementations consume physical paths.
        if tool_name != 'list_files':
            _resolve_workspace_input_paths(args, output_dir)
        _resolve_output_paths(args, output_dir)

    _record_tool_diag(context, "paths_resolved", tool=tool_name, args=_safe_diag_value(args))

    # Add context to args for tools that need it
    if tool_name in [
        "google_calendar", "gmail", "anysearch_search", "anysearch_extract",
        "anysearch_get_sub_domains", "exa_search", "langsearch_search", "tavily_search",
    ]:
        args["_context"] = context

    # Pass output_dir to Python tools which need a stable workspace root.
    if tool_name == "python_execute" and output_dir:
        args["output_dir"] = output_dir
    if tool_name in {"list_files", "file_manage", "extract_zip", "pdf_manage", "download_media"} and output_dir:
        args["_output_dir"] = output_dir

    # Pass timeout for native tools
    if tool_name in ["ocr_image", "ffmpeg_process", "smart_crop"]:
        args["_timeout_ms"] = context.get("tool_timeout_ms", 30000)

    # Strip internal keys that Claude may echo back from context
    args.pop('_timeout_ms', None) if tool_name not in ["ocr_image", "ffmpeg_process", "smart_crop"] else None

    try:
        import inspect
        inspect.signature(tool_func).bind(**args)
    except TypeError as e:
        raise ToolError(f"[MODEL_TOOL_ARGUMENT_ERROR] {tool_name}: {str(e)}")

    result = tool_func(**args)
    return _verify_tool_result(tool_name, args, result)


def _resolve_file_paths(args: Dict[str, Any], file_map: Dict[str, str]) -> None:
    """Resolve basename references to full file paths using the attached files map."""
    import os
    path_keys = ['image_path', 'input_path', 'pdf_path', 'file_path', 'path', 'source_path', 'zip_path', 'docx_path', 'pptx_path', 'xlsx_path']
    for key in path_keys:
        if key in args:
            value = args[key]
            if isinstance(value, str):
                # Direct match (value is already a basename in the map)
                if value in file_map:
                    args[key] = file_map[value]
                # Try matching basename of a full path Claude may have guessed
                elif os.path.basename(value) in file_map:
                    args[key] = file_map[os.path.basename(value)]

    # Also resolve arrays of paths (e.g. image_paths for create_pdf, file_paths for create_zip)
    array_path_keys = ['image_paths', 'file_paths', 'input_paths']
    for key in array_path_keys:
        if key in args and isinstance(args[key], list):
            resolved = []
            for p in args[key]:
                if isinstance(p, str):
                    if p in file_map:
                        resolved.append(file_map[p])
                    elif os.path.basename(p) in file_map:
                        resolved.append(file_map[os.path.basename(p)])
                    else:
                        resolved.append(p)
                else:
                    resolved.append(p)
            args[key] = resolved

    # Office modification operations may carry attached image/file paths one
    # level deeper under operations[*].params. Resolve those basenames too.
    operations = args.get('operations')
    if isinstance(operations, list):
        for op in operations:
            if not isinstance(op, dict):
                continue
            params = op.get('params')
            if not isinstance(params, dict):
                continue
            for nested_key in ('image_path', 'file_path', 'source_path', 'input_path'):
                value = params.get(nested_key)
                if not isinstance(value, str):
                    continue
                if value in file_map:
                    params[nested_key] = file_map[value]
                elif os.path.basename(value) in file_map:
                    params[nested_key] = file_map[os.path.basename(value)]


# RASTACODER_V11_GLOBAL_WORKSPACE_PATHS
# RASTACODER_V12_CENTRAL_PATH_CONTRACT
def _workspace_relative_path(value: str, output_dir: str) -> str:
    return resolve_model_path(value, output_dir, allow_android_roots=True)


def _resolve_workspace_input_paths(args: Dict[str, Any], output_dir: str) -> None:
    path_keys = [
        'image_path', 'input_path', 'pdf_path', 'file_path', 'path', 'source_path',
        'zip_path', 'docx_path', 'pptx_path', 'xlsx_path',
    ]
    for key in path_keys:
        value = args.get(key)
        if isinstance(value, str):
            args[key] = _workspace_relative_path(value, output_dir)
    for key in ('image_paths', 'file_paths', 'input_paths'):
        values = args.get(key)
        if isinstance(values, list):
            args[key] = [
                _workspace_relative_path(v, output_dir) if isinstance(v, str) else v
                for v in values
            ]
    operations = args.get('operations')
    if isinstance(operations, list):
        for op in operations:
            if not isinstance(op, dict) or not isinstance(op.get('params'), dict):
                continue
            params = op['params']
            for key in ('image_path', 'file_path', 'source_path', 'input_path'):
                if isinstance(params.get(key), str):
                    params[key] = _workspace_relative_path(params[key], output_dir)


def _resolve_output_paths(args: Dict[str, Any], output_dir: str) -> None:
    """Resolve generated outputs through the same virtual-workspace contract."""
    import os
    os.makedirs(output_dir, exist_ok=True)
    value = args.get('output_path')
    if isinstance(value, str):
        args['output_path'] = resolve_output_path(value, output_dir)


def _file_info(file_path: str, **kwargs) -> dict:
    """Get file metadata (size, name, extension)."""
    import os
    if not os.path.exists(file_path):
        raise ToolError(f"File not found: {file_path}")
    size_bytes = os.path.getsize(file_path)
    return {
        "name": os.path.basename(file_path),
        "path": file_path,
        "size_bytes": size_bytes,
        "size_mb": round(size_bytes / (1024 * 1024), 2),
        "extension": os.path.splitext(file_path)[1].lstrip('.'),
    }


def _ffmpeg_process(**kwargs) -> dict:
    """FFmpeg processing - delegates to native Flutter tool."""
    from ..bridge import get_bridge
    bridge = get_bridge()

    timeout_ms = kwargs.pop('_timeout_ms', 30000)
    # FFmpeg gets 10x the base timeout (video processing is slow)
    return bridge.call_native("ffmpeg", kwargs, timeout_ms=timeout_ms * 10)


def _ocr_image(**kwargs) -> dict:
    """OCR - delegates to native ML Kit tool."""
    from ..bridge import get_bridge
    bridge = get_bridge()

    timeout_ms = kwargs.pop('_timeout_ms', 30000)
    return bridge.call_native("ocr", kwargs, timeout_ms=timeout_ms)


def _smart_crop(**kwargs) -> dict:
    """Smart crop with face detection - delegates to native tool."""
    from ..bridge import get_bridge
    bridge = get_bridge()

    timeout_ms = kwargs.pop('_timeout_ms', 30000)
    # Smart crop gets 10x the base timeout (video processing is slow)
    return bridge.call_native("smart_crop", kwargs, timeout_ms=timeout_ms * 10)
