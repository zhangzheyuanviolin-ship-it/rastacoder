"""
Document Tools - PDF, Office documents (DOCX/PPTX/XLSX), file I/O, and ZIP archive creation
"""

import os
import zipfile
from io import BytesIO
from typing import Optional

from ..bridge import ToolError
from ..utils.file_limits import validate_file_for_processing, validate_pdf_for_processing, PROCESSING_LIMITS


def read_pdf(pdf_path: str, pages: str = "all") -> dict:
    """
    Extract text from a PDF file.

    Args:
        pdf_path: Path to the PDF file
        pages: Page range ("all", "1-5", "3", etc.)

    Returns:
        Dict with extracted text
    """
    from pypdf import PdfReader

    # Validate file
    validate_pdf_for_processing(pdf_path)

    try:
        reader = PdfReader(pdf_path)
        total_pages = len(reader.pages)

        # Parse page range
        if pages == "all":
            page_indices = range(total_pages)
        elif "-" in pages:
            start, end = pages.split("-")
            start = int(start) - 1
            end = min(int(end), total_pages)
            page_indices = range(start, end)
        else:
            page_num = int(pages) - 1
            if page_num >= total_pages:
                raise ToolError(f"Page {pages} doesn't exist. PDF has {total_pages} pages.")
            page_indices = [page_num]

        # Extract text
        text_parts = []
        for i in page_indices:
            page = reader.pages[i]
            text = page.extract_text()
            if text:
                text_parts.append(f"--- Page {i + 1} ---\n{text}")

        full_text = "\n\n".join(text_parts)

        # Truncate if too long
        if len(full_text) > 100000:
            full_text = full_text[:100000] + "\n\n[Content truncated...]"

        return {
            "path": pdf_path,
            "total_pages": total_pages,
            "pages_extracted": len(page_indices),
            "text": full_text
        }

    except Exception as e:
        raise ToolError(f"Failed to read PDF: {str(e)}")


def read_file(file_path: str) -> dict:
    """
    Read text content from a file.

    Args:
        file_path: Path to the file to read

    Returns:
        Dict with path, content, and size_bytes
    """
    validate_file_for_processing(file_path)

    try:
        size_bytes = os.path.getsize(file_path)
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()

        # Truncate if too long
        if len(content) > 100000:
            content = content[:100000] + "\n\n[Content truncated...]"

        return {
            "path": file_path,
            "content": content,
            "size_bytes": size_bytes,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read file: {str(e)}")


def write_file(output_path: str, content: str) -> dict:
    """
    Write text content to a file.

    Args:
        output_path: Where to save the file
        content: Text content to write

    Returns:
        Dict with output_path, success, and size_bytes
    """
    max_chars = PROCESSING_LIMITS['text_chars']
    if len(content) > max_chars:
        raise ToolError(
            f"Content too large: {len(content)} chars. "
            f"Maximum: {max_chars} chars."
        )

    try:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)

        size_bytes = os.path.getsize(output_path)

        return {
            "output_path": output_path,
            "success": True,
            "size_bytes": size_bytes,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to write file: {str(e)}")


def create_pdf(
    output_path: str,
    content: Optional[str] = None,
    title: Optional[str] = None,
    image_paths: Optional[list] = None,
) -> dict:
    """
    Create a PDF from text and/or images.

    Args:
        output_path: Where to save the PDF
        content: Optional text content for the PDF
        title: Optional document title
        image_paths: Optional list of image file paths to embed

    Returns:
        Dict with output path and page count
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image

    if not content and not image_paths:
        raise ToolError("create_pdf requires at least 'content' or 'image_paths'")

    try:
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # Create document
        page_width, page_height = letter
        margin = 72
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=margin,
            leftMargin=margin,
            topMargin=margin,
            bottomMargin=margin,
        )

        usable_width = page_width - 2 * margin
        usable_height = page_height - 2 * margin

        # Get styles
        styles = getSampleStyleSheet()

        # Build content
        story = []

        # Add title if provided
        if title:
            title_style = ParagraphStyle(
                'Title',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
            )
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.25 * inch))

        # Add text content paragraphs
        if content:
            body_style = styles['Normal']
            for paragraph in content.split('\n\n'):
                if paragraph.strip():
                    # Escape special characters
                    safe_text = paragraph.replace('&', '&amp;')
                    safe_text = safe_text.replace('<', '&lt;')
                    safe_text = safe_text.replace('>', '&gt;')
                    story.append(Paragraph(safe_text, body_style))
                    story.append(Spacer(1, 0.1 * inch))

        # Embed images
        if image_paths:
            for img_path in image_paths:
                if not os.path.isfile(img_path):
                    raise ToolError(f"Image file not found: {img_path}")

                # Get image dimensions and scale to fit page
                from PIL import Image as PILImage
                with PILImage.open(img_path) as pil_img:
                    img_w, img_h = pil_img.size

                # Scale to fit within usable area while maintaining aspect ratio
                scale = min(usable_width / img_w, usable_height / img_h, 1.0)
                display_w = img_w * scale
                display_h = img_h * scale

                story.append(Image(img_path, width=display_w, height=display_h))
                story.append(Spacer(1, 0.2 * inch))

        # Build PDF
        doc.build(story)

        return {
            "output_path": output_path,
            "success": True,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to create PDF: {str(e)}")


# RASTACODER_V4_TOOL_CONTRACT

def create_docx(output_path: str, content: str, title: str = None) -> dict:
    """Create a DOCX document from plain text."""
    from docx import Document

    if content is None:
        raise ToolError("create_docx requires 'content'")
    if len(content) > PROCESSING_LIMITS['text_chars']:
        raise ToolError(
            f"Content too large: {len(content)} chars. "
            f"Maximum: {PROCESSING_LIMITS['text_chars']} chars."
        )

    try:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        doc = Document()
        if title:
            doc.add_heading(str(title), level=1)
        # Preserve ordinary line structure while allowing blank lines.
        for paragraph in str(content).splitlines():
            doc.add_paragraph(paragraph)
        if not str(content).splitlines():
            doc.add_paragraph(str(content))
        doc.save(output_path)
        return {
            "output_path": output_path,
            "success": True,
            "format": "docx",
        }
    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to create DOCX: {str(e)}")


def _extract_document_text(input_path: str) -> str:
    """Extract text for text-oriented format conversion."""
    ext = os.path.splitext(input_path)[1].lower()
    if ext == '.txt':
        with open(input_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    if ext == '.docx':
        from docx import Document
        doc = Document(input_path)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append('\t'.join(cell.text for cell in row.cells))
        return '\n'.join(parts)
    if ext == '.pdf':
        from pypdf import PdfReader
        reader = PdfReader(input_path)
        return '\n\n'.join((p.extract_text() or '') for p in reader.pages)
    if ext in ('.html', '.htm'):
        from bs4 import BeautifulSoup
        with open(input_path, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'lxml')
        return soup.get_text('\n', strip=True)
    if ext == '.doc':
        raise ToolError(
            "Legacy .doc files are not supported by the bundled python-docx library. "
            "Please convert the file to .docx first."
        )
    raise ToolError(
        f"Unsupported input format: {ext or '[no extension]'}. "
        "Supported inputs: .txt, .docx, .pdf, .html, .htm."
    )


def convert_document(
    input_path: str,
    output_format: str,
    output_path: str = None,
) -> dict:
    """Convert TXT/DOCX/PDF/HTML by extracting and recreating text content.

    Complex source layout, embedded media, formulas, and advanced Office styling
    may be simplified because this is a text-oriented mobile conversion tool.
    """
    validate_file_for_processing(input_path, 'document')

    fmt = str(output_format or '').strip().lower().lstrip('.')
    fmt = {
        'word': 'docx', 'msword': 'docx', 'microsoft word': 'docx',
        'text': 'txt', 'htm': 'html',
    }.get(fmt, fmt)
    if fmt not in {'pdf', 'html', 'txt', 'docx'}:
        raise ToolError(
            f"Unsupported output format: {output_format}. "
            "Supported outputs: pdf, html, txt, docx."
        )

    base_name = os.path.splitext(input_path)[0]
    output_path = output_path or f"{base_name}.{fmt}"

    try:
        text = _extract_document_text(input_path)
        if len(text) > PROCESSING_LIMITS['text_chars']:
            text = text[:PROCESSING_LIMITS['text_chars']]

        if fmt == 'txt':
            output_dir = os.path.dirname(output_path)
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            return {"output_path": output_path, "success": True, "format": "txt"}

        if fmt == 'docx':
            return create_docx(output_path=output_path, content=text)

        if fmt == 'pdf':
            return create_pdf(output_path=output_path, content=text)

        # HTML output: escape source text so arbitrary TXT/DOCX content does not
        # become executable markup.
        from html import escape
        paragraphs = text.split('\n\n')
        html_body = ''.join(f'<p>{escape(p)}</p>' for p in paragraphs if p.strip())
        html = (
            '<!DOCTYPE html>\n<html>\n<head><meta charset="utf-8">'
            '<title>Converted Document</title></head>\n<body>\n'
            f'{html_body}\n</body>\n</html>'
        )
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
        return {"output_path": output_path, "success": True, "format": "html"}

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Conversion failed: {str(e)}")


def create_zip(
    output_path: str,
    file_paths: list,
    compression: str = "deflated",
) -> dict:
    """
    Create a ZIP archive from a list of files.

    Uses Python's built-in zipfile module (PSF License 2.0 — fully free/open).

    Args:
        output_path: Where to save the ZIP file
        file_paths: List of file paths to include in the archive
        compression: Compression method — "deflated" (smaller, default) or "stored" (no compression, faster)

    Returns:
        Dict with output path, file count, and total size
    """
    if not file_paths:
        raise ToolError("create_zip requires at least one file path")

    compression_map = {
        "deflated": zipfile.ZIP_DEFLATED,
        "stored": zipfile.ZIP_STORED,
    }
    if compression not in compression_map:
        raise ToolError(
            f"Unsupported compression: {compression}. Use 'deflated' or 'stored'."
        )

    try:
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # Verify all inputs exist. V7 accepts both files and directories.
        for fpath in file_paths:
            if not os.path.exists(fpath):
                raise ToolError(f"File or directory not found: {fpath}")

        # Track top-level basenames to handle collisions.
        seen_names = {}
        zip_compression = compression_map[compression]

        with zipfile.ZipFile(output_path, 'w', compression=zip_compression) as zf:
            for fpath in file_paths:
                basename = os.path.basename(os.path.normpath(fpath)) or 'item'
                if basename in seen_names:
                    seen_names[basename] += 1
                    root_name = f"{basename}_{seen_names[basename]}"
                else:
                    seen_names[basename] = 0
                    root_name = basename

                if os.path.isdir(fpath):
                    wrote_any = False
                    for root, dirs, files in os.walk(fpath):
                        rel_root = os.path.relpath(root, fpath)
                        archive_root = root_name if rel_root == '.' else os.path.join(root_name, rel_root)
                        if not files and not dirs:
                            zf.writestr(archive_root.rstrip('/') + '/', b'')
                        for filename in files:
                            source = os.path.join(root, filename)
                            arcname = os.path.join(archive_root, filename)
                            zf.write(source, arcname)
                            wrote_any = True
                    if not wrote_any:
                        zf.writestr(root_name.rstrip('/') + '/', b'')
                else:
                    name, ext = os.path.splitext(root_name)
                    if seen_names[basename] > 0:
                        arcname = f"{name}{ext}"
                    else:
                        arcname = root_name
                    zf.write(fpath, arcname)

        # Get final archive size
        archive_size = os.path.getsize(output_path)

        return {
            "output_path": output_path,
            "success": True,
            "file_count": len(file_paths),
            "size_bytes": archive_size,
            "size_mb": round(archive_size / (1024 * 1024), 2),
            "compression": compression,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to create ZIP: {str(e)}")


# ---------------------------------------------------------------------------
# DOCX read/write
# ---------------------------------------------------------------------------

def read_docx(docx_path: str, extract: str = "all") -> dict:
    """
    Extract text, tables, and metadata from a DOCX file.

    Args:
        docx_path: Path to the DOCX file
        extract: What to extract — "text", "tables", or "all"

    Returns:
        Dict with extracted content
    """
    from docx import Document

    validate_file_for_processing(docx_path, 'document')

    try:
        doc = Document(docx_path)
        result = {"path": docx_path}

        if extract in ("text", "all"):
            paragraphs = [p.text for p in doc.paragraphs]
            full_text = "\n\n".join(p for p in paragraphs if p.strip())
            if len(full_text) > PROCESSING_LIMITS['text_chars']:
                full_text = full_text[:PROCESSING_LIMITS['text_chars']] + "\n\n[Content truncated...]"
            result["text"] = full_text
            result["paragraph_count"] = len(paragraphs)

        if extract in ("tables", "all"):
            tables_data = []
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                tables_data.append(rows)
            result["tables"] = tables_data
            result["table_count"] = len(tables_data)

        return result

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read DOCX: {str(e)}")


def modify_docx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify DOCX transactionally. Every requested operation must have a real effect."""
    from docx import Document

    validate_file_for_processing(input_path, "document")
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_docx requires a non-empty operations list")

    try:
        doc = Document(input_path)
        applied = 0
        verification = []

        def paragraphs_with_tables():
            for paragraph in doc.paragraphs:
                yield paragraph
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            yield paragraph

        for index, op in enumerate(operations):
            if not isinstance(op, dict):
                raise ToolError(f"DOCX operation #{index + 1} must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params") if isinstance(op.get("params"), dict) else {}

            if action == "replace_text":
                old = str(params.get("old") or "")
                new = str(params.get("new") or "")
                if not old:
                    raise ToolError("replace_text requires non-empty params.old")
                replacements = 0
                for para in paragraphs_with_tables():
                    if old in para.text:
                        replacements += para.text.count(old)
                        para.text = para.text.replace(old, new)
                if replacements == 0:
                    raise ToolError(f"replace_text found no matches for: {old!r}")
                verification.append(("replace_text", new, replacements))

            elif action == "add_paragraph":
                text = str(params.get("text") or "")
                if not text:
                    raise ToolError("add_paragraph requires non-empty params.text")
                style = params.get("style")
                doc.add_paragraph(text, style=style)
                verification.append(("paragraph", len(doc.paragraphs) - 1, text))

            elif action == "update_table_cell":
                table_idx = int(params.get("table", 0))
                row_idx = int(params.get("row", 0))
                col_idx = int(params.get("col", 0))
                if min(table_idx, row_idx, col_idx) < 0 or table_idx >= len(doc.tables):
                    raise ToolError("update_table_cell table/row/col is out of range")
                table = doc.tables[table_idx]
                if row_idx >= len(table.rows) or col_idx >= len(table.rows[row_idx].cells):
                    raise ToolError("update_table_cell table/row/col is out of range")
                text = str(params.get("text") or "")
                table.rows[row_idx].cells[col_idx].text = text
                verification.append(("cell", table_idx, row_idx, col_idx, text))

            elif action == "add_heading":
                text = str(params.get("text") or "")
                if not text:
                    raise ToolError("add_heading requires non-empty params.text")
                level = int(params.get("level", 1))
                if level < 0 or level > 9:
                    raise ToolError("add_heading level must be between 0 and 9")
                doc.add_heading(text, level=level)
                verification.append(("paragraph", len(doc.paragraphs) - 1, text))

            elif action == "add_page_break":
                before = len(doc.paragraphs)
                doc.add_page_break()
                if len(doc.paragraphs) <= before:
                    raise ToolError("add_page_break produced no structural change")

            elif action == "add_table":
                rows = params.get("rows") or []
                if not isinstance(rows, list) or not rows:
                    raise ToolError("add_table requires params.rows as a non-empty 2D list")
                col_count = max(len(row) if isinstance(row, list) else 1 for row in rows)
                before = len(doc.tables)
                table = doc.add_table(rows=len(rows), cols=col_count)
                for r, row in enumerate(rows):
                    values = row if isinstance(row, list) else [row]
                    for c, value in enumerate(values):
                        table.cell(r, c).text = str(value)
                verification.append(("table_count", before + 1))

            elif action == "add_image":
                image_path = params.get("image_path")
                if not image_path or not os.path.isfile(image_path):
                    raise ToolError(f"add_image image not found: {image_path}")
                before = len(doc.inline_shapes)
                width_inches = params.get("width_inches")
                if width_inches is None:
                    doc.add_picture(image_path)
                else:
                    from docx.shared import Inches
                    doc.add_picture(image_path, width=Inches(float(width_inches)))
                verification.append(("inline_shapes", before + 1))

            else:
                raise ToolError(f"Unknown DOCX modify action: {action}")
            applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        same_path = os.path.abspath(input_path) == os.path.abspath(output_path)
        save_path = output_path + ".rastacoder_tmp.docx" if same_path else output_path
        try:
            doc.save(save_path)
            check = Document(save_path)
            for item in verification:
                kind = item[0]
                if kind == "paragraph":
                    _, para_idx, expected = item
                    if para_idx >= len(check.paragraphs) or check.paragraphs[para_idx].text != expected:
                        raise ToolError("DOCX paragraph verification failed after save")
                elif kind == "cell":
                    _, t, r, c, expected = item
                    if check.tables[t].rows[r].cells[c].text != expected:
                        raise ToolError("DOCX table-cell verification failed after save")
                elif kind == "table_count" and len(check.tables) < item[1]:
                    raise ToolError("DOCX table verification failed after save")
                elif kind == "inline_shapes" and len(check.inline_shapes) < item[1]:
                    raise ToolError("DOCX image verification failed after save")
            if same_path:
                os.replace(save_path, output_path)
        finally:
            if same_path and os.path.exists(save_path):
                try:
                    os.remove(save_path)
                except OSError:
                    pass

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
            "operations_requested": len(operations),
            "verified": True,
            "in_place": same_path,
        }
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"Failed to modify DOCX: {exc}")


# ---------------------------------------------------------------------------
# PPTX read/write

# ---------------------------------------------------------------------------
# PPTX read/write
# ---------------------------------------------------------------------------

def read_pptx(pptx_path: str, extract: str = "all") -> dict:
    """
    Extract text, slide content, speaker notes, and metadata from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file
        extract: What to extract — "text", "slides", "notes", or "all"

    Returns:
        Dict with extracted content
    """
    from pptx import Presentation

    validate_file_for_processing(pptx_path, 'document')

    try:
        prs = Presentation(pptx_path)

        if len(prs.slides) > PROCESSING_LIMITS.get('pptx_slides', 500):
            raise ToolError(
                f"Presentation has too many slides ({len(prs.slides)}). "
                f"Maximum: {PROCESSING_LIMITS.get('pptx_slides', 500)}"
            )

        result = {"path": pptx_path, "slide_count": len(prs.slides)}
        slides_data = []

        for i, slide in enumerate(prs.slides):
            slide_info = {"slide_number": i + 1, "shapes": []}

            # Extract text from shapes
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                    if shape_text.strip():
                        texts.append(shape_text)
                        slide_info["shapes"].append({
                            "name": shape.name,
                            "text": shape_text,
                        })

                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        table_rows.append([cell.text for cell in row.cells])
                    slide_info["shapes"].append({
                        "name": shape.name,
                        "table": table_rows,
                    })

            slide_info["text"] = "\n".join(texts)

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_info["notes"] = slide.notes_slide.notes_text_frame.text
            else:
                slide_info["notes"] = ""

            slides_data.append(slide_info)

        if extract in ("text", "all"):
            all_text = []
            for sd in slides_data:
                all_text.append(f"--- Slide {sd['slide_number']} ---\n{sd['text']}")
            full_text = "\n\n".join(all_text)
            if len(full_text) > PROCESSING_LIMITS['text_chars']:
                full_text = full_text[:PROCESSING_LIMITS['text_chars']] + "\n\n[Content truncated...]"
            result["text"] = full_text

        if extract in ("slides", "all"):
            result["slides"] = slides_data

        if extract in ("notes", "all"):
            result["notes"] = [
                {"slide": sd["slide_number"], "notes": sd["notes"]}
                for sd in slides_data if sd["notes"]
            ]

        return result

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read PPTX: {str(e)}")


def modify_pptx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify PPTX transactionally and reject silent no-op operations."""
    from pptx import Presentation

    validate_file_for_processing(input_path, "document")
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_pptx requires a non-empty operations list")

    try:
        prs = Presentation(input_path)
        applied = 0
        expected_slide_count = len(prs.slides)
        verification = []

        for index, op in enumerate(operations):
            if not isinstance(op, dict):
                raise ToolError(f"PPTX operation #{index + 1} must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params") if isinstance(op.get("params"), dict) else {}

            if action == "replace_text":
                old = str(params.get("old") or "")
                new = str(params.get("new") or "")
                if not old:
                    raise ToolError("replace_text requires non-empty params.old")
                replacements = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in shape.text_frame.paragraphs:
                            if old in para.text:
                                replacements += para.text.count(old)
                                para.text = para.text.replace(old, new)
                if replacements == 0:
                    raise ToolError(f"replace_text found no matches for: {old!r}")

            elif action == "add_slide":
                layout_idx = int(params.get("layout_index", 1))
                if layout_idx < 0:
                    raise ToolError("add_slide layout_index must be non-negative")
                layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
                slide = prs.slides.add_slide(layout)
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 0 and params.get("title") is not None:
                        ph.text = str(params.get("title"))
                    elif ph.placeholder_format.idx == 1 and params.get("content") is not None:
                        ph.text = str(params.get("content"))
                expected_slide_count += 1

            elif action == "update_slide_text":
                slide_num = int(params.get("slide", 1)) - 1
                shape_name = str(params.get("shape_name") or "")
                if not (0 <= slide_num < len(prs.slides)) or not shape_name:
                    raise ToolError("update_slide_text requires a valid slide and shape_name")
                found = False
                for shape in prs.slides[slide_num].shapes:
                    if shape.name == shape_name and shape.has_text_frame:
                        expected = str(params.get("text") or "")
                        shape.text_frame.text = expected
                        verification.append(("shape_text", slide_num, shape_name, expected))
                        found = True
                        break
                if not found:
                    raise ToolError(f"Shape not found or not text-editable: {shape_name}")

            elif action == "set_notes":
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("set_notes slide is out of range")
                expected = str(params.get("text") or "")
                prs.slides[slide_num].notes_slide.notes_text_frame.text = expected
                verification.append(("notes", slide_num, expected))

            elif action == "add_textbox":
                from pptx.util import Inches
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("add_textbox slide is out of range")
                expected = str(params.get("text") or "")
                box = prs.slides[slide_num].shapes.add_textbox(
                    Inches(float(params.get("left", 1))),
                    Inches(float(params.get("top", 1))),
                    Inches(float(params.get("width", 6))),
                    Inches(float(params.get("height", 1))),
                )
                box.text_frame.text = expected
                verification.append(("slide_contains", slide_num, expected))

            elif action == "add_image":
                from pptx.util import Inches
                slide_num = int(params.get("slide", 1)) - 1
                image_path = params.get("image_path")
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("add_image slide is out of range")
                if not image_path or not os.path.isfile(image_path):
                    raise ToolError(f"add_image image not found: {image_path}")
                before = len(prs.slides[slide_num].shapes)
                kwargs = {
                    "left": Inches(float(params.get("left", 1))),
                    "top": Inches(float(params.get("top", 1))),
                }
                if params.get("width") is not None:
                    kwargs["width"] = Inches(float(params["width"]))
                if params.get("height") is not None:
                    kwargs["height"] = Inches(float(params["height"]))
                prs.slides[slide_num].shapes.add_picture(image_path, **kwargs)
                verification.append(("shape_count", slide_num, before + 1))

            elif action == "delete_slide":
                slide_num = int(params.get("slide", 1)) - 1
                if not (0 <= slide_num < len(prs.slides)):
                    raise ToolError("delete_slide slide is out of range")
                slide_id = prs.slides._sldIdLst[slide_num]
                prs.part.drop_rel(slide_id.rId)
                prs.slides._sldIdLst.remove(slide_id)
                expected_slide_count -= 1

            else:
                raise ToolError(f"Unknown PPTX modify action: {action}")
            applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        same_path = os.path.abspath(input_path) == os.path.abspath(output_path)
        save_path = output_path + ".rastacoder_tmp.pptx" if same_path else output_path
        try:
            prs.save(save_path)
            check = Presentation(save_path)
            if len(check.slides) != expected_slide_count:
                raise ToolError("PPTX slide-count verification failed after save")
            for item in verification:
                kind = item[0]
                if kind == "shape_text":
                    _, slide_num, shape_name, expected = item
                    matches = [shape for shape in check.slides[slide_num].shapes if shape.name == shape_name and shape.has_text_frame]
                    if not matches or matches[0].text_frame.text != expected:
                        raise ToolError("PPTX shape-text verification failed after save")
                elif kind == "notes":
                    _, slide_num, expected = item
                    if check.slides[slide_num].notes_slide.notes_text_frame.text != expected:
                        raise ToolError("PPTX notes verification failed after save")
                elif kind == "slide_contains":
                    _, slide_num, expected = item
                    texts = [shape.text for shape in check.slides[slide_num].shapes if getattr(shape, "has_text_frame", False)]
                    if expected not in texts:
                        raise ToolError("PPTX textbox verification failed after save")
                elif kind == "shape_count":
                    _, slide_num, expected = item
                    if len(check.slides[slide_num].shapes) < expected:
                        raise ToolError("PPTX image verification failed after save")
            if same_path:
                os.replace(save_path, output_path)
        finally:
            if same_path and os.path.exists(save_path):
                try:
                    os.remove(save_path)
                except OSError:
                    pass

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
            "operations_requested": len(operations),
            "verified": True,
            "in_place": same_path,
        }
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"Failed to modify PPTX: {exc}")


# ---------------------------------------------------------------------------
# XLSX read/write

# ---------------------------------------------------------------------------
# XLSX read/write
# ---------------------------------------------------------------------------

def read_xlsx(xlsx_path: str, sheet: str = None, range: str = None, extract: str = "values") -> dict:
    """
    Extract cell data, sheet names, and formulas from an XLSX file.

    Args:
        xlsx_path: Path to the XLSX file
        sheet: Sheet name or index (omit for all sheets)
        range: Cell range like "A1:D10" (omit for all data)
        extract: What to extract — "values", "formulas", or "all"

    Returns:
        Dict with extracted data
    """
    from openpyxl import load_workbook

    validate_file_for_processing(xlsx_path, 'document')

    try:
        wb = load_workbook(xlsx_path, data_only=(extract == "values"))
        result = {
            "path": xlsx_path,
            "sheet_names": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }

        # Determine which sheets to process
        if sheet is not None:
            if sheet.isdigit():
                idx = int(sheet)
                if idx < len(wb.sheetnames):
                    sheets_to_read = [wb.sheetnames[idx]]
                else:
                    raise ToolError(f"Sheet index {idx} out of range. Available: {len(wb.sheetnames)} sheets.")
            elif sheet in wb.sheetnames:
                sheets_to_read = [sheet]
            else:
                raise ToolError(f"Sheet '{sheet}' not found. Available: {wb.sheetnames}")
        else:
            sheets_to_read = wb.sheetnames

        max_rows = PROCESSING_LIMITS.get('xlsx_rows', 100_000)
        sheets_data = {}

        for sheet_name in sheets_to_read:
            ws = wb[sheet_name]

            if range:
                cells = ws[range]
            else:
                cells = ws.iter_rows()

            rows_data = []
            row_count = 0
            for row in cells:
                if row_count >= max_rows:
                    rows_data.append(["[Truncated — max rows exceeded]"])
                    break
                if extract == "formulas":
                    rows_data.append([cell.value if cell.data_type == 'f' else cell.value for cell in row])
                else:
                    rows_data.append([cell.value for cell in row])
                row_count += 1

            sheets_data[sheet_name] = {
                "rows": rows_data,
                "row_count": row_count,
                "dimensions": ws.dimensions,
            }

        result["sheets"] = sheets_data
        wb.close()

        return result

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read XLSX: {str(e)}")


def modify_xlsx(input_path: str, output_path: str, operations: list) -> dict:
    """Modify XLSX transactionally and reject silent no-op operations."""
    from openpyxl import load_workbook

    validate_file_for_processing(input_path, "document")
    if not isinstance(operations, list) or not operations:
        raise ToolError("modify_xlsx requires a non-empty operations list")

    try:
        wb = load_workbook(input_path)
        applied = 0
        verification = []

        for index, op in enumerate(operations):
            if not isinstance(op, dict):
                raise ToolError(f"XLSX operation #{index + 1} must be an object")
            action = str(op.get("action") or "").strip()
            params = op.get("params") if isinstance(op.get("params"), dict) else {}

            if action in {"set_cell", "set_formula"}:
                sheet_name = str(params.get("sheet") or wb.sheetnames[0])
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet '{sheet_name}' not found.")
                cell = str(params.get("cell") or "").strip()
                if not cell:
                    raise ToolError(f"{action} requires params.cell")
                value = params.get("value") if action == "set_cell" else str(params.get("formula") or "")
                wb[sheet_name][cell] = value
                verification.append(("cell", sheet_name, cell, value))

            elif action == "add_row":
                sheet_name = str(params.get("sheet") or wb.sheetnames[0])
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet '{sheet_name}' not found.")
                values = params.get("values")
                if not isinstance(values, list):
                    raise ToolError("add_row requires params.values as a list")
                before = wb[sheet_name].max_row
                wb[sheet_name].append(values)
                verification.append(("min_rows", sheet_name, before + 1))

            elif action == "add_sheet":
                name = str(params.get("name") or "").strip()
                if not name:
                    raise ToolError("add_sheet requires params.name")
                if name in wb.sheetnames:
                    raise ToolError(f"Sheet already exists: {name}")
                wb.create_sheet(title=name)
                verification.append(("sheet_present", name))

            elif action == "delete_sheet":
                name = str(params.get("name") or "").strip()
                if not name or name not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {name}")
                if len(wb.sheetnames) <= 1:
                    raise ToolError("Cannot delete the only worksheet")
                del wb[name]
                verification.append(("sheet_absent", name))

            elif action == "rename_sheet":
                old_name = str(params.get("old_name") or params.get("sheet") or "").strip()
                new_name = str(params.get("new_name") or "").strip()[:31]
                if not old_name or old_name not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {old_name}")
                if not new_name:
                    raise ToolError("rename_sheet requires params.new_name")
                if new_name in wb.sheetnames and new_name != old_name:
                    raise ToolError(f"Destination sheet name already exists: {new_name}")
                wb[old_name].title = new_name
                verification.append(("sheet_absent", old_name))
                verification.append(("sheet_present", new_name))

            elif action in {"insert_row", "delete_row", "insert_column", "delete_column"}:
                sheet_name = str(params.get("sheet") or wb.active.title)
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {sheet_name}")
                index_value = int(params.get("index", 1))
                amount = int(params.get("amount", 1))
                if index_value < 1 or amount < 1:
                    raise ToolError(f"{action} requires index>=1 and amount>=1")
                ws = wb[sheet_name]
                if action == "insert_row":
                    ws.insert_rows(index_value, amount)
                elif action == "delete_row":
                    ws.delete_rows(index_value, amount)
                elif action == "insert_column":
                    ws.insert_cols(index_value, amount)
                else:
                    ws.delete_cols(index_value, amount)

            else:
                raise ToolError(f"Unknown XLSX modify action: {action}")
            applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        same_path = os.path.abspath(input_path) == os.path.abspath(output_path)
        save_path = output_path + ".rastacoder_tmp.xlsx" if same_path else output_path
        try:
            wb.save(save_path)
            wb.close()
            check = load_workbook(save_path, data_only=False)
            for item in verification:
                kind = item[0]
                if kind == "cell":
                    _, sheet_name, cell, expected = item
                    if check[sheet_name][cell].value != expected:
                        raise ToolError("XLSX cell verification failed after save")
                elif kind == "min_rows":
                    _, sheet_name, minimum = item
                    if check[sheet_name].max_row < minimum:
                        raise ToolError("XLSX row verification failed after save")
                elif kind == "sheet_present" and item[1] not in check.sheetnames:
                    raise ToolError("XLSX sheet-add/rename verification failed after save")
                elif kind == "sheet_absent" and item[1] in check.sheetnames:
                    raise ToolError("XLSX sheet-delete/rename verification failed after save")
            check.close()
            if same_path:
                os.replace(save_path, output_path)
        finally:
            if same_path and os.path.exists(save_path):
                try:
                    os.remove(save_path)
                except OSError:
                    pass

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
            "operations_requested": len(operations),
            "verified": True,
            "in_place": same_path,
        }
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"Failed to modify XLSX: {exc}")

